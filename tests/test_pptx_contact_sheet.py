import json
from pathlib import Path

from pptx import Presentation

from luddite.analysis.render_pptx_contact_sheet import (
    ContactSheetTarget,
    ThumbnailGeneration,
    check_contact_sheet_backend,
    default_targets,
    render_pptx_contact_sheet,
)


def _make_pptx(path: Path, slide_count: int = 2) -> None:
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    for _ in range(slide_count):
        presentation.slides.add_slide(blank)
    presentation.save(path)


def _make_slide_spec(path: Path, slide_count: int = 2) -> None:
    slides = [
        {
            "slide_id": f"slide_{slide_no:03d}",
            "slide_no": slide_no,
            "section_id": "section_01",
            "layout_intent": "diagram",
            "screen_headline": f"Headline {slide_no}",
            "screen_body": ["short body"],
            "speaker_notes_expanded": "notes",
            "overflow_notes": [],
            "proof_object": {
                "type": "diagram",
                "screen_position": "center_large",
                "diagram_nodes": [
                    "정부가 정책자금을 공급함",
                    "기업의 장기 투자 위험을 나눔",
                    "손실 분담 논쟁이 생김",
                ],
                "diagram_edges": [],
                "manual_insert_required": False,
                "copyright_risk": False,
            },
            "editor_instruction": None,
            "source_refs": [],
            "risk_flags": [],
            "needs_source": False,
            "needs_fact_check": False,
            "required_before_broadcast": False,
            "do_not_claim": [],
        }
        for slide_no in range(1, slide_count + 1)
    ]
    path.write_text(
        json.dumps(
            {
                "deck_id": "contact_sheet_fixture",
                "story_seed_title": "Fixture",
                "sections": [
                    {
                        "section_id": "section_01",
                        "section_no": 1,
                        "section_title": "Section",
                        "purpose": "Test",
                        "slides": slides,
                    }
                ],
                "slides": slides,
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )


def test_default_targets_include_adapter_and_direct_live() -> None:
    targets = default_targets(include_direct=True, direct_run_id="live_fixture")

    assert len(targets) == 4
    assert [target.source_kind for target in targets] == [
        "adapter",
        "adapter",
        "direct_live",
        "direct_live",
    ]
    assert targets[2].pptx_path.parts[-3:] == (
        "live_fixture",
        "ai_knowledge_institution",
        "direct_piti_slide_spec_draft.pptx",
    )


def test_contact_sheet_reports_backend_warning_without_crashing(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    spec_path = tmp_path / "spec.json"
    output_dir = tmp_path / "qa"
    _make_pptx(pptx_path)
    _make_slide_spec(spec_path)

    def unavailable_backend(
        target: ContactSheetTarget,
        output: Path,
    ) -> ThumbnailGeneration:
        return ThumbnailGeneration(
            status="backend_unavailable_or_failed",
            backend="none",
            thumbnails=[],
            contact_sheet_path=None,
            contact_sheet_pdf_path=None,
            pdf_path=None,
            warnings=["LibreOffice/soffice not found; thumbnail generation skipped."],
        )

    results = render_pptx_contact_sheet(
        targets=[
            ContactSheetTarget(
                deck_id="fixture_deck",
                pptx_path=pptx_path,
                slide_spec_path=spec_path,
                source_kind="test",
            )
        ],
        output_dir=output_dir,
        summary_output_path=output_dir / "summary.md",
        review_summary_output_path=tmp_path / "review_summary.md",
        thumbnail_generator=unavailable_backend,
    )

    assert results[0].status == "metadata_only_with_warning"
    assert (output_dir / "summary.md").exists()
    assert (tmp_path / "review_summary.md").exists()
    report = (output_dir / "fixture_deck_contact_sheet.md").read_text(encoding="utf-8")
    assert "contact_sheet_review_status" in report
    assert "unchecked" in report
    assert "No LLM/API calls" in report
    assert "thumbnail_missing" in report


def test_backend_check_reports_missing_and_found_commands(monkeypatch) -> None:
    monkeypatch.setattr(
        "luddite.analysis.render_pptx_contact_sheet._pillow_available",
        lambda: (True, None),
    )

    def missing_finder(candidates: list[str]) -> str | None:
        return None

    missing = check_contact_sheet_backend(command_finder=missing_finder)

    assert not missing.libreoffice_found
    assert not missing.pdftoppm_found
    assert missing.pillow_available
    assert not missing.thumbnail_backend_ready

    def found_finder(candidates: list[str]) -> str | None:
        if "soffice" in candidates:
            return "/usr/local/bin/soffice"
        if "pdftoppm" in candidates:
            return "/usr/local/bin/pdftoppm"
        return None

    found = check_contact_sheet_backend(command_finder=found_finder)

    assert found.libreoffice_found
    assert found.pdftoppm_found
    assert found.thumbnail_backend_ready


def test_backend_check_reports_pillow_missing(monkeypatch) -> None:
    monkeypatch.setattr(
        "luddite.analysis.render_pptx_contact_sheet._pillow_available",
        lambda: (False, "missing pillow"),
    )

    def found_finder(candidates: list[str]) -> str | None:
        return f"/mock/{candidates[0]}"

    backend = check_contact_sheet_backend(command_finder=found_finder)

    assert backend.libreoffice_found
    assert backend.pdftoppm_found
    assert not backend.pillow_available
    assert backend.pillow_error == "missing pillow"
    assert not backend.thumbnail_backend_ready


def test_contact_sheet_reports_missing_pptx_gracefully(tmp_path: Path) -> None:
    output_dir = tmp_path / "qa"
    results = render_pptx_contact_sheet(
        targets=[
            ContactSheetTarget(
                deck_id="missing_deck",
                pptx_path=tmp_path / "missing.pptx",
                slide_spec_path=None,
                source_kind="test",
            )
        ],
        output_dir=output_dir,
        summary_output_path=output_dir / "summary.md",
        review_summary_output_path=None,
    )

    assert results[0].status == "missing_pptx"
    summary = (output_dir / "summary.md").read_text(encoding="utf-8")
    assert "missing_pptx" in summary
    assert "Broadcast readiness remains false" in summary


def test_contact_sheet_uses_mock_thumbnail_generator(tmp_path: Path) -> None:
    pptx_path = tmp_path / "deck.pptx"
    spec_path = tmp_path / "spec.json"
    output_dir = tmp_path / "qa"
    _make_pptx(pptx_path)
    _make_slide_spec(spec_path)

    def generated_backend(
        target: ContactSheetTarget,
        output: Path,
    ) -> ThumbnailGeneration:
        deck_dir = output / target.deck_id
        deck_dir.mkdir(parents=True)
        thumbnails = [deck_dir / "slide-1.png", deck_dir / "slide-2.png"]
        for thumbnail in thumbnails:
            thumbnail.write_text("mock", encoding="utf-8")
        contact_sheet = deck_dir / "fixture_contact_sheet.png"
        contact_sheet.write_text("mock", encoding="utf-8")
        return ThumbnailGeneration(
            status="generated",
            backend="mock",
            thumbnails=thumbnails,
            contact_sheet_path=contact_sheet,
            contact_sheet_pdf_path=None,
            pdf_path=None,
            warnings=[],
        )

    results = render_pptx_contact_sheet(
        targets=[
            ContactSheetTarget(
                deck_id="generated_deck",
                pptx_path=pptx_path,
                slide_spec_path=spec_path,
                source_kind="test",
            )
        ],
        output_dir=output_dir,
        summary_output_path=output_dir / "summary.md",
        review_summary_output_path=None,
        thumbnail_generator=generated_backend,
    )

    assert results[0].status == "contact_sheet_generated"
    assert results[0].thumbnail_count == 2
    summary = (output_dir / "summary.md").read_text(encoding="utf-8")
    assert "fixture_contact_sheet.png" in summary
