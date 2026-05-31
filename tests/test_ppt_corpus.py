import json
import zipfile

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from luddite.ppt.corpus import (
    GOOGLE_SLIDES_MIME,
    PPTX_MIME,
    build_drive_manifest,
    build_inventory,
    build_ppt_corpus_insight_reports,
    extract_ppt_corpus_slides,
    infer_deck_type,
    normalize_inventory_rows,
)
from luddite.utils.jsonl import read_jsonl


def _write_jsonl(path, rows):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        "\n".join(json.dumps(row, ensure_ascii=False) for row in rows) + "\n",
        encoding="utf-8",
    )


def _make_fixture_pptx(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "왜 가격이 10배가 됐나?"
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8), Inches(2))
    box.text_frame.text = "공식자료 https://example.com/report?utm_source=x&keep=1"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "시장 비교"
    table = slide.shapes.add_table(2, 2, Inches(0.8), Inches(1.3), Inches(3), Inches(1.2)).table
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "10"
    table.cell(1, 1).text = "20"
    data = CategoryChartData()
    data.categories = ["A", "B"]
    data.add_series("growth", (10, 20))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(4.2),
        Inches(1.3),
        Inches(4),
        Inches(2.8),
        data,
    )
    prs.save(path)
    _add_notes_slide(path)


def _add_notes_slide(path):
    rels_name = "ppt/slides/_rels/slide1.xml.rels"
    rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
    notes_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
 xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree><p:sp><p:txBody><a:bodyPr/><a:lstStyle/>
    <a:p><a:r><a:t>[내용] https://notes.example.com/source</a:t></a:r></a:p>
  </p:txBody></p:sp></p:spTree></p:cSld>
</p:notes>
"""
    with zipfile.ZipFile(path) as deck:
        entries = {info.filename: deck.read(info.filename) for info in deck.infolist()}
    rels = entries[rels_name].decode("utf-8")
    relationship = (
        f'<Relationship Id="rIdNotesTest" Type="{rel_type}" '
        'Target="../notesSlides/notesSlide1.xml"/>'
    )
    entries[rels_name] = rels.replace(
        "</Relationships>", relationship + "</Relationships>"
    ).encode()
    entries["ppt/notesSlides/notesSlide1.xml"] = notes_xml.encode()
    tmp_path = path.with_suffix(".with-notes.pptx")
    with zipfile.ZipFile(tmp_path, "w") as deck:
        for name, payload in entries.items():
            deck.writestr(name, payload)
    tmp_path.replace(path)


def test_inventory_normalizes_schema_mimes_and_duplicate_guard() -> None:
    records = normalize_inventory_rows(
        [
            {
                "drive_file_id": "deck-1",
                "title": "전당포 주식회사_배형찬",
                "mime_type": PPTX_MIME,
                "local_path": "data/ppt_corpus/drive_raw/latest/f88.pptx",
            },
            {
                "drive_file_id": "deck-1",
                "title": "중복 파일",
                "mime_type": PPTX_MIME,
                "local_path": "data/ppt_corpus/drive_raw/latest/dupe.pptx",
            },
            {
                "drive_file_id": "slides-1",
                "title": "슈카월드 템플릿",
                "mime_type": GOOGLE_SLIDES_MIME,
            },
        ]
    )

    assert len(records) == 2
    assert records[0]["inferred_author"] == "배형찬"
    assert records[0]["inferred_topic_title"] == "전당포 주식회사"
    assert records[0]["inferred_deck_type"] == "topic_deck"
    assert records[1]["needs_download"]
    assert records[1]["skip_reason"] == "google_slides_export_required"
    assert records[1]["inferred_deck_type"] == "template"
    assert infer_deck_type("작업 매뉴얼") == "manual"


def test_build_drive_manifest_scans_roots_to_canonical_schema(tmp_path) -> None:
    latest = tmp_path / "latest"
    past = tmp_path / "past"
    pptx = latest / "슈카월드" / "20260517 방송용 (직원)" / "전당포 주식회사_배형찬.pptx"
    old_ppt = past / "과거 자료" / "김성원_월 200으로는 성공할 수 없는가.ppt"
    non_ppt = past / "과거 자료" / "정산.xlsx"
    for path in [pptx, old_ppt, non_ppt]:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(b"placeholder")
    output = tmp_path / "drive_ppt_inventory.jsonl"
    report = tmp_path / "inventory.md"

    rows = build_drive_manifest(
        latest_root=latest,
        past_root=past,
        output_jsonl=output,
        report_md=report,
    )

    assert len(rows) == 2
    assert rows[0]["drive_file_id"].startswith("local_drive_")
    assert rows[0]["inferred_topic_title"] == "전당포 주식회사"
    assert rows[0]["inferred_author"] == "배형찬"
    assert rows[1]["mime_type"].endswith("ms-powerpoint")
    assert read_jsonl(output)[0]["inferred_deck_type"] == "daily_deck"
    assert "정산.xlsx" in report.read_text(encoding="utf-8")


def test_extract_ppt_corpus_slides_writes_canonical_outputs(tmp_path) -> None:
    pptx = tmp_path / "fixture.pptx"
    _make_fixture_pptx(pptx)
    raw_manifest = tmp_path / "raw_inventory.jsonl"
    inventory = tmp_path / "drive_ppt_inventory.jsonl"
    _write_jsonl(
        raw_manifest,
        [
            {
                "drive_file_id": "fixture",
                "title": "테스트 시장_김성원",
                "mime_type": PPTX_MIME,
                "local_path": str(pptx),
            }
        ],
    )
    build_inventory(
        manifest=raw_manifest,
        output_jsonl=inventory,
        report_md=tmp_path / "inventory.md",
    )

    deck_rows, slide_rows, link_rows, media_rows = extract_ppt_corpus_slides(
        inventory_jsonl=inventory,
        deck_manifest_jsonl=tmp_path / "deck_manifest.jsonl",
        slides_jsonl=tmp_path / "slides.jsonl",
        links_jsonl=tmp_path / "links.jsonl",
        media_manifest_jsonl=tmp_path / "media_manifest.jsonl",
        report_md=tmp_path / "quality.md",
    )

    assert deck_rows[0]["parse_status"] == "parsed"
    assert len(slide_rows) == 2
    assert slide_rows[0]["possible_role"] == "opening_question"
    assert slide_rows[0]["notes_char_count"] > 0
    assert any(row["url"] == "https://example.com/report?keep=1" for row in link_rows)
    assert any(row["url"] == "https://notes.example.com/source" for row in link_rows)
    assert media_rows[1]["table_count"] == 1
    assert media_rows[1]["chart_count"] == 1


def test_extract_ppt_corpus_slides_records_corrupt_warning(tmp_path) -> None:
    corrupt = tmp_path / "corrupt.pptx"
    corrupt.write_bytes(b"not a real pptx")
    inventory = tmp_path / "drive_ppt_inventory.jsonl"
    _write_jsonl(
        inventory,
        [
            {
                "drive_file_id": "broken",
                "title": "깨진 PPT",
                "mime_type": PPTX_MIME,
                "local_path": str(corrupt),
            }
        ],
    )
    report = tmp_path / "quality.md"

    deck_rows, slide_rows, link_rows, media_rows = extract_ppt_corpus_slides(
        inventory_jsonl=inventory,
        deck_manifest_jsonl=tmp_path / "deck_manifest.jsonl",
        slides_jsonl=tmp_path / "slides.jsonl",
        links_jsonl=tmp_path / "links.jsonl",
        media_manifest_jsonl=tmp_path / "media_manifest.jsonl",
        report_md=report,
    )

    assert deck_rows[0]["parse_status"] == "failed"
    assert slide_rows == []
    assert link_rows == []
    assert media_rows == []
    assert "Parse Warnings" in report.read_text(encoding="utf-8")


def test_build_ppt_corpus_insight_reports_from_synthetic_rows(tmp_path) -> None:
    deck_manifest = tmp_path / "deck_manifest.jsonl"
    slides = tmp_path / "slides.jsonl"
    links = tmp_path / "links.jsonl"
    _write_jsonl(
        deck_manifest,
        [{"deck_id": "deck-1", "inferred_topic_title": "생활 물가", "slide_count": 2}],
    )
    _write_jsonl(
        slides,
        [
            {
                "deck_id": "deck-1",
                "slide_number": 1,
                "possible_role": "opening_question",
                "visible_char_count": 20,
                "visual_density_estimate": "low",
            },
            {
                "deck_id": "deck-1",
                "slide_number": 2,
                "possible_role": "big_number",
                "visible_char_count": 40,
                "visual_density_estimate": "medium",
            },
        ],
    )
    _write_jsonl(links, [{"deck_id": "deck-1", "domain": "korea.kr"}])

    outputs = build_ppt_corpus_insight_reports(
        deck_manifest_jsonl=deck_manifest,
        slides_jsonl=slides,
        links_jsonl=links,
        jibi_md=tmp_path / "docs" / "jibi.md",
        anny_md=tmp_path / "docs" / "anny.md",
        piti_md=tmp_path / "docs" / "piti.md",
        report_dir=tmp_path / "reports",
    )

    assert "No scoring rules changed" in outputs["jibi"].read_text(encoding="utf-8")
    assert "생활 물가" in outputs["anny"].read_text(encoding="utf-8")
    assert "average_slides_per_deck: 2.0" in outputs["piti"].read_text(encoding="utf-8")
    assert (tmp_path / "reports" / "07_recommended_code_changes.md").exists()
