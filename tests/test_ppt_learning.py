import json
import sqlite3

from pptx import Presentation
from pptx.util import Inches

from luddite.analysis.render_pptx_contact_sheet import ThumbnailGeneration
from luddite.ppt.learning import (
    SourcePageHttpResponse,
    build_drive_manifest,
    build_inventory,
    build_ppt_enrichment_queue,
    build_ppt_slide_visual_memos,
    build_ppt_story_arc_memos,
    build_seed_lesson_review_queue,
    canonical_domain,
    clean_number_token,
    clean_ppt_title,
    extract_slide_sources_from_parsed,
    fetch_ppt_source_memos,
    infer_evidence_types,
    match_broadcast_usage,
    normalize_url_for_queue,
    seed_lesson_from_parsed,
    source_access_hint,
    source_domain_category,
    write_combined_report,
)
from luddite.utils.jsonl import read_jsonl


def _write_manifest(path, rows):
    path.write_text(
        "\n".join(json.dumps(row, ensure_ascii=False) for row in rows) + "\n",
        encoding="utf-8",
    )


def _make_pptx(path, title, body):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(3))
    box.text_frame.text = body
    prs.save(path)


def _make_two_slide_pptx(path):
    prs = Presentation()
    for title, body in [
        ("Fixture hook", "Opening hook"),
        ("Fixture evidence", "10배 성장 chart"),
    ]:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(3))
        box.text_frame.text = body
    prs.save(path)


def _make_syuka_db(path, *, include_syuka=True, include_money=False):
    conn = sqlite3.connect(path)
    conn.execute(
        """
        CREATE TABLE videos (
            video_id TEXT PRIMARY KEY,
            title TEXT,
            upload_date TEXT,
            view_count INTEGER,
            like_count INTEGER,
            channel_name TEXT,
            channel_key TEXT,
            source_url TEXT
        )
        """
    )
    conn.execute("CREATE TABLE video_analysis (video_id TEXT, summary TEXT, keywords_json TEXT)")
    conn.execute("CREATE TABLE transcripts (video_id TEXT, dialogue TEXT)")
    if include_syuka:
        conn.execute(
            "INSERT INTO videos VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
            (
                "vid_f88",
                "베트남 전당포 F88, 상장을 향해 가다",
                "20260516",
                1234567,
                45000,
                "슈카월드",
                "syukaworld",
                "https://youtu.be/f88",
            ),
        )
        conn.execute(
            "INSERT INTO video_analysis VALUES (?, ?, ?)",
            (
                "vid_f88",
                "F88 전당포 상장과 오토바이 담보대출, 추심 리스크",
                '["F88", "전당포", "오토바이 담보대출"]',
            ),
        )
        conn.execute(
            "INSERT INTO transcripts VALUES (?, ?)",
            ("vid_f88", "베트남 F88은 오토바이 등록증 담보대출과 추심 리스크가 있습니다."),
        )
    if include_money:
        conn.execute(
            "INSERT INTO videos VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
            (
                "vid_money",
                "베트남 전당포 F88, 상장을 향해 가다",
                "20260516",
                100,
                10,
                "머니코믹스",
                "moneymoneycomics",
                "https://youtu.be/money",
            ),
        )
        conn.execute(
            "INSERT INTO video_analysis VALUES (?, ?, ?)",
            ("vid_money", "F88 전당포 상장과 오토바이 담보대출", '["F88", "전당포"]'),
        )
        conn.execute(
            "INSERT INTO transcripts VALUES (?, ?)",
            ("vid_money", "베트남 F88 전당포 상장 이야기입니다."),
        )
    conn.commit()
    conn.close()


class _FakeSourceClient:
    def __init__(self, responses):
        self.responses = responses

    def fetch(self, url: str, *, timeout: float) -> SourcePageHttpResponse:
        response = self.responses[url]
        if isinstance(response, Exception):
            raise response
        return response


def test_slide_source_extraction_classifies_content_image_and_evidence() -> None:
    parsed = {
        "slides": [
            {
                "slide_no": 1,
                "headline": "F88 매출은 10배 성장",
                "visible_text": "F88은 베트남 전당포 체인이고 매출은 10배 성장했다.",
                "notes": "[내용] https://asia.nikkei.com/f88\n[이미지] https://f88.vn/aboutf88",
                "source_notes": [
                    {
                        "label": "내용",
                        "value": "https://asia.nikkei.com/f88",
                        "urls": ["https://asia.nikkei.com/f88"],
                        "is_image": False,
                        "is_gpt_generated": False,
                    },
                    {
                        "label": "이미지",
                        "value": "https://f88.vn/aboutf88",
                        "urls": ["https://f88.vn/aboutf88"],
                        "is_image": True,
                        "is_gpt_generated": False,
                    },
                ],
                "text_urls": [],
                "notes_urls": ["https://asia.nikkei.com/f88", "https://f88.vn/aboutf88"],
                "relationship_urls": [],
                "all_urls": ["https://asia.nikkei.com/f88", "https://f88.vn/aboutf88"],
                "media_count": 1,
                "slide_type": "data",
            }
        ]
    }

    records = extract_slide_sources_from_parsed(
        {"ppt_id": "f88", "title": "전당포 주식회사"},
        parsed,
    )

    assert records[0]["extracted_urls"] == [
        "https://asia.nikkei.com/f88",
        "https://f88.vn/aboutf88",
    ]
    assert [entry["source_kind"] for entry in records[0]["source_entries"]] == [
        "content_url",
        "image_url",
    ]
    assert "number/statistic" in records[0]["evidence_types"]
    assert "company" in records[0]["evidence_types"]


def test_evidence_heuristic_and_f88_seed_lesson_signals() -> None:
    text = (
        "베트남 전당포 F88은 상장을 추진한다. 매출은 10배 성장했고 "
        "오토바이 등록증 담보대출, 창업자 일화, 추심 리스크가 있다."
    )
    evidence = infer_evidence_types(text)
    parsed = {
        "title": "전당포 주식회사",
        "urls": ["https://asia.nikkei.com/f88"],
        "slides": [
            {
                "slide_no": 1,
                "headline": "최근 상장에 도전하는 금융 회사가 있다",
                "visible_text": text,
                "notes": "[내용] https://asia.nikkei.com/f88",
                "source_notes": [
                    {
                        "label": "내용",
                        "value": "https://asia.nikkei.com/f88",
                        "urls": ["https://asia.nikkei.com/f88"],
                        "is_image": False,
                        "is_gpt_generated": False,
                    }
                ],
                "all_urls": ["https://asia.nikkei.com/f88"],
                "text_urls": [],
                "notes_urls": ["https://asia.nikkei.com/f88"],
                "relationship_urls": [],
                "media_count": 1,
                "slide_type": "data",
            }
        ],
    }

    lesson = seed_lesson_from_parsed(
        {"ppt_id": "f88", "title": "전당포 주식회사", "local_path": "f88.pptx"},
        parsed,
    )

    assert "number/statistic" in evidence
    assert "everyday_object" in evidence
    assert "risk/reversal" in evidence
    assert lesson["main_seed_hook"] == "베트남 전당포 체인 F88이 상장을 추진한다"
    assert "everyday_collateral_object" in lesson["jibi_candidate_signals"]
    assert "regulatory_or_collection_risk" in lesson["jibi_candidate_signals"]
    assert "추심 리스크와 제도권화" in lesson["story_expansion_path"]


def test_inventory_missing_ppt_does_not_crash(tmp_path) -> None:
    manifest = tmp_path / "manifest.jsonl"
    output = tmp_path / "inventory.jsonl"
    report = tmp_path / "inventory.md"
    _write_manifest(
        manifest,
        [
            {
                "ppt_id": "missing",
                "title": "없는 PPT",
                "researcher": "테스트",
                "local_path": str(tmp_path / "missing.pptx"),
                "folder_date_hint": "2026-05-15",
            }
        ],
    )

    records = build_inventory(manifest=manifest, output_jsonl=output, report_md=report)

    assert records[0]["path_status"] == "missing"
    assert records[0]["extracted_text_status"] == "missing"
    assert read_jsonl(output)[0]["ppt_id"] == "missing"
    assert "없는 PPT" in report.read_text(encoding="utf-8")


def test_drive_manifest_scans_ppt_roots_and_reports_non_ppt(tmp_path) -> None:
    latest = tmp_path / "latest"
    past = tmp_path / "past"
    latest_ppt = latest / "슈카월드" / "20260517 방송용 (직원)" / "전당포 주식회사_배형찬.pptx"
    past_ppt = past / "과거 자료" / "김성원_월 200으로는 성공할 수 없는가.pptx"
    non_ppt = past / "과거 자료" / "정산.xlsx"
    for path in [latest_ppt, past_ppt, non_ppt]:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(b"placeholder")
    manifest = tmp_path / "drive_ppts.jsonl"
    report = tmp_path / "non_ppt.md"

    rows = build_drive_manifest(
        latest_root=latest,
        past_root=past,
        output_jsonl=manifest,
        non_ppt_report_md=report,
    )

    assert len(rows) == 2
    assert rows[0]["title"] == "전당포 주식회사"
    assert rows[0]["researcher"] == "배형찬"
    assert rows[0]["folder_date_hint"] == "2026-05-17"
    assert rows[1]["title"] == "월 200으로는 성공할 수 없는가"
    assert rows[1]["researcher"] == "김성원"
    assert read_jsonl(manifest)[1]["source_root"] == "past"
    assert "정산.xlsx" in report.read_text(encoding="utf-8")


def test_custom_combined_report_uses_custom_jsonl_paths(tmp_path) -> None:
    inventory = tmp_path / "inventory.jsonl"
    sources = tmp_path / "sources.jsonl"
    matches = tmp_path / "matches.jsonl"
    lessons = tmp_path / "lessons.jsonl"
    report = tmp_path / "report.md"
    _write_manifest(
        inventory,
        [{"ppt_id": "f88", "title": "전당포 주식회사", "local_exists": True, "slide_count": 1}],
    )
    _write_manifest(
        sources,
        [{"ppt_id": "f88", "url_domains": ["asia.nikkei.com"], "evidence_types": ["company"]}],
    )
    _write_manifest(matches, [{"ppt_id": "f88", "broadcast_label": "not_found"}])
    _write_manifest(
        lessons,
        [
            {
                "ppt_id": "f88",
                "ppt_title": "전당포 주식회사",
                "main_seed_hook": "베트남 전당포 체인 F88이 상장을 추진한다",
                "jibi_candidate_signals": ["unfamiliar_industry"],
                "initial_seed_sources": [{"url": "https://asia.nikkei.com/f88"}],
            }
        ],
    )

    markdown = write_combined_report(
        inventory_jsonl=inventory,
        sources_jsonl=sources,
        matches_jsonl=matches,
        lessons_jsonl=lessons,
        report_md=report,
    )

    assert "asia.nikkei.com" in markdown
    assert "전당포 주식회사" in report.read_text(encoding="utf-8")


def test_visible_cleanup_number_title_and_domain_helpers() -> None:
    assert clean_number_token("10배") == "10배"
    assert clean_number_token("3,370만") == "3,370만"
    assert clean_number_token("1.25%") == "1.25%"
    assert clean_number_token("5763432159326606488") == ""
    assert clean_number_token("04") == ""

    title_info = clean_ppt_title("O 틱톡 방출령(1) (1)")
    assert title_info["clean_title"] == "틱톡 방출령"
    assert "leading_o" in title_info["title_flags"]
    assert "duplicate_copy" in title_info["title_flags"]
    assert "internal_reference" in clean_ppt_title("대표님 예시 자료")["title_flags"]

    assert canonical_domain("https://m.youtube.com/watch?v=abc") == "youtube.com"
    assert canonical_domain("twitter.com/user") == "x.com"
    assert source_domain_category("data.worldbank.org") == "official/data"
    assert source_domain_category("www.reuters.com") == "media/news"
    assert source_domain_category("youtu.be") == "social/video"
    assert (
        normalize_url_for_queue("https://www.reuters.com/world/?utm_source=x&foo=1#section")
        == "https://www.reuters.com/world?foo=1"
    )
    assert source_access_hint("www.bloomberg.com") == "likely_paywalled"


def test_seed_lesson_review_queue_buckets_gold_and_exclusions(tmp_path) -> None:
    inventory = tmp_path / "inventory.jsonl"
    sources = tmp_path / "sources.jsonl"
    matches = tmp_path / "matches.jsonl"
    lessons = tmp_path / "lessons.jsonl"
    output = tmp_path / "queue.jsonl"
    report = tmp_path / "queue.md"
    _write_manifest(
        inventory,
        [
            {"ppt_id": "f88", "title": "전당포 주식회사", "slide_count": 10},
            {"ppt_id": "internal", "title": "대표님 예시 자료", "slide_count": 10},
        ],
    )
    source_rows = []
    for idx in range(3):
        source_rows.append(
            {
                "ppt_id": "f88",
                "extracted_urls": [f"https://asia.nikkei.com/f88/{idx}"],
                "url_domains": ["asia.nikkei.com"],
                "source_entries": [{"source_kind": "content_url", "url_domain": "asia.nikkei.com"}],
                "evidence_types": ["number/statistic", "company", "image/proof_object"],
            }
        )
    source_rows.append(
        {
            "ppt_id": "f88",
            "extracted_urls": ["https://f88.vn/about"],
            "url_domains": ["f88.vn"],
            "source_entries": [{"source_kind": "image_url", "url_domain": "f88.vn"}],
            "evidence_types": ["image/proof_object", "risk/reversal", "local_context"],
        }
    )
    source_rows.append(
        {
            "ppt_id": "internal",
            "extracted_urls": ["https://example.com"],
            "url_domains": ["example.com"],
            "source_entries": [{"source_kind": "content_url", "url_domain": "example.com"}],
            "evidence_types": ["unknown"],
        }
    )
    _write_manifest(sources, source_rows)
    _write_manifest(
        matches,
        [
            {"ppt_id": "f88", "broadcast_label": "not_found"},
            {"ppt_id": "internal", "broadcast_label": "not_found"},
        ],
    )
    _write_manifest(
        lessons,
        [
            {
                "ppt_id": "f88",
                "ppt_title": "전당포 주식회사",
                "lesson_confidence": "high",
                "main_seed_hook": "베트남 전당포 체인 F88이 상장을 추진한다",
                "story_expansion_path": ["낯선 상장 후보"],
                "jibi_candidate_signals": [
                    "primary_source_rich",
                    "visual_proof_object",
                    "regulatory_or_collection_risk",
                    "korea_bridge_available",
                ],
                "numbers_used": ["5763432159326606488", "10배", "2025년"],
                "everyday_objects": ["오토바이", "차"],
            },
            {
                "ppt_id": "internal",
                "ppt_title": "대표님 예시 자료",
                "lesson_confidence": "high",
                "main_seed_hook": "내부 예시",
                "jibi_candidate_signals": ["primary_source_rich"],
                "numbers_used": ["10배"],
                "everyday_objects": [],
            },
        ],
    )

    rows = build_seed_lesson_review_queue(
        inventory_jsonl=inventory,
        sources_jsonl=sources,
        matches_jsonl=matches,
        lessons_jsonl=lessons,
        output_jsonl=output,
        report_md=report,
        show_progress=False,
    )

    assert rows[0]["review_bucket"] == "gold"
    assert rows[0]["numbers_used_clean"] == ["10배", "2025년"]
    assert rows[0]["numbers_noise_dropped"] == ["5763432159326606488"]
    assert rows[0]["everyday_objects_clean"] == ["오토바이"]
    assert rows[1]["review_bucket"] == "exclude_candidate"
    assert "internal_reference" in rows[1]["title_flags"]
    assert read_jsonl(output)[0]["review_bucket"] == "gold"
    assert "전당포 주식회사" in report.read_text(encoding="utf-8")


def test_ppt_enrichment_queue_prioritizes_urls_and_writes_flat_queue(tmp_path) -> None:
    inventory = tmp_path / "inventory.jsonl"
    sources = tmp_path / "sources.jsonl"
    matches = tmp_path / "matches.jsonl"
    lessons = tmp_path / "lessons.jsonl"
    review_queue = tmp_path / "review_queue.jsonl"
    output = tmp_path / "enrichment.jsonl"
    url_output = tmp_path / "url_queue.jsonl"
    report = tmp_path / "enrichment.md"
    _write_manifest(
        inventory,
        [
            {
                "ppt_id": "f88",
                "title": "전당포 주식회사",
                "local_path": "data/f88.pptx",
                "folder_date_hint": "2026-05-15",
                "slide_count": 12,
            }
        ],
    )
    _write_manifest(
        sources,
        [
            {
                "ppt_id": "f88",
                "slide_number": 2,
                "extracted_urls": ["https://data.worldbank.org/indicator/foo?utm_source=deck"],
                "url_domains": ["data.worldbank.org"],
                "source_entries": [
                    {
                        "source_kind": "content_url",
                        "source_note_group": "content",
                        "url": "https://data.worldbank.org/indicator/foo?utm_source=deck",
                        "url_domain": "data.worldbank.org",
                        "evidence_types": ["institution/regulation", "number/statistic"],
                    }
                ],
                "evidence_types": ["institution/regulation", "number/statistic"],
            },
            {
                "ppt_id": "f88",
                "slide_number": 3,
                "extracted_urls": [
                    "https://www.reuters.com/world/example?utm_campaign=x&id=1",
                    "https://www.bloomberg.com/news/articles/f88",
                ],
                "url_domains": ["www.reuters.com", "www.bloomberg.com"],
                "source_entries": [
                    {
                        "source_kind": "content_url",
                        "source_note_group": "content",
                        "url": "https://www.reuters.com/world/example?utm_campaign=x&id=1",
                        "url_domain": "www.reuters.com",
                        "evidence_types": ["company", "number/statistic"],
                    },
                    {
                        "source_kind": "content_url",
                        "source_note_group": "content",
                        "url": "https://www.bloomberg.com/news/articles/f88",
                        "url_domain": "www.bloomberg.com",
                        "evidence_types": ["company"],
                    },
                ],
                "evidence_types": ["company", "number/statistic"],
            },
            {
                "ppt_id": "f88",
                "slide_number": 7,
                "extracted_urls": ["https://www.youtube.com/watch?v=f88"],
                "url_domains": ["www.youtube.com"],
                "source_entries": [
                    {
                        "source_kind": "image_url",
                        "source_note_group": "image",
                        "url": "https://www.youtube.com/watch?v=f88",
                        "url_domain": "www.youtube.com",
                        "evidence_types": ["image/proof_object"],
                    }
                ],
                "evidence_types": ["image/proof_object"],
            },
        ],
    )
    _write_manifest(matches, [{"ppt_id": "f88", "broadcast_label": "not_found"}])
    _write_manifest(
        lessons,
        [
            {
                "ppt_id": "f88",
                "ppt_title": "전당포 주식회사",
                "lesson_confidence": "high",
                "main_seed_hook": "베트남 전당포 체인 F88이 상장을 추진한다",
                "story_expansion_path": ["낯선 상장 후보", "제도권화"],
                "jibi_candidate_signals": ["primary_source_rich"],
            }
        ],
    )
    _write_manifest(
        review_queue,
        [
            {
                "ppt_id": "f88",
                "ppt_title": "전당포 주식회사",
                "clean_title": "전당포 주식회사",
                "review_bucket": "gold",
                "review_reasons": ["high_confidence_clean_signals"],
                "main_seed_hook": "베트남 전당포 체인 F88이 상장을 추진한다",
                "story_expansion_path": ["낯선 상장 후보", "제도권화"],
                "jibi_candidate_signals": ["primary_source_rich"],
            }
        ],
    )

    ppt_rows, url_rows = build_ppt_enrichment_queue(
        inventory_jsonl=inventory,
        sources_jsonl=sources,
        matches_jsonl=matches,
        lessons_jsonl=lessons,
        review_queue_jsonl=review_queue,
        output_jsonl=output,
        url_queue_jsonl=url_output,
        report_md=report,
        max_urls_per_ppt=3,
        show_progress=False,
    )

    assert len(ppt_rows) == 1
    assert ppt_rows[0]["review_bucket"] == "gold"
    assert ppt_rows[0]["unique_url_count"] == 4
    assert ppt_rows[0]["queued_url_count"] == 3
    assert ppt_rows[0]["enrichment_status"] == "ready_public_fetch"
    assert ppt_rows[0]["priority_urls"][0]["domain"] == "worldbank.org"
    assert any(row["access_hint"] == "likely_paywalled" for row in url_rows)
    assert any(row["collection_hint"] == "manual_or_authenticated_session" for row in url_rows)
    assert read_jsonl(output)[0]["clean_title"] == "전당포 주식회사"
    assert len(read_jsonl(url_output)) == 3
    assert "No URL fetching" in report.read_text(encoding="utf-8")


def test_fetch_ppt_source_memos_keeps_short_memos_and_manual_requests(tmp_path) -> None:
    url_queue = tmp_path / "url_queue.jsonl"
    output = tmp_path / "memos.jsonl"
    status = tmp_path / "status.jsonl"
    manual = tmp_path / "manual.jsonl"
    report = tmp_path / "fetch.md"
    public_url = "https://example.com/public"
    teaser_url = "https://example.com/teaser"
    soft_error_url = "https://example.com/missing"
    paywall_url = "https://www.bloomberg.com/news/articles/f88"
    _write_manifest(
        url_queue,
        [
            {
                "batch_priority_rank": 1,
                "url_priority_rank": 1,
                "ppt_id": "f88",
                "ppt_title": "전당포 주식회사",
                "clean_title": "전당포 주식회사",
                "url": public_url,
                "normalized_url": public_url,
                "domain": "example.com",
                "source_category": "company/primary",
                "access_hint": "likely_public",
                "collection_hint": "public_fetch_first",
                "first_slide": 2,
                "slide_numbers": [2],
                "evidence_types": ["number/statistic", "company"],
                "source_kinds": ["content_url"],
            },
            {
                "batch_priority_rank": 1,
                "url_priority_rank": 2,
                "ppt_id": "f88",
                "ppt_title": "전당포 주식회사",
                "clean_title": "전당포 주식회사",
                "url": teaser_url,
                "normalized_url": teaser_url,
                "domain": "example.com",
                "source_category": "media/news",
                "access_hint": "likely_public",
                "collection_hint": "public_fetch_first",
                "first_slide": 3,
                "slide_numbers": [3],
                "evidence_types": ["company"],
                "source_kinds": ["content_url"],
            },
            {
                "batch_priority_rank": 1,
                "url_priority_rank": 3,
                "ppt_id": "f88",
                "ppt_title": "전당포 주식회사",
                "clean_title": "전당포 주식회사",
                "url": soft_error_url,
                "normalized_url": soft_error_url,
                "domain": "example.com",
                "source_category": "media/news",
                "access_hint": "likely_public",
                "collection_hint": "public_fetch_first",
                "first_slide": 5,
                "slide_numbers": [5],
                "evidence_types": ["company"],
                "source_kinds": ["content_url"],
            },
            {
                "batch_priority_rank": 1,
                "url_priority_rank": 4,
                "ppt_id": "f88",
                "ppt_title": "전당포 주식회사",
                "clean_title": "전당포 주식회사",
                "url": paywall_url,
                "normalized_url": paywall_url,
                "domain": "bloomberg.com",
                "source_category": "media/news",
                "access_hint": "likely_paywalled",
                "collection_hint": "manual_or_authenticated_session",
                "first_slide": 4,
                "slide_numbers": [4],
                "evidence_types": ["company"],
                "source_kinds": ["content_url"],
            },
        ],
    )
    public_html = """
    <html><head><title>F88 growth story</title>
    <meta name="description" content="F88 grew 10배 and reached 3.1 trillion in loans.">
    </head><body>
    <p>F88 is a Vietnam pawn chain with 900 stores and World Bank context.</p>
    <p>The company connects motorbike collateral to local finance risks in Korea-facing notes.</p>
    </body></html>
    """
    teaser_html = """
    <html><head><title>F88 subscriber-only update</title></head>
    <body><p>Subscribe to continue reading this premium article.</p></body></html>
    """
    soft_error_html = """
    <html><head><title>Page not found</title></head>
    <body><p>The page you requested was not found.</p></body></html>
    """
    client = _FakeSourceClient(
        {
            public_url: SourcePageHttpResponse(
                url=public_url,
                status=200,
                content_type="text/html; charset=utf-8",
                body=public_html.encode(),
            ),
            teaser_url: SourcePageHttpResponse(
                url=teaser_url,
                status=200,
                content_type="text/html; charset=utf-8",
                body=teaser_html.encode(),
            ),
            soft_error_url: SourcePageHttpResponse(
                url=soft_error_url,
                status=404,
                content_type="text/html; charset=utf-8",
                body=soft_error_html.encode(),
            ),
        }
    )

    memos, status_rows, manual_rows = fetch_ppt_source_memos(
        url_queue_jsonl=url_queue,
        output_jsonl=output,
        status_jsonl=status,
        manual_requests_jsonl=manual,
        report_md=report,
        http_client=client,
        limit=3,
        ppt_limit=1,
        show_progress=False,
        fetched_at="2026-05-26T00:00:00+00:00",
    )

    assert len(memos) == 3
    assert memos[0]["access_status"] == "fetched_public"
    assert memos[0]["memo_quality_status"] == "usable_public_memo"
    assert memos[0]["usable_as_story_evidence"] is True
    assert "body_text" not in memos[0]
    assert "10배" in memos[0]["numbers"]
    assert "World Bank" in memos[0]["institutions_regulation"]
    assert memos[1]["access_status"] == "paywalled_manual_needed"
    assert memos[1]["memo_quality_status"] == "teaser_only"
    assert memos[1]["usable_as_story_evidence"] is False
    assert memos[2]["memo_quality_status"] == "soft_error_page"
    assert memos[2]["usable_as_story_evidence"] is False
    assert len(manual_rows) == 2
    assert {row["access_status"] for row in manual_rows} == {
        "manual_or_auth_required",
        "paywalled_manual_needed",
    }
    assert read_jsonl(status)[0]["access_status"] == "fetched_public"
    assert read_jsonl(status)[2]["memo_quality_status"] == "soft_error_page"
    assert read_jsonl(output)[0]["copyright_note"].startswith("Full article body not stored")
    report_text = report.read_text(encoding="utf-8")
    assert "Full article bodies are not stored" in report_text
    assert "Memo Quality Distribution" in report_text


def test_build_ppt_slide_visual_memos_links_thumbnails_sources_and_memos(tmp_path) -> None:
    pptx = tmp_path / "deck.pptx"
    _make_two_slide_pptx(pptx)
    enrichment_queue = tmp_path / "enrichment.jsonl"
    slide_sources = tmp_path / "slide_sources.jsonl"
    source_memos = tmp_path / "source_memos.jsonl"
    visual_output = tmp_path / "slide_visuals.jsonl"
    story_output = tmp_path / "story_inputs.jsonl"
    report = tmp_path / "visual_report.md"
    render_dir = tmp_path / "render"
    images_dir = tmp_path / "images"
    sheets_dir = tmp_path / "sheets"
    url = "https://example.com/a"
    _write_manifest(
        enrichment_queue,
        [
            {
                "batch_priority_rank": 1,
                "ppt_id": "fixture_deck",
                "ppt_title": "Fixture Deck",
                "clean_title": "Fixture Deck",
                "local_path": str(pptx),
                "resolved_local_path": str(pptx),
                "slide_count": 2,
                "review_bucket": "gold",
                "main_seed_hook": "Fixture hook",
                "story_expansion_path": ["hook", "evidence"],
                "jibi_candidate_signals": ["visual_proof_object"],
            }
        ],
    )
    _write_manifest(
        slide_sources,
        [
            {
                "ppt_id": "fixture_deck",
                "slide_number": 1,
                "slide_title": "Fixture hook",
                "slide_type": "title",
                "media_count": 1,
                "extracted_urls": [],
                "url_domains": [],
                "source_entries": [],
                "evidence_types": ["image/proof_object"],
                "raw_text_excerpt": "Fixture hook",
                "notes_excerpt": "",
            },
            {
                "ppt_id": "fixture_deck",
                "slide_number": 2,
                "slide_title": "Fixture evidence",
                "slide_type": "image_centered",
                "media_count": 1,
                "extracted_urls": [url],
                "url_domains": ["example.com"],
                "source_entries": [
                    {
                        "source_kind": "content_url",
                        "url": url,
                        "url_domain": "example.com",
                    }
                ],
                "evidence_types": ["number/statistic", "chart/table"],
                "raw_text_excerpt": "10배 성장 chart",
                "notes_excerpt": url,
            },
        ],
    )
    _write_manifest(
        source_memos,
        [
            {
                "source_page_id": "memo_a",
                "ppt_id": "fixture_deck",
                "normalized_url": url,
                "domain": "example.com",
                "access_status": "fetched_public",
                "memo_quality_status": "usable_public_memo",
                "usable_as_story_evidence": True,
                "article_title": "Example source memo",
                "source_summary": "10배 성장과 chart evidence",
            }
        ],
    )

    def mock_generator(target, output_dir):
        deck_dir = output_dir / target.deck_id / "thumbnails"
        deck_dir.mkdir(parents=True, exist_ok=True)
        thumbnails = []
        for index in [1, 2]:
            path = deck_dir / f"slide-{index}.png"
            path.write_text("mock thumbnail", encoding="utf-8")
            thumbnails.append(path)
        sheet = output_dir / target.deck_id / f"{target.deck_id}_contact_sheet.png"
        sheet.parent.mkdir(parents=True, exist_ok=True)
        sheet.write_text("mock sheet", encoding="utf-8")
        return ThumbnailGeneration(
            status="generated",
            backend="mock",
            thumbnails=thumbnails,
            contact_sheet_path=sheet,
            contact_sheet_pdf_path=None,
            pdf_path=None,
            warnings=[],
        )

    visual_rows, story_rows = build_ppt_slide_visual_memos(
        enrichment_queue_jsonl=enrichment_queue,
        slide_sources_jsonl=slide_sources,
        source_memos_jsonl=source_memos,
        output_jsonl=visual_output,
        story_inputs_jsonl=story_output,
        report_md=report,
        render_output_dir=render_dir,
        slide_images_dir=images_dir,
        contact_sheets_dir=sheets_dir,
        ppt_limit=1,
        thumbnail_generator=mock_generator,
        show_progress=False,
    )

    assert len(visual_rows) == 2
    assert visual_rows[0]["story_role"] == "title_or_hook"
    assert visual_rows[0]["thumbnail_path"]
    assert visual_rows[1]["source_memo_ids"] == ["memo_a"]
    assert visual_rows[1]["source_memo_titles"] == ["Example source memo"]
    assert visual_rows[1]["source_memo_quality_statuses"] == ["usable_public_memo"]
    assert visual_rows[1]["usable_source_memo_ids"] == ["memo_a"]
    assert visual_rows[1]["usable_source_memo_titles"] == ["Example source memo"]
    assert "chart_or_table" in visual_rows[1]["visual_types"]
    assert "company_or_primary_page" in visual_rows[1]["visual_types"]
    assert len(story_rows) == 1
    assert story_rows[0]["thumbnail_count"] == 2
    assert story_rows[0]["contact_sheet_path"]
    assert read_jsonl(story_output)[0]["slide_outline"][1]["source_memo_titles"] == [
        "Example source memo"
    ]
    assert read_jsonl(story_output)[0]["slide_outline"][1]["usable_source_memo_titles"] == [
        "Example source memo"
    ]
    assert (images_dir / "fixture_deck" / "slide_001.png").exists()
    assert (sheets_dir / "fixture_deck.png").exists()
    assert "PPT Story Inputs" in report.read_text(encoding="utf-8")


def test_build_ppt_story_arc_memos_writes_json_and_markdown_reports(tmp_path) -> None:
    story_inputs = tmp_path / "story_inputs.jsonl"
    slide_visuals = tmp_path / "slide_visuals.jsonl"
    source_memos = tmp_path / "source_memos.jsonl"
    manual_requests = tmp_path / "manual.jsonl"
    lessons = tmp_path / "lessons.jsonl"
    matches = tmp_path / "matches.jsonl"
    output = tmp_path / "story_arcs.jsonl"
    report = tmp_path / "story_arc_report.md"
    report_dir = tmp_path / "story_arcs"
    url = "https://example.com/a"
    _write_manifest(
        story_inputs,
        [
            {
                "ppt_id": "fixture_deck",
                "ppt_title": "Fixture Deck",
                "clean_title": "Fixture Deck",
                "slide_count": 4,
                "thumbnail_count": 4,
                "render_status": "contact_sheet_generated",
                "contact_sheet_path": "outputs/ppt_learning/contact_sheets/fixture.png",
                "main_seed_hook": "Fixture hook becomes a story",
                "story_expansion_path": ["seed source", "numbers", "risk", "Korea bridge"],
                "jibi_candidate_signals": ["visual_proof_object"],
                "top_source_domains": ["example.com"],
            }
        ],
    )
    _write_manifest(
        slide_visuals,
        [
            {
                "ppt_id": "fixture_deck",
                "slide_number": 1,
                "story_role": "title_or_hook",
                "slide_title": "Fixture hook",
                "raw_text_excerpt": "Fixture hook",
                "visual_types": ["rendered_slide_thumbnail"],
                "thumbnail_path": "outputs/ppt_learning/slide_images/fixture/slide_001.png",
                "source_domains": [],
                "source_memo_ids": [],
                "source_memo_titles": [],
                "evidence_types": ["unknown"],
            },
            {
                "ppt_id": "fixture_deck",
                "slide_number": 2,
                "story_role": "number_or_data_evidence",
                "slide_title": "Numbers",
                "raw_text_excerpt": "10배 growth",
                "visual_types": [
                    "rendered_slide_thumbnail",
                    "embedded_image_or_screenshot",
                    "chart_or_table",
                    "company_or_primary_page",
                ],
                "thumbnail_path": "outputs/ppt_learning/slide_images/fixture/slide_002.png",
                "source_domains": ["example.com"],
                "source_memo_ids": ["memo_a", "memo_bad"],
                "source_memo_titles": ["Example source memo", "Page not found"],
                "source_memo_quality_statuses": ["usable_public_memo", "soft_error_page"],
                "usable_source_memo_ids": ["memo_a"],
                "usable_source_memo_titles": ["Example source memo"],
                "source_urls": [url],
                "evidence_types": ["number/statistic", "chart/table"],
            },
            {
                "ppt_id": "fixture_deck",
                "slide_number": 3,
                "story_role": "risk_or_reversal",
                "slide_title": "Risk",
                "raw_text_excerpt": "collection risk",
                "visual_types": ["rendered_slide_thumbnail", "risk_or_reversal"],
                "thumbnail_path": "outputs/ppt_learning/slide_images/fixture/slide_003.png",
                "source_domains": [],
                "source_memo_ids": [],
                "source_memo_titles": [],
                "evidence_types": ["risk/reversal"],
            },
            {
                "ppt_id": "fixture_deck",
                "slide_number": 4,
                "story_role": "korea_bridge",
                "slide_title": "Korea bridge",
                "raw_text_excerpt": "한국 연결",
                "visual_types": ["rendered_slide_thumbnail", "korea_bridge"],
                "thumbnail_path": "outputs/ppt_learning/slide_images/fixture/slide_004.png",
                "source_domains": [],
                "source_memo_ids": [],
                "source_memo_titles": [],
                "evidence_types": ["local_context"],
            },
        ],
    )
    _write_manifest(
        source_memos,
        [
            {
                "source_page_id": "memo_a",
                "ppt_id": "fixture_deck",
                "normalized_url": url,
                "domain": "example.com",
                "access_status": "fetched_public",
                "memo_quality_status": "usable_public_memo",
                "usable_as_story_evidence": True,
                "article_title": "Example source memo",
                "numbers": ["10배", "900개"],
                "institutions_regulation": ["World Bank"],
                "risks_reversals": ["리스크"],
                "korea_bridge": ["한국"],
                "everyday_objects": ["오토바이"],
            },
            {
                "source_page_id": "memo_bad",
                "ppt_id": "fixture_deck",
                "normalized_url": "https://example.com/missing",
                "domain": "example.com",
                "access_status": "fetch_failed",
                "memo_quality_status": "soft_error_page",
                "usable_as_story_evidence": False,
                "article_title": "Page not found",
                "numbers": ["404"],
                "institutions_regulation": [],
                "risks_reversals": [],
                "korea_bridge": [],
                "everyday_objects": [],
            }
        ],
    )
    _write_manifest(
        manual_requests,
        [
            {
                "ppt_id": "fixture_deck",
                "first_slide": 3,
                "domain": "bloomberg.com",
                "manual_reason": "likely_paywalled_domain",
                "normalized_url": "https://bloomberg.com/example",
            }
        ],
    )
    _write_manifest(
        lessons,
        [
            {
                "ppt_id": "fixture_deck",
                "main_seed_hook": "Fixture hook becomes a story",
                "story_expansion_path": ["seed source", "numbers", "risk", "Korea bridge"],
            }
        ],
    )
    _write_manifest(
        matches,
        [
            {
                "ppt_id": "fixture_deck",
                "broadcast_label": "likely_used",
                "matched_video_title": "Fixture broadcast",
            }
        ],
    )

    memos = build_ppt_story_arc_memos(
        story_inputs_jsonl=story_inputs,
        slide_visuals_jsonl=slide_visuals,
        source_memos_jsonl=source_memos,
        manual_requests_jsonl=manual_requests,
        lessons_jsonl=lessons,
        matches_jsonl=matches,
        output_jsonl=output,
        report_md=report,
        report_dir=report_dir,
        show_progress=False,
    )

    assert len(memos) == 1
    memo = memos[0]
    assert memo["story_arc_confidence"] == "high"
    assert memo["initial_seed_source"]["domain"] == "example.com"
    assert memo["initial_seed_source"]["usable_as_story_evidence"] is True
    assert memo["numbers_sequence"][0]["numbers"] == ["10배", "900개"]
    assert memo["source_memo_count"] == 1
    assert memo["raw_source_memo_count"] == 2
    assert memo["usable_source_memo_count"] == 1
    assert memo["unusable_source_memo_count"] == 1
    assert memo["source_memo_quality_counts"] == {
        "soft_error_page": 1,
        "usable_public_memo": 1,
    }
    assert "manual_or_auth_source_gap" in memo["jibi_teachable_signals"]
    assert memo["manual_or_auth_gaps"]["manual_request_count"] == 1
    assert memo["broadcast_label"] == "likely_used"
    assert read_jsonl(output)[0]["clean_title"] == "Fixture Deck"
    assert (report_dir / "fixture_deck.md").exists()
    assert "Visual Proof Sequence" in (report_dir / "fixture_deck.md").read_text(
        encoding="utf-8"
    )
    assert "No LLM calls" in report.read_text(encoding="utf-8")


def test_broadcast_match_uses_syukaworld_snapshot(tmp_path) -> None:
    pptx = tmp_path / "f88.pptx"
    _make_pptx(
        pptx,
        "전당포 주식회사",
        "베트남 F88 상장, 오토바이 담보대출, 추심 리스크",
    )
    manifest = tmp_path / "manifest.jsonl"
    _write_manifest(
        manifest,
        [
            {
                "ppt_id": "f88",
                "title": "전당포 주식회사",
                "local_path": str(pptx),
                "folder_date_hint": "2026-05-15",
            }
        ],
    )
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    _make_syuka_db(data_dir / "syuka_ops.db", include_syuka=True)

    records = match_broadcast_usage(
        manifest=manifest,
        syuka_data_dir=data_dir,
        output_jsonl=tmp_path / "matches.jsonl",
        report_md=tmp_path / "matches.md",
    )

    assert records[0]["matched_video_id"] == "vid_f88"
    assert records[0]["broadcast_label"] in {"broadcast_confirmed", "likely_used"}
    assert "syukaworld_channel" in records[0]["matched_signals"]


def test_broadcast_match_excludes_non_syukaworld_channel(tmp_path) -> None:
    pptx = tmp_path / "f88.pptx"
    _make_pptx(
        pptx,
        "전당포 주식회사",
        "베트남 F88 상장, 오토바이 담보대출, 추심 리스크",
    )
    manifest = tmp_path / "manifest.jsonl"
    _write_manifest(
        manifest,
        [
            {
                "ppt_id": "f88",
                "title": "전당포 주식회사",
                "local_path": str(pptx),
                "folder_date_hint": "2026-05-15",
            }
        ],
    )
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    _make_syuka_db(data_dir / "syuka_ops.db", include_syuka=False, include_money=True)

    records = match_broadcast_usage(
        manifest=manifest,
        syuka_data_dir=data_dir,
        output_jsonl=tmp_path / "matches.jsonl",
        report_md=tmp_path / "matches.md",
    )

    assert records[0]["broadcast_label"] == "excluded_non_syukaworld"
    assert records[0]["matched_video_id"] == "vid_money"
