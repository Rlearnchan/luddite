import json
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from luddite.analysis import extract_pptx_style


def _make_sample_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    paragraph = title.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "스타일 테스트 제목"
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = RGBColor(10, 20, 30)
    body = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(6), Inches(1.6))
    body.text_frame.text = "본문 한 줄"
    body.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
    slide.notes_slide.notes_text_frame.text = "[내용] https://example.com/source"
    prs.save(path)


def test_extract_shape_samples_reads_geometry_and_style(tmp_path: Path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    _make_sample_pptx(pptx_path)

    records, meta = extract_pptx_style.extract_shape_samples([pptx_path])

    assert meta["slide_sizes"][0]["width_in"] == 13.333
    assert records
    title_record = next(record for record in records if "스타일 테스트 제목" in record["text"])
    assert title_record["x_cm"] is not None
    assert title_record["font_size_pt"] == 40
    assert title_record["font_family"] == "Malgun Gothic"
    assert title_record["font_color"] == "#0A141E"
    assert title_record["bold"] is True
    assert title_record["notes_url_count"] == 1


def test_aggregate_profile_and_report(tmp_path: Path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    samples_path = tmp_path / "samples.jsonl"
    profile_path = tmp_path / "profile.json"
    report_path = tmp_path / "report.md"
    _make_sample_pptx(pptx_path)

    outputs = extract_pptx_style.run_extraction(
        [pptx_path],
        samples_jsonl=samples_path,
        profile_json=profile_path,
        report_md=report_path,
    )

    assert outputs.samples_jsonl.exists()
    assert outputs.profile_json.exists()
    assert outputs.report_md.exists()
    profile = json.loads(profile_path.read_text(encoding="utf-8"))
    assert profile["shape_count"] >= 2
    assert profile["common_font_sizes"]
    assert "theme_fonts" in profile
    assert "font_resolution" in profile
    assert "Extraction Caveats" in report_path.read_text(encoding="utf-8")
    assert "Theme / Master Font Candidates" in report_path.read_text(encoding="utf-8")


def test_extract_theme_master_fonts_reads_theme_candidates(tmp_path: Path) -> None:
    pptx_path = tmp_path / "sample.pptx"
    _make_sample_pptx(pptx_path)

    font_profile = extract_pptx_style.extract_theme_master_fonts([pptx_path])

    assert font_profile["theme_files"]
    assert "theme_fonts" in font_profile
    assert "font_resolution" in font_profile
    assert isinstance(font_profile["font_resolution"]["resolved_font_candidates"], list)


def test_extract_theme_master_fonts_handles_missing_theme(tmp_path: Path) -> None:
    pptx_path = tmp_path / "minimal.pptx"
    with zipfile.ZipFile(pptx_path, "w") as deck:
        deck.writestr("[Content_Types].xml", "<Types />")

    font_profile = extract_pptx_style.extract_theme_master_fonts([pptx_path])

    assert font_profile["theme_files"] == []
    assert font_profile["theme_fonts"]["major_latin"] is None
    assert font_profile["font_resolution"]["fallback_font"] == "Malgun Gothic"
