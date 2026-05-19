import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

from luddite.agents.piti import build_deck_plan_from_storyline, render_pptx
from luddite.parsers.parse_pptx import parse_presentation


def _storyline() -> dict:
    return {
        "storyline_id": "pptx_storyline_test",
        "title": "PPTX 테스트",
        "sections": [
            {
                "section_title": "도입",
                "slides": [
                    {
                        "slide_no": 1,
                        "slide_type": "title",
                        "headline": "PPTX 테스트",
                        "body": ["수정 가능한 뼈대"],
                        "source_urls": [],
                        "image_urls": [],
                        "notes": "title notes",
                        "needs_source": False,
                        "needs_fact_check": False,
                    },
                    {
                        "slide_no": 2,
                        "slide_type": "explainer",
                        "headline": "출처가 있는 장표",
                        "body": ["본문 메시지 하나"],
                        "source_urls": ["https://example.com/source"],
                        "image_urls": [],
                        "notes": "source note",
                        "needs_source": True,
                        "needs_fact_check": True,
                        "required_before_broadcast": True,
                        "visual_plan": {
                            "kind": "diagram",
                            "description": "diagram candidate",
                            "copyright_risk": False,
                        },
                    },
                    {
                        "slide_no": 3,
                        "slide_type": "production_checklist",
                        "headline": "방송 전 체크",
                        "body": ["추가 확인"],
                        "source_urls": ["https://example.com/check"],
                        "image_urls": [],
                        "notes": "checklist note",
                        "needs_source": False,
                        "needs_fact_check": True,
                    },
                ],
            },
            {
                "section_title": "마무리",
                "slides": [
                    {
                        "slide_no": 4,
                        "slide_type": "closing_question",
                        "headline": "마지막 질문",
                        "body": ["무엇을 볼까"],
                        "source_urls": [],
                        "image_urls": [],
                        "notes": "closing note",
                        "needs_source": False,
                        "needs_fact_check": False,
                    }
                ],
            },
        ],
    }


def _deck_plan() -> dict:
    return build_deck_plan_from_storyline.build_deck_plan_from_storyline(
        _storyline(),
        deck_id="pptx_test",
    )


def _write_style_profile(path: Path) -> None:
    path.write_text(
        json.dumps(
            {
                "slide_size": {"width_in": 13.333, "height_in": 7.5},
                "common_colors": [{"value": "#FF0000", "count": 10}],
                "font_resolution": {"fallback_font": "맑은 고딕"},
                "layout_patterns": {
                    "title": {
                        "x_cm_median": 2.38,
                        "y_cm_median": 3.118,
                        "w_cm_median": 29.106,
                        "h_cm_median": 6.632,
                        "font_size_pt_median": 54.0,
                    },
                    "headline_body": {
                        "x_cm_median": 1.591,
                        "y_cm_median": 3.357,
                        "w_cm_median": 31.092,
                        "h_cm_median": 3.616,
                        "font_size_pt_median": 28.0,
                        "font_color_top": "#FF0000",
                    },
                },
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )


def _runs_for_text(prs: Presentation, text: str):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if text in run.text:
                        yield run


def _first_run_for_text(prs: Presentation, text: str):
    return next(_runs_for_text(prs, text))


def _first_shape_for_text(prs: Presentation, text: str):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and text in shape.text:
                return shape
    raise AssertionError(f"No shape found for text: {text}")


def _first_paragraph_for_text(prs: Presentation, text: str):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                if text in paragraph.text:
                    return paragraph
    raise AssertionError(f"No paragraph found for text: {text}")


def _minimal_deck(slide: dict) -> dict:
    return {
        "deck_id": "format_test",
        "title": "format test",
        "sections": [{"section_no": 1, "section_title": "test", "slides": [slide]}],
        "slides": [slide],
    }


def test_render_deck_plan_to_pptx_creates_editable_file(tmp_path: Path) -> None:
    output_path = tmp_path / "draft.pptx"

    result = render_pptx.render_deck_plan_to_pptx(_deck_plan(), output_path)
    prs = Presentation(str(output_path))

    assert output_path.exists()
    assert result["passed"]
    assert result["slide_count"] == 4
    assert len(prs.slides) == 4
    assert any(
        "PPTX 테스트" in shape.text
        for shape in prs.slides[0].shapes
        if hasattr(shape, "text")
    )


def test_rendered_pptx_notes_preserve_sources_and_flags(tmp_path: Path) -> None:
    output_path = tmp_path / "draft.pptx"
    render_pptx.render_deck_plan_to_pptx(_deck_plan(), output_path)

    parsed = parse_presentation(output_path)
    notes_text = "\n".join(slide["notes"] for slide in parsed["slides"])

    assert "https://example.com/source" in notes_text
    assert "needs_source: true" in notes_text.lower()
    assert "needs_fact_check: true" in notes_text.lower()
    assert "required_before_broadcast: true" in notes_text.lower()
    assert "visual_plan" in notes_text


def test_production_checklist_moves_to_end_and_is_marked_internal(tmp_path: Path) -> None:
    output_path = tmp_path / "draft.pptx"
    render_pptx.render_deck_plan_to_pptx(_deck_plan(), output_path)

    parsed = parse_presentation(output_path)
    last_slide = parsed["slides"][-1]

    assert "[내부 체크리스트]" in last_slide["visible_text"]
    assert "internal checklist" in last_slide["notes"].lower()


def test_render_report_created(tmp_path: Path) -> None:
    deck_path = tmp_path / "deck.json"
    output_path = tmp_path / "draft.pptx"
    report_path = tmp_path / "report.md"
    deck_plan = _deck_plan()
    deck_path.write_text(json.dumps(deck_plan, ensure_ascii=False), encoding="utf-8")
    loaded = json.loads(deck_path.read_text(encoding="utf-8"))
    result = render_pptx.render_deck_plan_to_pptx(loaded, output_path)

    render_pptx.write_render_report(report_path, [result])

    text = report_path.read_text(encoding="utf-8")
    assert "Piti PPTX Render Report" in text
    assert "ready_for_ppt_generation: true (scaffold only)" in text
    assert "Adaptive body font" in text


def test_render_deck_plan_to_pptx_applies_style_profile(tmp_path: Path) -> None:
    output_path = tmp_path / "styled.pptx"
    style_path = tmp_path / "style.json"
    _write_style_profile(style_path)

    style = render_pptx.load_style_profile(style_path)
    result = render_pptx.render_deck_plan_to_pptx(
        _deck_plan(),
        output_path,
        style_profile=style,
    )
    parsed = parse_presentation(output_path)

    assert result["style_profile_loaded"] is True
    assert result["applied_font_family"] == "맑은 고딕"
    assert result["applied_layout_count"] >= 1
    assert result["adaptive_font_applied"] is True
    assert result["visual_placeholder_shortened"] is True
    assert result["slides_with_text_placeholder_overlap"] == []
    assert result["proof_text_overlap_count"] == 0
    assert result["headline_red_count"] > 0
    assert result["body_black_count"] > 0
    assert result["parse_back_slide_count"] == 4
    assert parsed["slide_count"] == 4
    notes_text = "\n".join(slide["notes"] for slide in parsed["slides"])
    assert "style_profile_loaded: True" in notes_text
    assert "visual_plan" in notes_text
    assert "proof_object" in notes_text
    visible_text = "\n".join(slide["visible_text"] for slide in parsed["slides"])
    assert "diagram candidate" not in visible_text
    assert "draft skeleton" not in visible_text
    assert "needs_fact_check" not in visible_text
    assert result["screen_footer_hidden_count"] == 4

    prs = Presentation(str(output_path))
    headline_run = _first_run_for_text(prs, "출처가 있는 장표")
    headline_shape = _first_shape_for_text(prs, "출처가 있는 장표")
    body_shape = _first_shape_for_text(prs, "본문 메시지 하나")
    body_run = _first_run_for_text(prs, "본문 메시지 하나")
    body_paragraph = _first_paragraph_for_text(prs, "본문 메시지 하나")
    headline_paragraph = _first_paragraph_for_text(prs, "출처가 있는 장표")
    assert headline_run.font.size.pt == 28
    assert headline_run.font.color.rgb == RGBColor(255, 0, 0)
    assert headline_run.font.bold is False
    assert headline_paragraph.line_spacing is None
    assert 1.4 <= headline_shape.left.cm <= 1.8
    assert 0.8 <= headline_shape.top.cm <= 1.2
    assert body_run.font.color.rgb == RGBColor(0, 0, 0)
    assert body_shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert body_paragraph.line_spacing == 1.5
    assert result["headline_bold_count"] == 0
    assert result["headline_nonbold_count"] > 0
    assert result["body_vertical_middle_count"] > 0
    assert result["body_vertical_top_count"] == 0
    assert result["debug_label_visible_count"] == 0
    assert result["actual_body_black_count"] > 0
    assert result["body_line_spacing_applied_count"] > 0
    assert result["body_line_spacing_value"] == 1.5
    assert result["body_line_spacing_missing_count"] == 0
    assert result["parse_back_line_spacing_1_5_count"] > 0


def test_quote_bilingual_mode_uses_english_black_and_korean_red(tmp_path: Path) -> None:
    output_path = tmp_path / "quote.pptx"
    style_path = tmp_path / "style.json"
    _write_style_profile(style_path)
    slide = {
        "slide_no": 1,
        "layout_type": "quote",
        "slide_type": "quote",
        "headline": "영한 교차 인용",
        "body": [
            "They collectively held a record $160bn",
            "11월 말 기준 개인들이 보유한 미국 주식 규모는 1,600억달러",
        ],
        "source_urls": ["https://example.com/quote"],
        "image_urls": [],
        "speaker_notes": "quote notes",
        "visual_plan": {"kind": "none"},
    }

    render_pptx.render_deck_plan_to_pptx(
        _minimal_deck(slide),
        output_path,
        style_profile=render_pptx.load_style_profile(style_path),
    )

    prs = Presentation(str(output_path))
    english_run = _first_run_for_text(prs, "They collectively")
    korean_run = _first_run_for_text(prs, "11월 말 기준")
    english_paragraph = _first_paragraph_for_text(prs, "They collectively")
    korean_paragraph = _first_paragraph_for_text(prs, "11월 말 기준")
    assert english_run.font.size.pt == 28
    assert english_run.font.color.rgb == RGBColor(0, 0, 0)
    assert korean_run.font.size.pt == 28
    assert korean_run.font.color.rgb == RGBColor(255, 0, 0)
    assert english_paragraph.line_spacing == 1.5
    assert korean_paragraph.line_spacing == 1.5


def test_chart_table_placeholder_uses_chart_typography(tmp_path: Path) -> None:
    output_path = tmp_path / "chart.pptx"
    style_path = tmp_path / "style.json"
    _write_style_profile(style_path)
    slide = {
        "slide_no": 1,
        "layout_type": "chart_placeholder",
        "slide_type": "data",
        "headline": "한국인의 사랑을 듬뿍 받는 테슬라",
        "body": [
            "미국주식 보관금액 상위 10종목",
            "테슬라 264",
            "차트 설명은 노트로 보낸다",
            "엔비디아 195",
        ],
        "source_urls": ["https://example.com/seibro"],
        "image_urls": [],
        "speaker_notes": "chart notes",
        "visual_plan": {"kind": "chart_candidate"},
    }

    result = render_pptx.render_deck_plan_to_pptx(
        _minimal_deck(slide),
        output_path,
        style_profile=render_pptx.load_style_profile(style_path),
    )

    prs = Presentation(str(output_path))
    title_run = _first_run_for_text(prs, "미국주식 보관금액")
    headline_run = _first_run_for_text(prs, "한국인의 사랑")
    data_run = _first_run_for_text(prs, "테슬라 264")
    data_paragraph = _first_paragraph_for_text(prs, "테슬라 264")
    source_run = _first_run_for_text(prs, "(출처:")
    source_paragraph = _first_paragraph_for_text(prs, "(출처:")
    visible_text = "\n".join(
        slide["visible_text"] for slide in parse_presentation(output_path)["slides"]
    )
    assert result["chart_table_style_applied_count"] == 1
    assert result["chart_table_skeleton_count"] == 1
    assert result["chart_body_text_leak_count"] == 0
    assert result["proof_object_type_counts"] == {"chart": 1}
    assert headline_run.font.color.rgb == RGBColor(255, 0, 0)
    assert headline_run.font.bold is False
    assert title_run.font.size.pt == 28
    assert title_run.font.bold is True
    assert title_run.font.underline is True
    assert title_run.font.color.rgb == RGBColor(0, 0, 0)
    assert data_run.font.size.pt == 18
    assert data_run.font.bold is True
    assert data_paragraph.line_spacing is None
    assert source_run.font.size.pt == 20
    assert source_run.font.underline is True
    assert source_paragraph.line_spacing is None
    assert result["chart_title_bold_underline_count"] == 1
    assert "차트 설명은 노트로 보낸다" not in visible_text


def test_section_title_is_not_forced_red_by_headline_rule(tmp_path: Path) -> None:
    output_path = tmp_path / "section.pptx"
    style_path = tmp_path / "style.json"
    _write_style_profile(style_path)
    slide = {
        "slide_no": 1,
        "layout_type": "section_title",
        "slide_type": "section_title",
        "headline": "섹션 제목은 별도 규칙",
        "body": [],
        "source_urls": [],
        "image_urls": [],
        "speaker_notes": "section notes",
        "visual_plan": {"kind": "none"},
    }

    render_pptx.render_deck_plan_to_pptx(
        _minimal_deck(slide),
        output_path,
        style_profile=render_pptx.load_style_profile(style_path),
    )

    prs = Presentation(str(output_path))
    section_run = _first_run_for_text(prs, "섹션 제목")
    assert section_run.font.bold is True
    assert section_run.font.color.rgb == RGBColor(24, 28, 35)


def test_manual_placeholder_hidden_and_image_left_counted(tmp_path: Path) -> None:
    output_path = tmp_path / "image.pptx"
    style_path = tmp_path / "style.json"
    _write_style_profile(style_path)
    slides = [
        {
            "slide_no": 1,
            "layout_type": "headline_body",
            "slide_type": "explainer",
            "headline": "수동 이미지는 화면에 숨긴다",
            "body": ["본문은 검은색"],
            "source_urls": [],
            "image_urls": [],
            "speaker_notes": "manual visual",
            "visual_plan": {"kind": "manual", "description": "manual insert"},
        },
        {
            "slide_no": 2,
            "layout_type": "image_placeholder",
            "slide_type": "image",
            "headline": "이미지는 왼쪽에 둔다",
            "body": ["오른쪽에 해석을 둔다"],
            "source_urls": [],
            "image_urls": [],
            "speaker_notes": "image visual",
            "visual_plan": {"kind": "screenshot_candidate"},
        },
    ]
    deck = {"deck_id": "image_test", "sections": [{"slides": slides}], "slides": slides}

    result = render_pptx.render_deck_plan_to_pptx(
        deck,
        output_path,
        style_profile=render_pptx.load_style_profile(style_path),
    )
    prs = Presentation(str(output_path))
    parsed = parse_presentation(output_path)
    visible_text = "\n".join(slide["visible_text"] for slide in parsed["slides"])
    notes_text = "\n".join(slide["notes"] for slide in parsed["slides"])

    assert result["manual_placeholder_hidden_count"] == 1
    assert result["image_left_layout_count"] == 1
    assert result["proof_object_type_counts"] == {"screenshot": 1}
    assert "[수동" not in visible_text
    assert "manual_placeholder_hidden: True" in notes_text
    assert "proof_object_type: screenshot" in notes_text
    right_text_shape = _first_shape_for_text(prs, "오른쪽에 해석을 둔다")
    assert right_text_shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
    assert result["editor_instruction_blue_count"] >= 1
    assert result["editor_instruction_screen_count"] >= 1


def test_article_quote_proof_object_reserves_left_area(tmp_path: Path) -> None:
    output_path = tmp_path / "article_quote.pptx"
    style_path = tmp_path / "style.json"
    _write_style_profile(style_path)
    slide = {
        "slide_no": 1,
        "layout_type": "quote",
        "slide_type": "quote",
        "headline": "기사 인용으로 보여준다",
        "body": ["They collectively held a record $160bn", "한국어 해석은 빨간색"],
        "source_urls": ["https://example.com/article"],
        "image_urls": [],
        "speaker_notes": "article quote notes",
        "visual_plan": {"kind": "manual", "description": "article quote proof"},
    }

    result = render_pptx.render_deck_plan_to_pptx(
        _minimal_deck(slide),
        output_path,
        style_profile=render_pptx.load_style_profile(style_path),
    )
    parsed = parse_presentation(output_path)
    visible_text = "\n".join(item["visible_text"] for item in parsed["slides"])
    notes_text = "\n".join(item["notes"] for item in parsed["slides"])

    assert result["article_quote_skeleton_count"] == 1
    assert result["article_quote_count"] == 1
    assert result["proof_object_type_counts"] == {"article_quote": 1}
    assert result["proof_object_area_reserved_count"] == 1
    assert result["proof_text_overlap_count"] == 0
    assert "[인용]" in visible_text
    assert "Example" in visible_text
    assert "example.com" not in visible_text
    assert "proof_object_type: article_quote" in notes_text
    editor_run = _first_run_for_text(prs := Presentation(str(output_path)), "[인용]")
    source_run = _first_run_for_text(prs, "Example")
    assert editor_run.font.color.rgb == RGBColor(0, 112, 192)
    assert source_run.font.color.rgb == RGBColor(0, 0, 0)
    assert result["visible_url_count"] == 0


def test_text_only_and_proof_object_metrics_are_reported(tmp_path: Path) -> None:
    output_path = tmp_path / "mixed.pptx"
    report_path = tmp_path / "report.md"
    style_path = tmp_path / "style.json"
    _write_style_profile(style_path)
    slides = [
        {
            "slide_no": 1,
            "layout_type": "headline_body",
            "slide_type": "explainer",
            "headline": "텍스트만 있는 장표",
            "body": ["짧은 해석만 남긴다"],
            "source_urls": [],
            "image_urls": [],
            "speaker_notes": "text only",
            "visual_plan": {"kind": "none"},
        },
        {
            "slide_no": 2,
            "layout_type": "question",
            "slide_type": "bridge",
            "headline": "도식이 필요한 장표",
            "body": ["오른쪽에 해석을 둔다"],
            "source_urls": [],
            "image_urls": [],
            "speaker_notes": "diagram proof",
            "visual_plan": {"kind": "diagram"},
        },
    ]
    deck = {"deck_id": "proof_metrics", "sections": [{"slides": slides}], "slides": slides}

    result = render_pptx.render_deck_plan_to_pptx(
        deck,
        output_path,
        style_profile=render_pptx.load_style_profile(style_path),
    )
    render_pptx.write_render_report(report_path, [result])
    report_text = report_path.read_text(encoding="utf-8")

    assert result["proof_object_slide_count"] == 1
    assert result["proof_object_type_counts"] == {"diagram": 1}
    assert result["text_only_slide_count"] == 1
    assert result["proof_object_required_but_missing_count"] == 0
    assert result["image_left_layout_count"] == 1
    assert result["diagram_skeleton_count"] == 1
    assert "proof_object_slide_count" in report_text
    assert "proof_object_type_counts" in report_text
    visible_text = "\n".join(
        slide["visible_text"] for slide in parse_presentation(output_path)["slides"]
    )
    assert "기존 방식" in visible_text
    assert "AI 즉답" in visible_text


def test_source_backed_manual_slide_uses_source_card_template(tmp_path: Path) -> None:
    output_path = tmp_path / "source_card.pptx"
    style_path = tmp_path / "style.json"
    _write_style_profile(style_path)
    slide = {
        "slide_no": 1,
        "layout_type": "headline_body",
        "slide_type": "explainer",
        "headline": "BBC가 던진 질문",
        "body": ["출처가 있는 주장은 카드로 보여준다", "긴 설명은 노트에 둔다"],
        "source_urls": ["https://www.bbc.com/news/example"],
        "image_urls": [],
        "speaker_notes": "source-backed manual slide",
        "visual_plan": {"kind": "manual", "description": "manual source card"},
    }

    result = render_pptx.render_deck_plan_to_pptx(
        _minimal_deck(slide),
        output_path,
        style_profile=render_pptx.load_style_profile(style_path),
    )
    parsed = parse_presentation(output_path)
    visible_text = "\n".join(item["visible_text"] for item in parsed["slides"])
    notes_text = "\n".join(item["notes"] for item in parsed["slides"])

    assert result["proof_object_type_counts"] == {"source_card": 1}
    assert result["layout_template_counts"] == {"source_card_or_article_quote": 1}
    assert result["source_card_or_article_quote_count"] == 1
    assert result["source_card_count"] == 1
    assert result["article_quote_count"] == 0
    assert result["source_card_repeated_headline_count"] == 0
    assert result["text_only_slide_count"] == 0
    assert result["text_only_slide_count_before_after"] == {"before": 1, "after": 0}
    assert result["proof_object_slide_count_before_after"] == {"before": 0, "after": 1}
    assert result["source_backed_text_only_should_have_card_count"] == 0
    assert "[출처]" in visible_text
    assert "BBC" in visible_text
    assert visible_text.count("BBC가 던진 질문") == 1
    assert "proof_object_type: source_card" in notes_text


def test_reference_layout_report_exposes_template_metrics(tmp_path: Path) -> None:
    output_path = tmp_path / "templates.pptx"
    report_path = tmp_path / "report.md"
    style_path = tmp_path / "style.json"
    _write_style_profile(style_path)
    slides = [
        {
            "slide_no": 1,
            "layout_type": "chart_placeholder",
            "slide_type": "data",
            "headline": "차트가 주인공",
            "body": ["차트 제목", "A 10", "B 7"],
            "source_urls": ["https://example.com/chart"],
            "image_urls": [],
            "speaker_notes": "chart",
            "visual_plan": {"kind": "chart_candidate"},
        },
        {
            "slide_no": 2,
            "layout_type": "question",
            "slide_type": "bridge",
            "headline": "텍스트 자체가 장면",
            "body": ["그래서 무엇을 물어야 하나"],
            "source_urls": [],
            "image_urls": [],
            "speaker_notes": "calculation",
            "visual_plan": {"kind": "none"},
        },
        {
            "slide_no": 3,
            "layout_type": "image_placeholder",
            "slide_type": "image",
            "headline": "이미지 왼쪽",
            "body": ["오른쪽 해석"],
            "source_urls": [],
            "image_urls": [],
            "speaker_notes": "image",
            "visual_plan": {"kind": "screenshot_candidate"},
        },
    ]
    deck = {"deck_id": "template_metrics", "sections": [{"slides": slides}], "slides": slides}

    result = render_pptx.render_deck_plan_to_pptx(
        deck,
        output_path,
        style_profile=render_pptx.load_style_profile(style_path),
    )
    render_pptx.write_render_report(report_path, [result])
    report_text = report_path.read_text(encoding="utf-8")

    assert result["layout_template_counts"] == {
        "chart_table_reference": 1,
        "text_only_calculation": 1,
        "image_left_quote_right": 1,
    }
    assert result["chart_table_reference_count"] == 1
    assert result["image_left_quote_right_count"] == 1
    assert result["text_only_calculation_count"] == 1
    assert "layout_template_counts" in report_text
    assert "source_backed_text_only_should_have_card_count" in report_text
