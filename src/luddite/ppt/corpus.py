"""Read-only Syukaworld PPT corpus inventory, extraction, and reports."""

from __future__ import annotations

import hashlib
import re
import unicodedata
from collections import Counter, defaultdict
from datetime import UTC, date, datetime
from pathlib import Path
from typing import Annotated, Any
from urllib.parse import urlsplit

import typer
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from rich.console import Console

from luddite import paths
from luddite.parsers.parse_pptx import parse_presentation
from luddite.utils.jsonl import read_jsonl, write_jsonl
from luddite.utils.urls import canonicalize_url

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PPT_MIME = "application/vnd.ms-powerpoint"
GOOGLE_SLIDES_MIME = "application/vnd.google-apps.presentation"
PRESENTATION_MIMES = {PPTX_MIME, PPT_MIME, GOOGLE_SLIDES_MIME}
PRESENTATION_EXTENSIONS = {".pptx", ".ppt"}

PPT_CORPUS_DATA_DIR = paths.DATA_DIR / "ppt_corpus"
PPT_CORPUS_EXTRACTED_DIR = PPT_CORPUS_DATA_DIR / "extracted"
PPT_CORPUS_REPORT_DIR = paths.REPORTS_DIR / "ppt_corpus"
PPT_CORPUS_CONTACT_SHEETS_DIR = paths.OUTPUTS_DIR / "contact_sheets" / "ppt_corpus"
PPT_CORPUS_EVAL_DIR = paths.OUTPUTS_DIR / "eval"
INSIGHTS_DIR = paths.DOCS_DIR / "insights"

LEGACY_PPT_LEARNING_DATA_DIR = paths.DATA_DIR / "ppt_learning"
LEGACY_PPT_LEARNING_OUTPUT_DIR = paths.OUTPUTS_DIR / "ppt_learning"

DRIVE_INVENTORY_JSONL = PPT_CORPUS_DATA_DIR / "drive_ppt_inventory.jsonl"
DECK_MANIFEST_JSONL = PPT_CORPUS_EXTRACTED_DIR / "deck_manifest.jsonl"
SLIDES_JSONL = PPT_CORPUS_EXTRACTED_DIR / "slides.jsonl"
LINKS_JSONL = PPT_CORPUS_EXTRACTED_DIR / "links.jsonl"
MEDIA_MANIFEST_JSONL = PPT_CORPUS_EXTRACTED_DIR / "media_manifest.jsonl"

DRIVE_INVENTORY_REPORT_MD = PPT_CORPUS_REPORT_DIR / f"drive_ppt_inventory_{date.today()}.md"
EXTRACTION_QUALITY_REPORT_MD = PPT_CORPUS_REPORT_DIR / f"extraction_quality_{date.today()}.md"

JIBI_INSIGHTS_MD = INSIGHTS_DIR / "jibi_ppt_corpus_selection_lessons.md"
ANNY_ARCHETYPES_MD = INSIGHTS_DIR / "anny_story_archetypes_from_ppt_corpus.md"
PITI_STYLE_GUIDE_MD = INSIGHTS_DIR / "piti_syukaworld_slide_style_guide.md"

JIBI_REPORT_MD = PPT_CORPUS_REPORT_DIR / "03_jibi_selection_lessons.md"
ANNY_REPORT_MD = PPT_CORPUS_REPORT_DIR / "04_anny_story_archetypes.md"
PITI_REPORT_MD = PPT_CORPUS_REPORT_DIR / "05_piti_slide_style.md"
RECOMMENDED_CODE_CHANGES_MD = PPT_CORPUS_REPORT_DIR / "07_recommended_code_changes.md"

drive_manifest_app = typer.Typer(no_args_is_help=False)
inventory_app = typer.Typer(no_args_is_help=False)
extract_slides_app = typer.Typer(no_args_is_help=False)
quality_report_app = typer.Typer(no_args_is_help=False)
insight_reports_app = typer.Typer(no_args_is_help=False)
console = Console()

KNOWN_RESEARCHERS = {
    "김동찬",
    "김성원",
    "김현수",
    "배형찬",
    "이동현",
    "이상민",
    "정재민",
    "최윤영",
    "한상혁",
}

POSITIVE_JIBI_LESSONS = [
    "known_brand_plus_structural_shift",
    "daily_life_problem_to_system",
    "hidden_cost_or_hidden_owner",
    "absurd_institution_or_rule",
    "official_stat_with_counterintuitive_question",
    "overseas_case_with_korea_bridge",
    "old_story_reactivated_by_current_news",
    "market_price_as_social_signal",
]

NEGATIVE_JIBI_LESSONS = [
    "single_company_ir_without_audience_bridge",
    "broad_ai_discourse_without_use_case",
    "sports_primary_without_business_bridge",
    "policy_release_without_scene",
    "foreign_explainer_without_korea_bridge",
    "title_hook_without_mechanism",
]

ANNY_ARCHETYPES = [
    "생활 불편 -> 숨은 비용 -> 제도/기업의 손익계산서",
    "이상한 해외 사례 -> 알고 보니 한국도 같은 구조",
    "유명 브랜드/인물 -> 산업 구조 변화",
    "숫자 하나 -> 통계의 함정 -> 진짜 문제",
    "옛날 이야기 -> 현재 뉴스 -> 반복되는 구조",
    "정책 발표 -> 현장에서는 왜 다르게 작동하나",
    "시장 가격 -> 사람들이 아직 못 본 위험 신호",
]

SLIDE_PATTERN_TAXONOMY = [
    "opening_question",
    "big_number",
    "news_capture",
    "official_document_capture",
    "timeline",
    "comparison_table",
    "map_or_geography",
    "market_chart",
    "before_after",
    "quote_or_claim",
    "meme_or_visual_joke",
    "section_divider",
    "mechanism_diagram",
    "summary_takeaway",
]


def _now_iso() -> str:
    return datetime.now(UTC).isoformat()


def _repo_relative(path: Path) -> str:
    try:
        return str(path.relative_to(paths.REPO_ROOT))
    except ValueError:
        return str(path)


def _resolve_local_path(value: str) -> Path:
    if not value:
        return Path("")
    path = Path(value)
    if path.is_absolute():
        return path
    return paths.REPO_ROOT / path


def _read_jsonl_if_exists(path: Path) -> list[dict[str, Any]]:
    if not path.exists():
        return []
    return read_jsonl(path)


def _compact(text: Any) -> str:
    text = unicodedata.normalize("NFC", str(text or ""))
    text = text.replace("\xa0", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _slug(text: str, *, limit: int = 80) -> str:
    normalized = _compact(text).lower()
    normalized = re.sub(r"[^\w가-힣]+", "_", normalized)
    normalized = re.sub(r"_+", "_", normalized).strip("_")
    return normalized[:limit] or "deck"


def _stable_id(prefix: str, *parts: Any) -> str:
    raw = "|".join(_compact(part) for part in parts if _compact(part))
    digest = hashlib.sha1(raw.encode("utf-8")).hexdigest()[:12] if raw else "0" * 12
    return f"{prefix}_{digest}"


def _path_modified_iso(path: Path) -> str:
    try:
        return datetime.fromtimestamp(path.stat().st_mtime, UTC).isoformat()
    except OSError:
        return ""


def _mime_from_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return PPTX_MIME
    if suffix == ".ppt":
        return PPT_MIME
    return ""


def _domain(url: str) -> str:
    try:
        return urlsplit(url).netloc.lower().removeprefix("www.")
    except ValueError:
        return ""


def infer_author_and_topic(title: str) -> tuple[str, str]:
    """Infer researcher and clean topic from common Syukaworld deck names."""
    stem = Path(_compact(title)).stem
    parts = [part.strip() for part in re.split(r"[_-]", stem) if part.strip()]
    if len(parts) >= 2 and parts[-1] in KNOWN_RESEARCHERS:
        return parts[-1], _compact("_".join(parts[:-1]))
    if len(parts) >= 2 and parts[0] in KNOWN_RESEARCHERS:
        return parts[0], _compact("_".join(parts[1:]))

    paren_author = re.search(r"\(([^()]+)\)$", stem)
    if paren_author and paren_author.group(1).strip() in KNOWN_RESEARCHERS:
        author = paren_author.group(1).strip()
        topic = _compact(stem[: paren_author.start()])
        return author, topic

    return "", stem


def infer_deck_type(title: str, *, path_hint: str = "", mime_type: str = "") -> str:
    text = f"{title} {path_hint}".lower()
    korean_text = f"{title} {path_hint}"
    if any(token in text for token in ["template", "sample"]) or any(
        token in korean_text for token in ["템플릿", "양식", "예시 자료", "예시자료"]
    ):
        return "template"
    if any(token in text for token in ["manual", "guide", "playbook"]) or any(
        token in korean_text for token in ["매뉴얼", "가이드", "작업법"]
    ):
        return "manual"
    if any(token in korean_text for token in ["국내 이슈 정리", "이슈 정리", "방송용"]):
        return "daily_deck"
    if mime_type in PRESENTATION_MIMES or Path(path_hint).suffix.lower() in PRESENTATION_EXTENSIONS:
        return "topic_deck"
    return "unknown"


def _needs_download(mime_type: str, local_path: str) -> bool:
    if mime_type == GOOGLE_SLIDES_MIME:
        return True
    return not bool(local_path)


def _skip_reason(mime_type: str, local_path: str) -> str:
    if mime_type and mime_type not in PRESENTATION_MIMES:
        return "not_presentation"
    if mime_type == GOOGLE_SLIDES_MIME and not local_path:
        return "google_slides_export_required"
    return ""


def canonical_inventory_record(row: dict[str, Any]) -> dict[str, Any]:
    title = _compact(
        row.get("title")
        or row.get("name")
        or row.get("ppt_title")
        or Path(str(row.get("local_path") or "")).stem
    )
    local_path = str(row.get("local_path") or row.get("resolved_local_path") or "")
    inferred_author, inferred_topic = infer_author_and_topic(title)
    inferred_author = _compact(
        row.get("researcher") or row.get("inferred_author") or inferred_author
    )
    inferred_topic = _compact(row.get("inferred_topic_title") or inferred_topic or title)
    mime_type = _compact(
        row.get("mime_type") or row.get("mimeType") or _mime_from_path(Path(local_path))
    )
    drive_file_id = _compact(
        row.get("drive_file_id")
        or row.get("file_id")
        or row.get("id")
        or row.get("ppt_id")
        or _stable_id("drive_file", local_path, title)
    )
    deck_id = _compact(
        row.get("deck_id")
        or row.get("ppt_id")
        or f"deck_{_slug(inferred_topic)}_{hashlib.sha1(drive_file_id.encode()).hexdigest()[:8]}"
    )
    path_hint = _compact(row.get("drive_path_hint") or row.get("path") or local_path)
    needs_download = bool(row.get("needs_download", _needs_download(mime_type, local_path)))
    skip_reason = _compact(row.get("skip_reason") or _skip_reason(mime_type, local_path))
    return {
        "deck_id": deck_id,
        "drive_file_id": drive_file_id,
        "title": title,
        "mime_type": mime_type,
        "created_at": _compact(row.get("created_at") or row.get("createdTime") or ""),
        "modified_at": _compact(row.get("modified_at") or row.get("modifiedTime") or ""),
        "url": _compact(row.get("url") or row.get("webViewLink") or ""),
        "inferred_author": inferred_author,
        "inferred_topic_title": inferred_topic,
        "inferred_deck_type": infer_deck_type(
            inferred_topic or title, path_hint=path_hint, mime_type=mime_type
        ),
        "needs_download": needs_download,
        "skip_reason": skip_reason,
        "local_path": local_path,
        "resolved_local_path": str(_resolve_local_path(local_path)) if local_path else "",
        "source_root": _compact(row.get("source_root") or ""),
        "drive_path_hint": path_hint,
        "file_size_bytes": int(row.get("file_size_bytes") or row.get("size") or 0),
        "inventory_status": "skipped" if skip_reason == "not_presentation" else "ready",
        "inventory_created_at": _now_iso(),
    }


def _normalize_inventory_rows_with_stats(
    rows: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], Counter[str]]:
    records: list[dict[str, Any]] = []
    stats: Counter[str] = Counter()
    seen_drive_ids: set[str] = set()
    seen_local_paths: set[str] = set()
    for row in rows:
        record = canonical_inventory_record(row)
        drive_file_id = str(record.get("drive_file_id") or "")
        local_path = str(record.get("local_path") or "")
        if drive_file_id and drive_file_id in seen_drive_ids:
            stats["duplicate_drive_file_id"] += 1
            continue
        if local_path and local_path in seen_local_paths:
            stats["duplicate_local_path"] += 1
            continue
        if drive_file_id:
            seen_drive_ids.add(drive_file_id)
        if local_path:
            seen_local_paths.add(local_path)
        stats[str(record.get("inferred_deck_type") or "unknown")] += 1
        if record.get("needs_download"):
            stats["needs_download"] += 1
        if record.get("skip_reason"):
            stats[f"skip:{record['skip_reason']}"] += 1
        records.append(record)
    return records, stats


def normalize_inventory_rows(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    records, _stats = _normalize_inventory_rows_with_stats(rows)
    return records


def _local_inventory_row(path: Path, root: Path, source_root: str) -> dict[str, Any]:
    try:
        relative = path.relative_to(root)
    except ValueError:
        relative = path
    author, topic = infer_author_and_topic(path.stem)
    return {
        "drive_file_id": _stable_id("local_drive", source_root, relative),
        "title": topic or path.stem,
        "mime_type": _mime_from_path(path),
        "created_at": "",
        "modified_at": _path_modified_iso(path),
        "url": "",
        "researcher": author,
        "local_path": _repo_relative(path),
        "source_root": source_root,
        "drive_path_hint": str(relative),
        "file_size_bytes": path.stat().st_size if path.exists() else 0,
    }


def _fallback_root(primary: Path, legacy: Path) -> Path:
    if primary.exists():
        return primary
    return legacy


def build_drive_manifest(
    *,
    latest_root: Path = PPT_CORPUS_DATA_DIR / "drive_raw" / "latest",
    past_root: Path = PPT_CORPUS_DATA_DIR / "drive_raw" / "past",
    output_jsonl: Path = DRIVE_INVENTORY_JSONL,
    report_md: Path = DRIVE_INVENTORY_REPORT_MD,
) -> list[dict[str, Any]]:
    latest_root = _fallback_root(latest_root, LEGACY_PPT_LEARNING_DATA_DIR / "drive_raw" / "latest")
    past_root = _fallback_root(past_root, LEGACY_PPT_LEARNING_DATA_DIR / "drive_raw" / "past")

    raw_rows: list[dict[str, Any]] = []
    non_presentation_files: list[str] = []
    for source_root, root in [("latest", latest_root), ("past", past_root)]:
        if not root.exists():
            continue
        for path in sorted(root.rglob("*")):
            if not path.is_file():
                continue
            if path.suffix.lower() in PRESENTATION_EXTENSIONS:
                raw_rows.append(_local_inventory_row(path, root, source_root))
            else:
                non_presentation_files.append(_repo_relative(path))

    records, stats = _normalize_inventory_rows_with_stats(raw_rows)
    write_jsonl(output_jsonl, records)
    report_md.parent.mkdir(parents=True, exist_ok=True)
    report_md.write_text(
        _inventory_report_markdown(
            records=records,
            stats=stats,
            source=f"latest={_repo_relative(latest_root)}, past={_repo_relative(past_root)}",
            output_jsonl=output_jsonl,
            non_presentation_files=non_presentation_files,
        ),
        encoding="utf-8",
    )
    return records


def build_inventory(
    *,
    manifest: Path = DRIVE_INVENTORY_JSONL,
    output_jsonl: Path = DRIVE_INVENTORY_JSONL,
    report_md: Path = DRIVE_INVENTORY_REPORT_MD,
) -> list[dict[str, Any]]:
    if (
        not manifest.exists()
        and manifest == DRIVE_INVENTORY_JSONL
        and (LEGACY_PPT_LEARNING_DATA_DIR / "drive_ppts.jsonl").exists()
    ):
        manifest = LEGACY_PPT_LEARNING_DATA_DIR / "drive_ppts.jsonl"
    raw_rows = _read_jsonl_if_exists(manifest)
    records, stats = _normalize_inventory_rows_with_stats(raw_rows)
    write_jsonl(output_jsonl, records)
    report_md.parent.mkdir(parents=True, exist_ok=True)
    report_md.write_text(
        _inventory_report_markdown(
            records=records,
            stats=stats,
            source=_repo_relative(manifest),
            output_jsonl=output_jsonl,
            non_presentation_files=[],
        ),
        encoding="utf-8",
    )
    return records


def _inventory_report_markdown(
    *,
    records: list[dict[str, Any]],
    stats: Counter[str],
    source: str,
    output_jsonl: Path,
    non_presentation_files: list[str],
) -> str:
    type_counts = Counter(str(row.get("inferred_deck_type") or "unknown") for row in records)
    mime_counts = Counter(str(row.get("mime_type") or "unknown") for row in records)
    lines = [
        "# Syukaworld PPT Corpus Inventory",
        "",
        "Read-only inventory report. Google Drive originals, Sheets, and production "
        "DBs were not modified.",
        "",
        "## Summary",
        "",
        f"- source: `{source}`",
        f"- output_jsonl: `{_repo_relative(output_jsonl)}`",
        f"- deck_count: {len(records)}",
        f"- needs_download_count: {sum(1 for row in records if row.get('needs_download'))}",
        f"- duplicate_drive_file_id_skipped: {stats.get('duplicate_drive_file_id', 0)}",
        f"- duplicate_local_path_skipped: {stats.get('duplicate_local_path', 0)}",
        f"- non_presentation_file_count: {len(non_presentation_files)}",
        "",
        "## Deck Types",
        "",
    ]
    lines.extend(f"- {deck_type}: {count}" for deck_type, count in sorted(type_counts.items()))
    lines.extend(["", "## MIME Types", ""])
    lines.extend(f"- {mime}: {count}" for mime, count in sorted(mime_counts.items()))
    lines.extend(
        ["", "## Sample Decks", "", "| deck | type | author | local |", "| --- | --- | --- | --- |"]
    )
    for row in records[:25]:
        lines.append(
            "| "
            f"{row.get('inferred_topic_title') or row.get('title')} | "
            f"{row.get('inferred_deck_type')} | "
            f"{row.get('inferred_author') or ''} | "
            f"`{row.get('local_path') or ''}` |"
        )
    if non_presentation_files:
        lines.extend(["", "## Excluded Non-Presentation Files", ""])
        lines.extend(f"- `{item}`" for item in non_presentation_files[:200])
        if len(non_presentation_files) > 200:
            lines.append(f"- ... {len(non_presentation_files) - 200} more")
    return "\n".join(lines) + "\n"


def _shape_counts_by_slide(path: Path) -> dict[int, dict[str, int]]:
    counts: dict[int, dict[str, int]] = {}
    presentation = Presentation(str(path))
    for slide_no, slide in enumerate(presentation.slides, start=1):
        slide_counts = {"image_count": 0, "chart_count": 0, "table_count": 0}
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                slide_counts["table_count"] += 1
            if getattr(shape, "has_chart", False):
                slide_counts["chart_count"] += 1
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_counts["image_count"] += 1
        counts[slide_no] = slide_counts
    return counts


def _possible_slide_role(
    *,
    slide_no: int,
    slide_count: int,
    headline: str,
    visible_text: str,
    url_count: int,
    image_count: int,
    chart_count: int,
    table_count: int,
    slide_type: str,
) -> str:
    compact = visible_text.replace("\n", " ")
    if slide_no == 1:
        return "opening_question" if "?" in compact else "opening_hook"
    if chart_count:
        return "market_chart"
    if table_count:
        return "comparison_table"
    if re.search(r"[0-9][0-9,.]*(%|원|달러|배|조|억|만|개)", compact):
        return "big_number"
    if url_count and image_count:
        return "news_capture"
    if image_count and len(compact) < 140:
        return "meme_or_visual_joke"
    if slide_type == "section_title":
        return "section_divider"
    if slide_no >= max(1, int(slide_count * 0.85)):
        return "summary_takeaway"
    if any(token in compact for token in ["비교", "반면", "VS", "vs"]):
        return "before_after"
    if any(token in compact for token in ["“", "”", '"', "발언", "말했다"]):
        return "quote_or_claim"
    return "mechanism_diagram" if image_count else "explainer"


def _visual_density(text: str, *, image_count: int, chart_count: int, table_count: int) -> str:
    line_count = len([line for line in text.splitlines() if line.strip()])
    visual_count = image_count + chart_count + table_count
    if visual_count >= 2 or line_count >= 9 or len(text) >= 500:
        return "high"
    if visual_count == 1 or line_count >= 4 or len(text) >= 180:
        return "medium"
    return "low"


def _visual_flags(
    text: str,
    *,
    url_count: int,
    image_count: int,
    chart_count: int,
    table_count: int,
) -> list[str]:
    flags: list[str] = []
    if image_count:
        flags.append("image")
    if chart_count:
        flags.append("chart")
    if table_count:
        flags.append("table")
    if url_count:
        flags.append("source_link")
    if re.search(r"[0-9][0-9,.]*(%|원|달러|배|조|억|만|개)", text):
        flags.append("large_number")
    if len(text) >= 500:
        flags.append("dense_text")
    return flags


def _deck_failure_row(row: dict[str, Any], status: str, error: str = "") -> dict[str, Any]:
    return {
        "deck_id": row.get("deck_id"),
        "drive_file_id": row.get("drive_file_id"),
        "title": row.get("title"),
        "inferred_topic_title": row.get("inferred_topic_title"),
        "inferred_deck_type": row.get("inferred_deck_type"),
        "local_path": row.get("local_path", ""),
        "parse_status": status,
        "parse_error": error,
        "slide_count": 0,
        "url_count": 0,
        "unique_url_count": 0,
        "media_count": 0,
        "extracted_at": _now_iso(),
    }


def extract_ppt_corpus_slides(
    *,
    inventory_jsonl: Path = DRIVE_INVENTORY_JSONL,
    deck_manifest_jsonl: Path = DECK_MANIFEST_JSONL,
    slides_jsonl: Path = SLIDES_JSONL,
    links_jsonl: Path = LINKS_JSONL,
    media_manifest_jsonl: Path = MEDIA_MANIFEST_JSONL,
    report_md: Path = EXTRACTION_QUALITY_REPORT_MD,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, Any]], list[dict[str, Any]]]:
    inventory = _read_jsonl_if_exists(inventory_jsonl)
    deck_rows: list[dict[str, Any]] = []
    slide_rows: list[dict[str, Any]] = []
    link_rows: list[dict[str, Any]] = []
    media_rows: list[dict[str, Any]] = []

    for raw_row in inventory:
        row = canonical_inventory_record(raw_row)
        local_path = _resolve_local_path(str(row.get("local_path") or ""))
        if row.get("needs_download") and not local_path.exists():
            deck_rows.append(
                _deck_failure_row(row, "needs_download", str(row.get("skip_reason") or ""))
            )
            continue
        if not local_path.exists():
            deck_rows.append(_deck_failure_row(row, "missing_local_file"))
            continue
        if local_path.suffix.lower() != ".pptx":
            deck_rows.append(
                _deck_failure_row(row, "unsupported_format", local_path.suffix.lower())
            )
            continue

        try:
            parsed = parse_presentation(local_path)
            shape_counts = _shape_counts_by_slide(local_path)
        except Exception as exc:  # pragma: no cover - defensive for local corrupt corpus
            deck_rows.append(_deck_failure_row(row, "failed", str(exc)))
            continue

        deck_rows.append(
            {
                "deck_id": row["deck_id"],
                "drive_file_id": row["drive_file_id"],
                "title": row["title"],
                "inferred_topic_title": row["inferred_topic_title"],
                "inferred_deck_type": row["inferred_deck_type"],
                "local_path": row["local_path"],
                "parse_status": "parsed",
                "slide_count": int(parsed.get("slide_count") or 0),
                "url_count": int(parsed.get("url_count") or 0),
                "unique_url_count": int(parsed.get("unique_url_count") or 0),
                "media_count": int(parsed.get("media_count") or 0),
                "slides_with_urls": int(parsed.get("slides_with_urls") or 0),
                "top_domains": parsed.get("top_domains", {}),
                "first_title": parsed.get("first_title", ""),
                "last_title": parsed.get("last_title", ""),
                "extracted_at": _now_iso(),
            }
        )

        slide_count = int(parsed.get("slide_count") or 0)
        for slide in parsed.get("slides", []):
            slide_no = int(slide.get("slide_no") or 0)
            counts = shape_counts.get(slide_no, {})
            image_count = int(counts.get("image_count") or 0)
            chart_count = int(counts.get("chart_count") or 0)
            table_count = int(counts.get("table_count") or 0)
            visible_text = str(slide.get("visible_text") or "")
            url_count = int(slide.get("url_count") or len(slide.get("all_urls") or []))
            role = _possible_slide_role(
                slide_no=slide_no,
                slide_count=slide_count,
                headline=str(slide.get("headline") or ""),
                visible_text=visible_text,
                url_count=url_count,
                image_count=image_count,
                chart_count=chart_count,
                table_count=table_count,
                slide_type=str(slide.get("slide_type") or ""),
            )
            flags = _visual_flags(
                visible_text,
                url_count=url_count,
                image_count=image_count,
                chart_count=chart_count,
                table_count=table_count,
            )
            slide_rows.append(
                {
                    "deck_id": row["deck_id"],
                    "drive_file_id": row["drive_file_id"],
                    "slide_number": slide_no,
                    "xml_slide_number": int(slide.get("xml_slide_no") or slide_no),
                    "title_candidate": str(slide.get("headline") or ""),
                    "visible_text": visible_text,
                    "visible_char_count": int(slide.get("visible_char_count") or len(visible_text)),
                    "notes": str(slide.get("notes") or ""),
                    "notes_char_count": int(slide.get("notes_char_count") or 0),
                    "hyperlink_count": url_count,
                    "image_count": image_count,
                    "chart_count": chart_count,
                    "table_count": table_count,
                    "media_count": int(slide.get("media_count") or 0),
                    "possible_role": role,
                    "visual_density_estimate": _visual_density(
                        visible_text,
                        image_count=image_count,
                        chart_count=chart_count,
                        table_count=table_count,
                    ),
                    "visual_flags": flags,
                    "parse_status": "parsed",
                }
            )
            media_rows.append(
                {
                    "deck_id": row["deck_id"],
                    "drive_file_id": row["drive_file_id"],
                    "slide_number": slide_no,
                    "image_count": image_count,
                    "chart_count": chart_count,
                    "table_count": table_count,
                    "relationship_media_count": int(slide.get("media_count") or 0),
                    "visual_flags": flags,
                }
            )
            link_rows.extend(_link_rows_for_slide(row, slide))

    write_jsonl(deck_manifest_jsonl, deck_rows)
    write_jsonl(slides_jsonl, slide_rows)
    write_jsonl(links_jsonl, link_rows)
    write_jsonl(media_manifest_jsonl, media_rows)
    report_md.parent.mkdir(parents=True, exist_ok=True)
    report_md.write_text(
        _extraction_quality_markdown(
            deck_rows=deck_rows,
            slide_rows=slide_rows,
            link_rows=link_rows,
            media_rows=media_rows,
            deck_manifest_jsonl=deck_manifest_jsonl,
            slides_jsonl=slides_jsonl,
            links_jsonl=links_jsonl,
            media_manifest_jsonl=media_manifest_jsonl,
        ),
        encoding="utf-8",
    )
    return deck_rows, slide_rows, link_rows, media_rows


def _link_rows_for_slide(row: dict[str, Any], slide: dict[str, Any]) -> list[dict[str, Any]]:
    link_rows: list[dict[str, Any]] = []
    seen: set[tuple[int, str, str]] = set()
    slide_no = int(slide.get("slide_no") or 0)

    for source_channel in ["text_urls", "notes_urls", "relationship_urls"]:
        for url in slide.get(source_channel) or []:
            normalized = canonicalize_url(str(url))
            key = (slide_no, source_channel, normalized)
            if normalized and key not in seen:
                seen.add(key)
                link_rows.append(
                    {
                        "deck_id": row["deck_id"],
                        "drive_file_id": row["drive_file_id"],
                        "slide_number": slide_no,
                        "url": normalized,
                        "domain": _domain(normalized),
                        "source_channel": source_channel,
                        "label": "",
                        "is_image": False,
                        "is_gpt_generated": False,
                    }
                )

    for note in slide.get("source_notes") or []:
        for url in note.get("urls") or []:
            normalized = canonicalize_url(str(url))
            key = (slide_no, "source_note", normalized)
            if normalized and key not in seen:
                seen.add(key)
                link_rows.append(
                    {
                        "deck_id": row["deck_id"],
                        "drive_file_id": row["drive_file_id"],
                        "slide_number": slide_no,
                        "url": normalized,
                        "domain": _domain(normalized),
                        "source_channel": "source_note",
                        "label": str(note.get("label") or ""),
                        "is_image": bool(note.get("is_image")),
                        "is_gpt_generated": bool(note.get("is_gpt_generated")),
                    }
                )
    return link_rows


def _extraction_quality_markdown(
    *,
    deck_rows: list[dict[str, Any]],
    slide_rows: list[dict[str, Any]],
    link_rows: list[dict[str, Any]],
    media_rows: list[dict[str, Any]],
    deck_manifest_jsonl: Path,
    slides_jsonl: Path,
    links_jsonl: Path,
    media_manifest_jsonl: Path,
) -> str:
    parse_counts = Counter(str(row.get("parse_status") or "unknown") for row in deck_rows)
    role_counts = Counter(str(row.get("possible_role") or "unknown") for row in slide_rows)
    density_counts = Counter(
        str(row.get("visual_density_estimate") or "unknown") for row in slide_rows
    )
    lines = [
        "# PPT Corpus Extraction Quality",
        "",
        "Report-only extraction quality summary. No Google Drive, Sheet, or production DB writes.",
        "",
        "## Outputs",
        "",
        f"- deck_manifest_jsonl: `{_repo_relative(deck_manifest_jsonl)}`",
        f"- slides_jsonl: `{_repo_relative(slides_jsonl)}`",
        f"- links_jsonl: `{_repo_relative(links_jsonl)}`",
        f"- media_manifest_jsonl: `{_repo_relative(media_manifest_jsonl)}`",
        "",
        "## Summary",
        "",
        f"- deck_count: {len(deck_rows)}",
        f"- slide_count: {len(slide_rows)}",
        f"- link_count: {len(link_rows)}",
        f"- media_manifest_rows: {len(media_rows)}",
    ]
    lines.extend(f"- parse_status.{key}: {value}" for key, value in sorted(parse_counts.items()))
    lines.extend(["", "## Slide Roles", ""])
    lines.extend(f"- {key}: {value}" for key, value in role_counts.most_common())
    lines.extend(["", "## Visual Density", ""])
    lines.extend(f"- {key}: {value}" for key, value in sorted(density_counts.items()))
    failed = [row for row in deck_rows if row.get("parse_status") != "parsed"]
    if failed:
        lines.extend(["", "## Parse Warnings", ""])
        for row in failed[:50]:
            lines.append(
                f"- {row.get('inferred_topic_title') or row.get('title')}: "
                f"{row.get('parse_status')} {row.get('parse_error') or ''}".strip()
            )
    return "\n".join(lines) + "\n"


def build_ppt_corpus_quality_report(
    *,
    deck_manifest_jsonl: Path = DECK_MANIFEST_JSONL,
    slides_jsonl: Path = SLIDES_JSONL,
    links_jsonl: Path = LINKS_JSONL,
    media_manifest_jsonl: Path = MEDIA_MANIFEST_JSONL,
    report_md: Path = EXTRACTION_QUALITY_REPORT_MD,
) -> str:
    markdown = _extraction_quality_markdown(
        deck_rows=_read_jsonl_if_exists(deck_manifest_jsonl),
        slide_rows=_read_jsonl_if_exists(slides_jsonl),
        link_rows=_read_jsonl_if_exists(links_jsonl),
        media_rows=_read_jsonl_if_exists(media_manifest_jsonl),
        deck_manifest_jsonl=deck_manifest_jsonl,
        slides_jsonl=slides_jsonl,
        links_jsonl=links_jsonl,
        media_manifest_jsonl=media_manifest_jsonl,
    )
    report_md.parent.mkdir(parents=True, exist_ok=True)
    report_md.write_text(markdown, encoding="utf-8")
    return markdown


def build_ppt_corpus_insight_reports(
    *,
    deck_manifest_jsonl: Path = DECK_MANIFEST_JSONL,
    slides_jsonl: Path = SLIDES_JSONL,
    links_jsonl: Path = LINKS_JSONL,
    jibi_md: Path = JIBI_INSIGHTS_MD,
    anny_md: Path = ANNY_ARCHETYPES_MD,
    piti_md: Path = PITI_STYLE_GUIDE_MD,
    report_dir: Path = PPT_CORPUS_REPORT_DIR,
) -> dict[str, Path]:
    deck_rows = _read_jsonl_if_exists(deck_manifest_jsonl)
    slide_rows = _read_jsonl_if_exists(slides_jsonl)
    link_rows = _read_jsonl_if_exists(links_jsonl)
    legacy_sources_path = LEGACY_PPT_LEARNING_OUTPUT_DIR / "drive_ppt_slide_sources.jsonl"
    legacy_source_rows: list[dict[str, Any]] = []
    if not slide_rows and legacy_sources_path.exists():
        legacy_source_rows = _read_jsonl_if_exists(legacy_sources_path)
        slide_rows = _legacy_slide_rows(legacy_source_rows)
    if not link_rows and legacy_source_rows:
        link_rows = _legacy_link_rows(legacy_source_rows)
    if not deck_rows and (LEGACY_PPT_LEARNING_OUTPUT_DIR / "drive_ppt_inventory.jsonl").exists():
        deck_rows = _legacy_deck_rows(
            _read_jsonl_if_exists(LEGACY_PPT_LEARNING_OUTPUT_DIR / "drive_ppt_inventory.jsonl")
        )

    jibi_markdown = _jibi_insights_markdown(deck_rows, slide_rows, link_rows)
    anny_markdown = _anny_archetypes_markdown(deck_rows, slide_rows)
    piti_markdown = _piti_style_markdown(deck_rows, slide_rows)
    recommended_changes = _recommended_code_changes_markdown()

    for path, markdown in [
        (jibi_md, jibi_markdown),
        (anny_md, anny_markdown),
        (piti_md, piti_markdown),
        (report_dir / "03_jibi_selection_lessons.md", jibi_markdown),
        (report_dir / "04_anny_story_archetypes.md", anny_markdown),
        (report_dir / "05_piti_slide_style.md", piti_markdown),
        (report_dir / "07_recommended_code_changes.md", recommended_changes),
    ]:
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(markdown, encoding="utf-8")

    return {
        "jibi": jibi_md,
        "anny": anny_md,
        "piti": piti_md,
        "recommended_code_changes": report_dir / "07_recommended_code_changes.md",
    }


def _legacy_deck_rows(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    deck_rows: list[dict[str, Any]] = []
    for row in rows:
        title = str(row.get("title") or row.get("ppt_title") or "")
        path_hint = str(row.get("local_path") or row.get("drive_path_hint") or "")
        deck_rows.append(
            {
                "deck_id": row.get("ppt_id"),
                "title": title,
                "inferred_topic_title": title,
                "inferred_deck_type": infer_deck_type(
                    title,
                    path_hint=path_hint,
                    mime_type=_mime_from_path(Path(path_hint)),
                ),
                "slide_count": row.get("slide_count", 0),
                "parse_status": row.get("extracted_text_status")
                or row.get("path_status")
                or "unknown",
            }
        )
    return deck_rows


def _legacy_slide_rows(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    legacy_rows: list[dict[str, Any]] = []
    for row in rows:
        text = "\n".join(
            [str(row.get("slide_title") or ""), str(row.get("raw_text_excerpt") or "")]
        )
        legacy_rows.append(
            {
                "deck_id": row.get("ppt_id"),
                "slide_number": row.get("slide_number") or 0,
                "title_candidate": row.get("slide_title") or "",
                "visible_text": text,
                "visible_char_count": len(text),
                "hyperlink_count": len(row.get("extracted_urls") or []),
                "image_count": int("image/proof_object" in (row.get("evidence_types") or [])),
                "chart_count": int("number/statistic" in (row.get("evidence_types") or [])),
                "table_count": 0,
                "possible_role": row.get("slide_type") or row.get("evidence_type") or "explainer",
                "visual_density_estimate": _visual_density(
                    text, image_count=0, chart_count=0, table_count=0
                ),
                "visual_flags": row.get("evidence_types") or [],
            }
        )
    return legacy_rows


def _legacy_link_rows(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    link_rows: list[dict[str, Any]] = []
    seen: set[tuple[str, int, str]] = set()
    for row in rows:
        deck_id = str(row.get("ppt_id") or "")
        slide_number = int(row.get("slide_number") or 0)
        for url in row.get("extracted_urls") or []:
            normalized = canonicalize_url(str(url))
            key = (deck_id, slide_number, normalized)
            if normalized and key not in seen:
                seen.add(key)
                link_rows.append(
                    {
                        "deck_id": deck_id,
                        "slide_number": slide_number,
                        "url": normalized,
                        "domain": _domain(normalized),
                        "source_channel": "legacy_ppt_learning",
                    }
                )
    return link_rows


def _rows_by_deck(slide_rows: list[dict[str, Any]]) -> dict[str, list[dict[str, Any]]]:
    grouped: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for row in slide_rows:
        grouped[str(row.get("deck_id") or "")].append(row)
    for rows in grouped.values():
        rows.sort(key=lambda row: int(row.get("slide_number") or 0))
    return grouped


def _jibi_insights_markdown(
    deck_rows: list[dict[str, Any]],
    slide_rows: list[dict[str, Any]],
    link_rows: list[dict[str, Any]],
) -> str:
    role_counts = Counter(str(row.get("possible_role") or "unknown") for row in slide_rows)
    domain_counts = Counter(str(row.get("domain") or "") for row in link_rows if row.get("domain"))
    deck_type_counts = Counter(str(row.get("inferred_deck_type") or "unknown") for row in deck_rows)
    lines = [
        "# Jibi PPT Corpus Selection Lessons",
        "",
        "Report-only lessons from finished Syukaworld PPT structure. No scoring rules changed.",
        "",
        "## Corpus Signals",
        "",
        f"- deck_count: {len(deck_rows)}",
        f"- slide_count: {len(slide_rows)}",
        f"- link_count: {len(link_rows)}",
        "",
        "## Deck Types",
        "",
    ]
    lines.extend(f"- {key}: {value}" for key, value in sorted(deck_type_counts.items()))
    lines.extend(["", "## Slide Hook/Evidence Roles", ""])
    lines.extend(f"- {key}: {value}" for key, value in role_counts.most_common(20))
    lines.extend(["", "## Top Source Domains", ""])
    if domain_counts:
        lines.extend(f"- {key}: {value}" for key, value in domain_counts.most_common(20))
    else:
        lines.append("- not available in current extraction")
    lines.extend(["", "## Positive Lesson Candidates", ""])
    lines.extend(f"- {lesson}" for lesson in POSITIVE_JIBI_LESSONS)
    lines.extend(["", "## Negative Lesson Candidates", ""])
    lines.extend(f"- {lesson}" for lesson in NEGATIVE_JIBI_LESSONS)
    lines.extend(
        [
            "",
            "## Guardrails",
            "",
            "- Treat these as priors for review, not production scoring changes.",
            "- Require outcome join before converting a lesson into a Jibi calibration rule.",
            "- Keep PPT-created-but-not-used decks as potential negative examples.",
        ]
    )
    return "\n".join(lines) + "\n"


def _anny_archetypes_markdown(
    deck_rows: list[dict[str, Any]],
    slide_rows: list[dict[str, Any]],
) -> str:
    grouped = _rows_by_deck(slide_rows)
    examples: list[str] = []
    deck_lookup = {str(row.get("deck_id")): row for row in deck_rows}
    for deck_id, rows in list(grouped.items())[:12]:
        first_roles = [str(row.get("possible_role") or "") for row in rows[:8]]
        title = deck_lookup.get(deck_id, {}).get("inferred_topic_title") or deck_id
        examples.append(f"- {title}: {' -> '.join(role for role in first_roles if role)}")

    lines = [
        "# Anny Story Archetypes From PPT Corpus",
        "",
        "Report-only archetype inventory for storyline expansion. Anny schema is unchanged.",
        "",
        "## Archetype Candidates",
        "",
    ]
    lines.extend(f"- {item}" for item in ANNY_ARCHETYPES)
    lines.extend(["", "## Observed Early-Slide Spines", ""])
    lines.extend(examples or ["- not available in current extraction"])
    lines.extend(
        [
            "",
            "## Contract Candidate",
            "",
            "```json",
            "{",
            '  "story_archetype": "daily_life_to_hidden_cost",',
            '  "opening_question": "...",',
            '  "must_have_evidence": ["official_stat", "case", "korea_bridge"],',
            '  "recommended_slide_spine": [',
            '    "hook_scene", "why_now", "mechanism", "case", "number",',
            '    "korea_bridge", "takeaway"',
            "  ],",
            '  "avoid_frames": ["moral_scolding", "generic_AI_discourse"]',
            "}",
            "```",
        ]
    )
    return "\n".join(lines) + "\n"


def _piti_style_markdown(
    deck_rows: list[dict[str, Any]],
    slide_rows: list[dict[str, Any]],
) -> str:
    density_counts = Counter(
        str(row.get("visual_density_estimate") or "unknown") for row in slide_rows
    )
    role_counts = Counter(str(row.get("possible_role") or "unknown") for row in slide_rows)
    slide_counts = [
        int(row.get("slide_count") or 0) for row in deck_rows if int(row.get("slide_count") or 0)
    ]
    avg_slide_count = sum(slide_counts) / len(slide_counts) if slide_counts else 0.0
    avg_chars = (
        sum(
            int(row.get("visible_char_count") or len(str(row.get("visible_text") or "")))
            for row in slide_rows
        )
        / len(slide_rows)
        if slide_rows
        else 0.0
    )
    lines = [
        "# Piti Syukaworld Slide Style Guide",
        "",
        "Report-only slide grammar summary. Piti generation is unchanged.",
        "",
        "## Density",
        "",
        f"- deck_count: {len(deck_rows)}",
        f"- slide_count: {len(slide_rows)}",
        f"- average_slides_per_deck: {avg_slide_count:.1f}",
        f"- average_visible_chars_per_slide: {avg_chars:.1f}",
    ]
    lines.extend(
        f"- visual_density.{key}: {value}" for key, value in sorted(density_counts.items())
    )
    lines.extend(["", "## Slide Pattern Taxonomy", ""])
    lines.extend(f"- {pattern}" for pattern in SLIDE_PATTERN_TAXONOMY)
    lines.extend(["", "## Observed Pattern Counts", ""])
    lines.extend(f"- {key}: {value}" for key, value in role_counts.most_common(20))
    lines.extend(
        [
            "",
            "## Style Guardrails",
            "",
            "- Prefer contact-sheet review before treating visual pattern counts as ground truth.",
            "- Keep OCR/LLM visual analysis out of the first deterministic extraction pass.",
            "- Use examples only after checking copyright and source-link provenance.",
        ]
    )
    return "\n".join(lines) + "\n"


def _recommended_code_changes_markdown() -> str:
    return "\n".join(
        [
            "# Recommended Code Changes",
            "",
            "Report-only PR slicing for PPT corpus work.",
            "",
            "- PR A: Add PPT corpus inventory and extractor.",
            "- PR B: Add deck story profile report.",
            "- PR C: Add Jibi positive/negative lessons from PPT corpus.",
            "- PR D: Add Anny story archetype examples.",
            "- PR E: Add Piti slide style guide and pattern taxonomy.",
            "",
            "Do not change production Jibi scoring, Anny handoff, or Piti generation "
            "until reports are reviewed.",
            "",
        ]
    )


@drive_manifest_app.callback(invoke_without_command=True)
def drive_manifest_main(
    latest_root: Annotated[
        Path,
        typer.Option("--latest-root", help="Unzipped latest Drive PPT root."),
    ] = PPT_CORPUS_DATA_DIR / "drive_raw" / "latest",
    past_root: Annotated[
        Path,
        typer.Option("--past-root", help="Unzipped past Drive PPT root."),
    ] = PPT_CORPUS_DATA_DIR / "drive_raw" / "past",
    output_jsonl: Annotated[
        Path,
        typer.Option("--output-jsonl", help="Canonical PPT corpus inventory JSONL."),
    ] = DRIVE_INVENTORY_JSONL,
    report_md: Annotated[
        Path,
        typer.Option("--report-md", help="Inventory Markdown report."),
    ] = DRIVE_INVENTORY_REPORT_MD,
) -> None:
    records = build_drive_manifest(
        latest_root=latest_root,
        past_root=past_root,
        output_jsonl=output_jsonl,
        report_md=report_md,
    )
    console.print(
        f"[green]Wrote {len(records)} PPT corpus inventory rows to {output_jsonl}[/green]"
    )


@inventory_app.callback(invoke_without_command=True)
def inventory_main(
    manifest: Annotated[
        Path,
        typer.Option("--manifest", help="Input Drive/local manifest JSONL."),
    ] = DRIVE_INVENTORY_JSONL,
    output_jsonl: Annotated[
        Path,
        typer.Option("--output-jsonl", help="Canonical PPT corpus inventory JSONL."),
    ] = DRIVE_INVENTORY_JSONL,
    report_md: Annotated[
        Path,
        typer.Option("--report-md", help="Inventory Markdown report."),
    ] = DRIVE_INVENTORY_REPORT_MD,
) -> None:
    records = build_inventory(manifest=manifest, output_jsonl=output_jsonl, report_md=report_md)
    console.print(
        f"[green]Wrote {len(records)} PPT corpus inventory rows to {output_jsonl}[/green]"
    )


@extract_slides_app.callback(invoke_without_command=True)
def extract_slides_main(
    inventory_jsonl: Annotated[
        Path,
        typer.Option("--inventory-jsonl", help="Canonical PPT corpus inventory JSONL."),
    ] = DRIVE_INVENTORY_JSONL,
    deck_manifest_jsonl: Annotated[
        Path,
        typer.Option("--deck-manifest-jsonl", help="Deck manifest output JSONL."),
    ] = DECK_MANIFEST_JSONL,
    slides_jsonl: Annotated[
        Path,
        typer.Option("--slides-jsonl", help="Slide extraction output JSONL."),
    ] = SLIDES_JSONL,
    links_jsonl: Annotated[
        Path,
        typer.Option("--links-jsonl", help="Hyperlink output JSONL."),
    ] = LINKS_JSONL,
    media_manifest_jsonl: Annotated[
        Path,
        typer.Option("--media-manifest-jsonl", help="Media manifest output JSONL."),
    ] = MEDIA_MANIFEST_JSONL,
    report_md: Annotated[
        Path,
        typer.Option("--report-md", help="Extraction quality Markdown report."),
    ] = EXTRACTION_QUALITY_REPORT_MD,
) -> None:
    deck_rows, slide_rows, link_rows, media_rows = extract_ppt_corpus_slides(
        inventory_jsonl=inventory_jsonl,
        deck_manifest_jsonl=deck_manifest_jsonl,
        slides_jsonl=slides_jsonl,
        links_jsonl=links_jsonl,
        media_manifest_jsonl=media_manifest_jsonl,
        report_md=report_md,
    )
    console.print(
        f"[green]Wrote PPT corpus extraction: decks={len(deck_rows)}, "
        f"slides={len(slide_rows)}, links={len(link_rows)}, media={len(media_rows)}[/green]"
    )


@quality_report_app.callback(invoke_without_command=True)
def quality_report_main(
    deck_manifest_jsonl: Annotated[
        Path,
        typer.Option("--deck-manifest-jsonl", help="Deck manifest input JSONL."),
    ] = DECK_MANIFEST_JSONL,
    slides_jsonl: Annotated[
        Path,
        typer.Option("--slides-jsonl", help="Slide extraction input JSONL."),
    ] = SLIDES_JSONL,
    links_jsonl: Annotated[
        Path,
        typer.Option("--links-jsonl", help="Hyperlink input JSONL."),
    ] = LINKS_JSONL,
    media_manifest_jsonl: Annotated[
        Path,
        typer.Option("--media-manifest-jsonl", help="Media manifest input JSONL."),
    ] = MEDIA_MANIFEST_JSONL,
    report_md: Annotated[
        Path,
        typer.Option("--report-md", help="Extraction quality Markdown report."),
    ] = EXTRACTION_QUALITY_REPORT_MD,
) -> None:
    markdown = build_ppt_corpus_quality_report(
        deck_manifest_jsonl=deck_manifest_jsonl,
        slides_jsonl=slides_jsonl,
        links_jsonl=links_jsonl,
        media_manifest_jsonl=media_manifest_jsonl,
        report_md=report_md,
    )
    console.print(f"[green]Wrote PPT corpus quality report to {report_md}[/green]")
    console.print(f"[green]{len(markdown.splitlines())} report lines[/green]")


@insight_reports_app.callback(invoke_without_command=True)
def insight_reports_main(
    deck_manifest_jsonl: Annotated[
        Path,
        typer.Option("--deck-manifest-jsonl", help="Deck manifest input JSONL."),
    ] = DECK_MANIFEST_JSONL,
    slides_jsonl: Annotated[
        Path,
        typer.Option("--slides-jsonl", help="Slide extraction input JSONL."),
    ] = SLIDES_JSONL,
    links_jsonl: Annotated[
        Path,
        typer.Option("--links-jsonl", help="Hyperlink input JSONL."),
    ] = LINKS_JSONL,
    report_dir: Annotated[
        Path,
        typer.Option("--report-dir", help="PPT corpus report directory."),
    ] = PPT_CORPUS_REPORT_DIR,
) -> None:
    outputs = build_ppt_corpus_insight_reports(
        deck_manifest_jsonl=deck_manifest_jsonl,
        slides_jsonl=slides_jsonl,
        links_jsonl=links_jsonl,
        report_dir=report_dir,
    )
    console.print("[green]Wrote PPT corpus insight reports:[/green]")
    for label, path in outputs.items():
        console.print(f"[green]- {label}: {path}[/green]")
