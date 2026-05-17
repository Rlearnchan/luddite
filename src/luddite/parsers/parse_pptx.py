"""Parse local PPTX decks into slide-level corpus JSONL."""

from __future__ import annotations

import hashlib
import re
import unicodedata
import zipfile
from collections import Counter
from datetime import UTC, datetime
from pathlib import Path
from typing import Annotated, Any
from xml.etree import ElementTree as ET

import typer
from pptx import Presentation

from luddite import paths
from luddite.utils.jsonl import write_jsonl
from luddite.utils.security import detect_risk_flags, extract_source_notes, redact_sensitive_text
from luddite.utils.urls import extract_urls

app = typer.Typer(no_args_is_help=False)

REL_NS = {"rel": "http://schemas.openxmlformats.org/package/2006/relationships"}
A_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
URLISH_REL_TYPES = ("hyperlink",)
MEDIA_TARGET_RE = re.compile(r"(?:^|/)(media|charts|embeddings)/", re.IGNORECASE)


def _now_iso() -> str:
    return datetime.now(UTC).isoformat()


def _relative(path: Path) -> str:
    try:
        return str(path.relative_to(paths.REPO_ROOT))
    except ValueError:
        return str(path)


def _stable_id(prefix: str, path: Path) -> str:
    digest = hashlib.sha1(_relative(path).encode("utf-8")).hexdigest()[:12]
    return f"{prefix}_{digest}"


def _normalize_text(text: str) -> str:
    text = unicodedata.normalize("NFC", text)
    text = text.replace("\xa0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _slide_part_number(partname: str) -> int | None:
    match = re.search(r"slide(\d+)\.xml$", partname)
    return int(match.group(1)) if match else None


def _read_zip_text(deck: zipfile.ZipFile, name: str) -> str:
    return deck.read(name).decode("utf-8", errors="ignore")


def _extract_xml_text(xml_text: str) -> str:
    root = ET.fromstring(xml_text)
    paragraphs: list[str] = []
    for paragraph in root.findall(".//a:p", A_NS):
        values = [node.text or "" for node in paragraph.findall(".//a:t", A_NS)]
        text = "".join(values).strip()
        if text:
            paragraphs.append(text)
    joined = "\n".join(paragraphs)
    return _normalize_text(redact_sensitive_text(joined))


def _relationships(deck: zipfile.ZipFile, rels_name: str) -> list[dict[str, str]]:
    if rels_name not in deck.namelist():
        return []
    root = ET.fromstring(_read_zip_text(deck, rels_name))
    relationships: list[dict[str, str]] = []
    for relationship in root.findall("rel:Relationship", REL_NS):
        relationships.append(
            {
                "id": relationship.attrib.get("Id", ""),
                "type": relationship.attrib.get("Type", ""),
                "target": relationship.attrib.get("Target", ""),
                "target_mode": relationship.attrib.get("TargetMode", ""),
            }
        )
    return relationships


def _notes_path_from_rels(rels: list[dict[str, str]]) -> str | None:
    for rel in rels:
        if rel["type"].endswith("/notesSlide"):
            target = rel["target"].replace("../", "")
            return f"ppt/{target}" if not target.startswith("ppt/") else target
    return None


def _external_urls_from_rels(rels: list[dict[str, str]]) -> list[str]:
    text = "\n".join(
        rel["target"]
        for rel in rels
        if rel["target_mode"].lower() == "external"
        or any(rel["type"].endswith(kind) for kind in URLISH_REL_TYPES)
    )
    return extract_urls(text)


def _media_count_from_rels(rels: list[dict[str, str]]) -> int:
    return sum(1 for rel in rels if MEDIA_TARGET_RE.search(rel["target"].replace("..", "")))


def _visible_slide_text(slide: Any) -> str:
    lines: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if text:
                lines.append(text)
    return _normalize_text(redact_sensitive_text("\n".join(lines)))


def _headline_and_body(text: str) -> tuple[str, list[str]]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        return "", []
    return lines[0], lines[1:]


def _infer_slide_type(
    slide_no: int,
    slide_count: int,
    headline: str,
    body: list[str],
    text: str,
    notes: str,
    url_count: int,
    media_count: int,
) -> str:
    line_count = len([line for line in text.splitlines() if line.strip()])
    compact_text = text.replace("\n", " ")
    if slide_no == 1:
        return "title"
    if line_count <= 2 and len(compact_text) <= 90 and url_count == 0:
        return "section_title"
    if slide_no > slide_count * 0.85 and "?" in compact_text:
        return "closing_question"
    if any(token in compact_text for token in ["슈카콜라", "ㅋㅋ", "아닐까", "^^"]):
        return "punchline"
    if url_count >= 4:
        return "source_heavy"
    if media_count > 0 and len(compact_text) <= 120:
        return "image_centered"
    if re.search(r"[0-9][0-9,.]*(%|원|달러|배|조|억|만|개)", compact_text):
        return "data"
    if any(token in compact_text for token in ["“", "”", "\"", "발언", "원문"]):
        return "quote"
    if any(token in compact_text for token in ["비교", "반면", "vs", "VS"]):
        return "comparison"
    if slide_no <= 5 or "?" in headline:
        return "hook"
    return "explainer"


def _domain_counts(urls: list[str]) -> dict[str, int]:
    domains: Counter[str] = Counter()
    for url in urls:
        match = re.match(r"https?://([^/]+)", url)
        if match:
            domains[match.group(1).lower()] += 1
    return dict(domains.most_common())


def _dedupe_urls(urls: list[str]) -> list[str]:
    unique: list[str] = []
    for url in urls:
        if url not in unique:
            unique.append(url)
    return unique


def parse_presentation(path: Path) -> dict[str, Any]:
    presentation = Presentation(str(path))
    slide_count = len(presentation.slides)
    slides: list[dict[str, Any]] = []
    deck_urls: list[str] = []
    total_media = 0

    with zipfile.ZipFile(path) as deck:
        names = set(deck.namelist())
        media_files = [name for name in names if name.startswith("ppt/media/")]

        for index, slide in enumerate(presentation.slides, start=1):
            partname = str(slide.part.partname).lstrip("/")
            slide_number = _slide_part_number(partname) or index
            rels_name = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
            rels = _relationships(deck, rels_name)
            notes_path = _notes_path_from_rels(rels)
            notes = ""
            if notes_path and notes_path in names:
                notes = _extract_xml_text(_read_zip_text(deck, notes_path))

            visible_text = _visible_slide_text(slide)
            headline, body = _headline_and_body(visible_text)
            relationship_urls = _external_urls_from_rels(rels)
            source_notes = extract_source_notes(notes)
            text_urls = extract_urls(visible_text)
            notes_urls = extract_urls(notes)
            all_urls = []
            for url in [*text_urls, *notes_urls, *relationship_urls]:
                if url not in all_urls:
                    all_urls.append(url)
            deck_urls.extend(all_urls)

            slide_media_count = _media_count_from_rels(rels)
            total_media += slide_media_count
            slide_type = _infer_slide_type(
                slide_no=index,
                slide_count=slide_count,
                headline=headline,
                body=body,
                text=visible_text,
                notes=notes,
                url_count=len(all_urls),
                media_count=slide_media_count,
            )
            image_urls = [
                url
                for note in source_notes
                if note["is_image"]
                for url in note["urls"]
            ]
            source_urls = [
                url
                for note in source_notes
                if not note["is_image"]
                for url in note["urls"]
            ]
            if not source_urls and not source_notes:
                source_urls = notes_urls

            slides.append(
                {
                    "slide_no": index,
                    "xml_slide_no": slide_number,
                    "slide_type": slide_type,
                    "is_section_title": slide_type == "section_title",
                    "headline": headline,
                    "body": body,
                    "visible_text": visible_text,
                    "visible_char_count": len(visible_text),
                    "notes": notes,
                    "notes_char_count": len(notes),
                    "source_notes": source_notes,
                    "source_urls": source_urls,
                    "image_urls": image_urls,
                    "text_urls": text_urls,
                    "notes_urls": notes_urls,
                    "relationship_urls": relationship_urls,
                    "all_urls": all_urls,
                    "url_count": len(all_urls),
                    "media_count": slide_media_count,
                    "risk_flags": detect_risk_flags(f"{visible_text}\n{notes}"),
                }
            )

    title = unicodedata.normalize("NFC", path.stem)
    return {
        "corpus_id": _stable_id("pptx", path),
        "type": "pptx",
        "title": title,
        "file_name": path.name,
        "local_path": _relative(path),
        "parse_status": "parsed",
        "parsed_at": _now_iso(),
        "urls": _dedupe_urls(deck_urls),
        "slide_count": slide_count,
        "media_count": len(media_files) or total_media,
        "relationship_media_count": total_media,
        "url_count": len(deck_urls),
        "unique_url_count": len(_dedupe_urls(deck_urls)),
        "slides_with_urls": sum(1 for slide in slides if slide["all_urls"]),
        "top_domains": _domain_counts(deck_urls),
        "first_title": slides[0]["visible_text"] if slides else "",
        "last_title": slides[-1]["visible_text"] if slides else "",
        "slides": slides,
        "risk_flags": detect_risk_flags(title),
    }


def iter_pptx_files(input_dir: Path) -> list[Path]:
    if input_dir.is_file() and input_dir.suffix.lower() == ".pptx":
        return [input_dir]
    if not input_dir.exists():
        return []
    return sorted(
        input_dir.glob("*.pptx"),
        key=lambda path: unicodedata.normalize("NFC", path.name),
    )


def parse_directory(input_dir: Path, output_path: Path) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    for file_path in iter_pptx_files(input_dir):
        try:
            records.append(parse_presentation(file_path))
        except Exception as exc:  # pragma: no cover - defensive corpus reporting
            records.append(
                {
                    "corpus_id": _stable_id("pptx", file_path),
                    "type": "pptx",
                    "title": unicodedata.normalize("NFC", file_path.stem),
                    "file_name": file_path.name,
                    "local_path": _relative(file_path),
                    "parse_status": "failed",
                    "parsed_at": _now_iso(),
                    "error": str(exc),
                    "risk_flags": ["needs_human_review"],
                }
            )
    write_jsonl(output_path, records)
    return records


@app.callback(invoke_without_command=True)
def main(
    input_dir: Annotated[
        Path,
        typer.Option("--input-dir", help="Directory or single .pptx file to parse."),
    ] = paths.LATEST_PPT_RAW_DIR,
    output: Annotated[
        Path,
        typer.Option("--output", help="JSONL output path."),
    ] = paths.PPT_PARSED_JSONL,
) -> None:
    records = parse_directory(input_dir, output)
    typer.echo(f"Wrote {len(records)} PPTX records to {output}")


if __name__ == "__main__":
    app()
