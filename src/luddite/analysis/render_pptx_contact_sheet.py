"""Render review-only contact sheet QA reports for generated PPTX drafts."""

from __future__ import annotations

import json
import shutil
import subprocess
from collections import Counter
from collections.abc import Callable
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from typing import Annotated, Any

import typer
from pptx import Presentation
from rich.console import Console

from luddite import paths
from luddite.agents.piti.render_visual_qa import VisualQaDeck, evaluate_slide_spec

app = typer.Typer(no_args_is_help=False)
console = Console()

DEFAULT_OUTPUT_DIR = paths.OUTPUTS_DIR / "qa" / "pptx_contact_sheet"
DEFAULT_SUMMARY_PATH = DEFAULT_OUTPUT_DIR / "pptx_contact_sheet_summary.md"
DEFAULT_REVIEW_SUMMARY_PATH = paths.DOCS_DIR / "reviews" / "pptx_contact_sheet_summary.md"
DEFAULT_DIRECT_RUN_ID = "live_m132_20260520_all"
DEFAULT_DIRECT_OUTPUT_ROOT = paths.MODEL_DRY_RUNS_DIR / "anny_slide_spec_experiments_live"


@dataclass(frozen=True)
class ContactSheetTarget:
    deck_id: str
    pptx_path: Path
    slide_spec_path: Path | None
    source_kind: str


@dataclass(frozen=True)
class ThumbnailGeneration:
    status: str
    backend: str
    thumbnails: list[Path]
    contact_sheet_path: Path | None
    contact_sheet_pdf_path: Path | None
    pdf_path: Path | None
    warnings: list[str]


@dataclass(frozen=True)
class ContactSheetSlide:
    slide_no: int
    thumbnail_path: Path | None
    screen_headline: str
    layout_intent: str
    proof_object_type: str
    visual_qa_flags: list[str]
    contact_sheet_review_status: str = "unchecked"
    reviewer_note: str = ""


@dataclass(frozen=True)
class ContactSheetResult:
    target: ContactSheetTarget
    status: str
    slide_count: int
    spec_slide_count: int | None
    thumbnail_count: int
    thumbnail_status: str
    backend: str
    pdf_path: Path | None
    contact_sheet_path: Path | None
    contact_sheet_pdf_path: Path | None
    report_path: Path
    warnings: list[str]
    slides: list[ContactSheetSlide]


ThumbnailGenerator = Callable[[ContactSheetTarget, Path], ThumbnailGeneration]


def _display_path(path: Path | None) -> str:
    if path is None:
        return "-"
    try:
        return str(path.resolve().relative_to(paths.REPO_ROOT))
    except ValueError:
        return str(path)


def _safe_deck_id(path: Path) -> str:
    return "".join(
        character if character.isalnum() or character in "-_" else "_"
        for character in path.stem
    )


def _load_json(path: Path) -> dict[str, Any]:
    with path.open(encoding="utf-8") as source:
        payload = json.load(source)
    if not isinstance(payload, dict):
        raise ValueError(f"Expected JSON object at {path}")
    return payload


def _as_list(value: Any) -> list[Any]:
    return value if isinstance(value, list) else []


def _proof_type(slide: dict[str, Any]) -> str:
    proof = slide.get("proof_object")
    if not isinstance(proof, dict):
        return "none"
    return str(proof.get("type") or "none")


def _slide_no(slide: dict[str, Any]) -> int:
    value = slide.get("slide_no")
    return int(value) if isinstance(value, int | float | str) and str(value).isdigit() else 0


def _slides_by_no(spec: dict[str, Any] | None) -> dict[int, dict[str, Any]]:
    if spec is None:
        return {}
    slides = [slide for slide in _as_list(spec.get("slides")) if isinstance(slide, dict)]
    return {_slide_no(slide): slide for slide in slides}


def _visual_qa_flags_by_no(deck: VisualQaDeck | None) -> dict[int, list[str]]:
    if deck is None:
        return {}
    return {slide.slide_no: slide.visual_qa_flags for slide in deck.slides}


def _slide_count_for_pptx(path: Path) -> tuple[int, str | None]:
    try:
        presentation = Presentation(path)
    except Exception as error:  # pragma: no cover - python-pptx raises several concrete types.
        return 0, f"pptx_read_failed: {error}"
    return len(presentation.slides), None


def _find_first_command(candidates: list[str]) -> str | None:
    for candidate in candidates:
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
    return None


def _run_command(args: list[str]) -> tuple[bool, str]:
    completed = subprocess.run(
        args,
        check=False,
        capture_output=True,
        text=True,
        timeout=120,
    )
    output = "\n".join(part for part in [completed.stdout, completed.stderr] if part)
    return completed.returncode == 0, output.strip()


def _convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> tuple[Path | None, str | None, str]:
    backend = _find_first_command(["soffice", "libreoffice"])
    if not backend:
        return None, "LibreOffice/soffice not found; thumbnail generation skipped.", "none"
    output_dir.mkdir(parents=True, exist_ok=True)
    ok, output = _run_command(
        [
            backend,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ]
    )
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if not ok or not pdf_path.exists():
        detail = output or "unknown LibreOffice conversion error"
        return None, f"pptx_to_pdf_failed: {detail}", Path(backend).name
    return pdf_path, None, Path(backend).name


def _rasterize_pdf(pdf_path: Path, output_dir: Path) -> tuple[list[Path], str | None, str]:
    backend = _find_first_command(["pdftoppm"])
    if not backend:
        return [], "pdftoppm not found; PDF was generated but thumbnails were skipped.", "none"
    output_dir.mkdir(parents=True, exist_ok=True)
    prefix = output_dir / "slide"
    ok, output = _run_command(
        [
            backend,
            "-png",
            "-r",
            "120",
            str(pdf_path),
            str(prefix),
        ]
    )
    thumbnails = sorted(output_dir.glob("slide-*.png"))
    if not ok or not thumbnails:
        detail = output or "unknown pdftoppm rasterization error"
        return [], f"pdf_thumbnail_failed: {detail}", Path(backend).name
    return thumbnails, None, Path(backend).name


def _compose_contact_sheet(
    *,
    thumbnails: list[Path],
    output_path: Path,
    columns: int = 4,
) -> tuple[Path | None, Path | None, str | None]:
    if not thumbnails:
        return None, None, "no thumbnails available for contact sheet composition"
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        return None, None, "Pillow not available; per-slide thumbnails were generated only."

    opened = [Image.open(path).convert("RGB") for path in thumbnails]
    thumb_width = 360
    label_height = 34
    gutter = 20
    scaled: list[Image.Image] = []
    for image in opened:
        ratio = thumb_width / image.width
        scaled.append(image.resize((thumb_width, int(image.height * ratio))))
    cell_height = max(image.height for image in scaled) + label_height
    rows = (len(scaled) + columns - 1) // columns
    sheet_width = columns * thumb_width + (columns + 1) * gutter
    sheet_height = rows * cell_height + (rows + 1) * gutter
    sheet = Image.new("RGB", (sheet_width, sheet_height), "white")
    draw = ImageDraw.Draw(sheet)
    font = ImageFont.load_default()
    for index, image in enumerate(scaled):
        row = index // columns
        column = index % columns
        x = gutter + column * (thumb_width + gutter)
        y = gutter + row * (cell_height + gutter)
        draw.text((x, y), f"Slide {index + 1}", fill=(40, 40, 40), font=font)
        sheet.paste(image, (x, y + label_height))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)
    pdf_path = output_path.with_suffix(".pdf")
    sheet.save(pdf_path)
    return output_path, pdf_path, None


def default_thumbnail_generator(
    target: ContactSheetTarget,
    output_dir: Path,
) -> ThumbnailGeneration:
    warnings: list[str] = []
    pdf_dir = output_dir / target.deck_id / "pdf"
    thumb_dir = output_dir / target.deck_id / "thumbnails"
    pdf_path, warning, pdf_backend = _convert_pptx_to_pdf(target.pptx_path, pdf_dir)
    if warning:
        warnings.append(warning)
    if pdf_path is None:
        return ThumbnailGeneration(
            status="backend_unavailable_or_failed",
            backend=pdf_backend,
            thumbnails=[],
            contact_sheet_path=None,
            contact_sheet_pdf_path=None,
            pdf_path=None,
            warnings=warnings,
        )
    thumbnails, warning, raster_backend = _rasterize_pdf(pdf_path, thumb_dir)
    if warning:
        warnings.append(warning)
    if not thumbnails:
        return ThumbnailGeneration(
            status="pdf_only",
            backend=f"{pdf_backend}+{raster_backend}",
            thumbnails=[],
            contact_sheet_path=None,
            contact_sheet_pdf_path=None,
            pdf_path=pdf_path,
            warnings=warnings,
        )
    contact_sheet_path, contact_sheet_pdf_path, warning = _compose_contact_sheet(
        thumbnails=thumbnails,
        output_path=output_dir / target.deck_id / f"{target.deck_id}_contact_sheet.png",
    )
    if warning:
        warnings.append(warning)
    status = "generated" if contact_sheet_path else "thumbnails_only"
    return ThumbnailGeneration(
        status=status,
        backend=f"{pdf_backend}+{raster_backend}",
        thumbnails=thumbnails,
        contact_sheet_path=contact_sheet_path,
        contact_sheet_pdf_path=contact_sheet_pdf_path,
        pdf_path=pdf_path,
        warnings=warnings,
    )


def default_targets(
    *,
    include_direct: bool = True,
    direct_run_id: str = DEFAULT_DIRECT_RUN_ID,
    direct_output_root: Path = DEFAULT_DIRECT_OUTPUT_ROOT,
) -> list[ContactSheetTarget]:
    targets = [
        ContactSheetTarget(
            deck_id="adapter_ai_knowledge_institution",
            pptx_path=paths.PPTX_OUTPUT_DIR
            / "ai_knowledge_institution_slide_spec_styled_draft.pptx",
            slide_spec_path=paths.PITI_SLIDE_SPECS_DIR / "ai_knowledge_institution_slide_spec.json",
            source_kind="adapter",
        ),
        ContactSheetTarget(
            deck_id="adapter_productive_finance_policy",
            pptx_path=paths.PPTX_OUTPUT_DIR
            / "productive_finance_policy_slide_spec_styled_draft.pptx",
            slide_spec_path=paths.PITI_SLIDE_SPECS_DIR
            / "productive_finance_policy_slide_spec.json",
            source_kind="adapter",
        ),
    ]
    if include_direct:
        direct_root = direct_output_root / direct_run_id
        targets.extend(
            [
                ContactSheetTarget(
                    deck_id=f"{direct_run_id}_ai_knowledge_institution",
                    pptx_path=direct_root
                    / "ai_knowledge_institution"
                    / "direct_piti_slide_spec_draft.pptx",
                    slide_spec_path=direct_root
                    / "ai_knowledge_institution"
                    / "parsed_piti_slide_spec.json",
                    source_kind="direct_live",
                ),
                ContactSheetTarget(
                    deck_id=f"{direct_run_id}_productive_finance_policy",
                    pptx_path=direct_root
                    / "productive_finance_policy"
                    / "direct_piti_slide_spec_draft.pptx",
                    slide_spec_path=direct_root
                    / "productive_finance_policy"
                    / "parsed_piti_slide_spec.json",
                    source_kind="direct_live",
                ),
            ]
        )
    return targets


def _custom_targets(pptx_paths: list[Path]) -> list[ContactSheetTarget]:
    return [
        ContactSheetTarget(
            deck_id=_safe_deck_id(path),
            pptx_path=path,
            slide_spec_path=None,
            source_kind="custom",
        )
        for path in pptx_paths
    ]


def _load_slide_spec(
    slide_spec_path: Path | None,
    output_dir: Path,
) -> tuple[dict[str, Any] | None, VisualQaDeck | None, list[str]]:
    if slide_spec_path is None:
        return None, None, []
    if not slide_spec_path.exists():
        return None, None, [f"slide_spec_missing: {_display_path(slide_spec_path)}"]
    try:
        spec = _load_json(slide_spec_path)
    except (OSError, ValueError, json.JSONDecodeError) as error:
        return None, None, [f"slide_spec_unreadable: {error}"]
    deck = evaluate_slide_spec(slide_spec_path, spec, output_dir / "visual_qa")
    return spec, deck, []


def _build_slide_rows(
    *,
    slide_count: int,
    spec: dict[str, Any] | None,
    visual_qa_deck: VisualQaDeck | None,
    thumbnails: list[Path],
) -> list[ContactSheetSlide]:
    spec_slides = _slides_by_no(spec)
    qa_flags = _visual_qa_flags_by_no(visual_qa_deck)
    row_count = max(slide_count, len(spec_slides), len(thumbnails))
    rows: list[ContactSheetSlide] = []
    for index in range(row_count):
        slide_no = index + 1
        slide = spec_slides.get(slide_no, {})
        thumbnail_path = thumbnails[index] if index < len(thumbnails) else None
        flags = list(qa_flags.get(slide_no, []))
        if thumbnail_path is None:
            flags.append("thumbnail_missing")
        rows.append(
            ContactSheetSlide(
                slide_no=slide_no,
                thumbnail_path=thumbnail_path,
                screen_headline=str(slide.get("screen_headline") or ""),
                layout_intent=str(slide.get("layout_intent") or ""),
                proof_object_type=_proof_type(slide),
                visual_qa_flags=flags,
            )
        )
    return rows


def _write_deck_report(result: ContactSheetResult) -> None:
    result.report_path.parent.mkdir(parents=True, exist_ok=True)
    lines = [
        f"# PPTX Contact Sheet QA: {result.target.deck_id}",
        "",
        f"- generated_at: {datetime.now(UTC).isoformat()}",
        f"- source_kind: {result.target.source_kind}",
        f"- pptx: {_display_path(result.target.pptx_path)}",
        f"- slide_spec: {_display_path(result.target.slide_spec_path)}",
        f"- status: {result.status}",
        f"- slide_count: {result.slide_count}",
        (
            "- spec_slide_count: "
            f"{result.spec_slide_count if result.spec_slide_count is not None else '-'}"
        ),
        f"- thumbnail_status: {result.thumbnail_status}",
        f"- thumbnail_backend: {result.backend}",
        f"- contact_sheet: {_display_path(result.contact_sheet_path)}",
        f"- contact_sheet_pdf: {_display_path(result.contact_sheet_pdf_path)}",
        f"- pdf: {_display_path(result.pdf_path)}",
        "- This is visual review surface only.",
        "- No PPT content was modified.",
        "- No LLM/API calls.",
        "- Broadcast readiness remains false.",
        "",
        "## Warnings",
        "",
    ]
    if result.warnings:
        lines.extend(f"- {warning}" for warning in result.warnings)
    else:
        lines.append("- none")
    lines.extend(
        [
            "",
            "## Slide Review",
            "",
            (
                "| slide_no | thumbnail | screen_headline | layout_intent | "
                "proof_object.type | visual QA flags | contact_sheet_review_status | "
                "reviewer_note |"
            ),
            "|---:|---|---|---|---|---|---|---|",
        ]
    )
    for slide in result.slides:
        thumbnail = _display_path(slide.thumbnail_path)
        flags = ", ".join(slide.visual_qa_flags) if slide.visual_qa_flags else "-"
        lines.append(
            "| {slide_no} | {thumbnail} | {headline} | {layout} | {proof} | "
            "{flags} | {status} | {note} |".format(
                slide_no=slide.slide_no,
                thumbnail=thumbnail,
                headline=str(slide.screen_headline).replace("|", "\\|"),
                layout=str(slide.layout_intent).replace("|", "\\|"),
                proof=str(slide.proof_object_type).replace("|", "\\|"),
                flags=flags.replace("|", "\\|"),
                status=slide.contact_sheet_review_status,
                note=slide.reviewer_note,
            )
        )
    result.report_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _status_for_result(
    *,
    warnings: list[str],
    thumbnail_generation: ThumbnailGeneration | None,
    slide_count: int,
) -> str:
    if any(warning.startswith("missing_pptx") for warning in warnings):
        return "missing_pptx"
    if slide_count == 0:
        return "pptx_unreadable"
    if thumbnail_generation is None:
        return "metadata_only"
    if thumbnail_generation.status == "generated":
        return "contact_sheet_generated"
    if thumbnail_generation.thumbnails:
        return "thumbnails_generated"
    if thumbnail_generation.pdf_path:
        return "pdf_generated"
    return "metadata_only_with_warning"


def render_contact_sheet_target(
    *,
    target: ContactSheetTarget,
    output_dir: Path,
    thumbnail_generator: ThumbnailGenerator | None = None,
) -> ContactSheetResult:
    deck_dir = output_dir / target.deck_id
    warnings: list[str] = []
    spec, visual_qa_deck, spec_warnings = _load_slide_spec(target.slide_spec_path, deck_dir)
    warnings.extend(spec_warnings)
    slide_count = 0
    if not target.pptx_path.exists():
        warnings.append(f"missing_pptx: {_display_path(target.pptx_path)}")
        pptx_warning = None
    else:
        slide_count, pptx_warning = _slide_count_for_pptx(target.pptx_path)
        if pptx_warning:
            warnings.append(pptx_warning)

    thumbnail_generation: ThumbnailGeneration | None = None
    if target.pptx_path.exists() and slide_count:
        generator = thumbnail_generator or default_thumbnail_generator
        thumbnail_generation = generator(target, output_dir)
        warnings.extend(thumbnail_generation.warnings)

    thumbnails = thumbnail_generation.thumbnails if thumbnail_generation else []
    spec_slide_count = len(_slides_by_no(spec)) if spec is not None else None
    if spec_slide_count is not None and slide_count and spec_slide_count != slide_count:
        warnings.append(
            f"slide_count_mismatch: pptx={slide_count}, slide_spec={spec_slide_count}"
        )
    if thumbnails and slide_count and len(thumbnails) != slide_count:
        warnings.append(
            f"thumbnail_count_mismatch: pptx={slide_count}, thumbnails={len(thumbnails)}"
        )
    if not thumbnails:
        warnings.append("thumbnail_missing")

    result = ContactSheetResult(
        target=target,
        status=_status_for_result(
            warnings=warnings,
            thumbnail_generation=thumbnail_generation,
            slide_count=slide_count,
        ),
        slide_count=slide_count,
        spec_slide_count=spec_slide_count,
        thumbnail_count=len(thumbnails),
        thumbnail_status=thumbnail_generation.status if thumbnail_generation else "not_attempted",
        backend=thumbnail_generation.backend if thumbnail_generation else "none",
        pdf_path=thumbnail_generation.pdf_path if thumbnail_generation else None,
        contact_sheet_path=(
            thumbnail_generation.contact_sheet_path if thumbnail_generation else None
        ),
        contact_sheet_pdf_path=(
            thumbnail_generation.contact_sheet_pdf_path if thumbnail_generation else None
        ),
        report_path=output_dir / f"{target.deck_id}_contact_sheet.md",
        warnings=warnings,
        slides=_build_slide_rows(
            slide_count=slide_count,
            spec=spec,
            visual_qa_deck=visual_qa_deck,
            thumbnails=thumbnails,
        ),
    )
    _write_deck_report(result)
    return result


def _write_summary(
    *,
    results: list[ContactSheetResult],
    output_path: Path,
) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    status_counts = Counter(result.status for result in results)
    failed = [
        result
        for result in results
        if result.status in {"missing_pptx", "pptx_unreadable"}
    ]
    lines = [
        "# PPTX Contact Sheet QA Summary",
        "",
        f"- generated_at: {datetime.now(UTC).isoformat()}",
        f"- deck_count: {len(results)}",
        f"- slide_count: {sum(result.slide_count for result in results)}",
        f"- generated_contact_sheets: {sum(1 for result in results if result.contact_sheet_path)}",
        f"- thumbnail_generation_status: {dict(status_counts)}",
        f"- failed_decks: {len(failed)}",
        "- This is visual review surface only.",
        "- No PPT content was modified.",
        "- No LLM/API calls.",
        "- Production readiness remains false.",
        "- Broadcast readiness remains false.",
        "",
        "## Decks",
        "",
        (
            "| deck | source | status | slides | thumbnails | backend | "
            "contact sheet | report | warnings |"
        ),
        "|---|---|---|---:|---:|---|---|---|---|",
    ]
    for result in results:
        warnings = ", ".join(result.warnings) if result.warnings else "-"
        lines.append(
            "| {deck} | {source} | {status} | {slides} | {thumbs} | {backend} | "
            "{sheet} | {report} | {warnings} |".format(
                deck=result.target.deck_id,
                source=result.target.source_kind,
                status=result.status,
                slides=result.slide_count,
                thumbs=result.thumbnail_count,
                backend=result.backend,
                sheet=_display_path(result.contact_sheet_path),
                report=_display_path(result.report_path),
                warnings=warnings.replace("|", "\\|"),
            )
        )
    lines.extend(
        [
            "",
            "## Review Notes",
            "",
            "- Each deck report has slide-level `contact_sheet_review_status: unchecked`.",
            "- `reviewer_note` is intentionally blank for human review.",
            "- Lightweight heuristic flags are limited to missing/render/count issues.",
            (
                "- No OCR, AI layout judgment, style scoring, image insertion, "
                "or chart generation is performed."
            ),
        ]
    )
    output_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def render_pptx_contact_sheet(
    *,
    targets: list[ContactSheetTarget] | None = None,
    output_dir: Path = DEFAULT_OUTPUT_DIR,
    summary_output_path: Path = DEFAULT_SUMMARY_PATH,
    review_summary_output_path: Path | None = DEFAULT_REVIEW_SUMMARY_PATH,
    thumbnail_generator: ThumbnailGenerator | None = None,
) -> list[ContactSheetResult]:
    output_dir.mkdir(parents=True, exist_ok=True)
    selected_targets = targets or default_targets()
    results = [
        render_contact_sheet_target(
            target=target,
            output_dir=output_dir,
            thumbnail_generator=thumbnail_generator,
        )
        for target in selected_targets
    ]
    _write_summary(results=results, output_path=summary_output_path)
    if review_summary_output_path is not None:
        _write_summary(results=results, output_path=review_summary_output_path)
    return results


@app.callback(invoke_without_command=True)
def main(
    pptx_paths: Annotated[
        list[Path] | None,
        typer.Option("--pptx", help="Optional explicit PPTX path. May be repeated."),
    ] = None,
    output_dir: Annotated[
        Path,
        typer.Option("--output-dir", help="Directory for contact sheet QA outputs."),
    ] = DEFAULT_OUTPUT_DIR,
    summary_output_path: Annotated[
        Path,
        typer.Option("--summary-output", help="Markdown summary report path."),
    ] = DEFAULT_SUMMARY_PATH,
    review_summary_output_path: Annotated[
        Path | None,
        typer.Option("--review-summary-output", help="GitHub-visible summary mirror."),
    ] = DEFAULT_REVIEW_SUMMARY_PATH,
    include_direct: Annotated[
        bool,
        typer.Option(
            "--include-direct/--adapter-only",
            help="Include Anny direct live PPTX drafts in the default target set.",
        ),
    ] = True,
    direct_run_id: Annotated[
        str,
        typer.Option("--direct-run-id", help="Anny direct live run id for default targets."),
    ] = DEFAULT_DIRECT_RUN_ID,
    direct_output_root: Annotated[
        Path,
        typer.Option("--direct-output-root", help="Root containing direct run outputs."),
    ] = DEFAULT_DIRECT_OUTPUT_ROOT,
) -> None:
    """Write review-only PPTX contact sheet QA reports."""
    targets = (
        _custom_targets(pptx_paths)
        if pptx_paths
        else default_targets(
            include_direct=include_direct,
            direct_run_id=direct_run_id,
            direct_output_root=direct_output_root,
        )
    )
    results = render_pptx_contact_sheet(
        targets=targets,
        output_dir=output_dir,
        summary_output_path=summary_output_path,
        review_summary_output_path=review_summary_output_path,
    )
    console.print(
        "[green]Wrote PPTX contact sheet QA for "
        f"{len(results)} deck(s) to {output_dir}.[/green]"
    )


if __name__ == "__main__":
    app()
