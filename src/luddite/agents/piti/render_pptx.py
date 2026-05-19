"""Render Piti deck plans into editable draft PPTX skeletons."""

from __future__ import annotations

import json
import re
from collections import Counter
from dataclasses import dataclass
from datetime import UTC, date, datetime
from pathlib import Path
from typing import Annotated, Any
from urllib.parse import urlparse

import typer
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from rich.console import Console

from luddite import paths
from luddite.agents.piti.build_deck_plan_from_storyline import validate_deck_plan
from luddite.parsers.parse_pptx import parse_presentation
from luddite.utils.urls import canonicalize_url

app = typer.Typer(no_args_is_help=False)
console = Console()

SLIDE_W = 13.333
SLIDE_H = 7.5

DEFAULT_REPORT_PATH = (
    paths.REPORTS_DIR / f"piti_pptx_render_report_{date.today().isoformat()}.md"
)
DEFAULT_STYLE_PROFILE_PATH = paths.STYLE_PROFILES_DIR / "syukaworld_ppt_style_profile.json"


@dataclass(frozen=True)
class PptxRenderCase:
    deck_plan_path: Path
    output_path: Path
    style_profile_path: Path | None = None


@dataclass(frozen=True)
class PptxStyle:
    profile_path: Path | None
    loaded: bool
    slide_width_in: float
    slide_height_in: float
    font_family: str
    fallback_font: str
    accent: RGBColor
    layout_patterns: dict[str, dict[str, Any]]
    raw_profile: dict[str, Any]


@dataclass(frozen=True)
class Rect:
    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height


@dataclass(frozen=True)
class ProofObject:
    type: str
    required: bool
    screen_position: str
    display_label: str
    source_url: str | None
    image_url: str | None
    copyright_risk: bool
    manual_insert_required: bool

    def as_dict(self) -> dict[str, Any]:
        return {
            "type": self.type,
            "required": self.required,
            "screen_position": self.screen_position,
            "display_label": self.display_label,
            "source_url": self.source_url,
            "image_url": self.image_url,
            "copyright_risk": self.copyright_risk,
            "manual_insert_required": self.manual_insert_required,
        }


DEFAULT_CASES = [
    PptxRenderCase(
        deck_plan_path=paths.PITI_DECK_PLANS_DIR / "ai_knowledge_institution_deck_plan.json",
        output_path=paths.PPTX_OUTPUT_DIR / "ai_knowledge_institution_draft.pptx",
    ),
    PptxRenderCase(
        deck_plan_path=paths.PITI_DECK_PLANS_DIR / "productive_finance_policy_deck_plan.json",
        output_path=paths.PPTX_OUTPUT_DIR / "productive_finance_policy_draft.pptx",
    ),
]

STYLED_CASES = [
    PptxRenderCase(
        deck_plan_path=paths.PITI_DECK_PLANS_DIR / "ai_knowledge_institution_deck_plan.json",
        output_path=paths.PPTX_OUTPUT_DIR / "ai_knowledge_institution_styled_draft.pptx",
        style_profile_path=DEFAULT_STYLE_PROFILE_PATH,
    ),
    PptxRenderCase(
        deck_plan_path=paths.PITI_DECK_PLANS_DIR / "productive_finance_policy_deck_plan.json",
        output_path=paths.PPTX_OUTPUT_DIR / "productive_finance_policy_styled_draft.pptx",
        style_profile_path=DEFAULT_STYLE_PROFILE_PATH,
    ),
]


def _display_path(path: Path) -> str:
    try:
        return str(path.resolve().relative_to(paths.REPO_ROOT))
    except ValueError:
        return str(path)


THEME = {
    "ink": RGBColor(24, 28, 35),
    "muted": RGBColor(86, 96, 112),
    "line": RGBColor(210, 216, 224),
    "paper": RGBColor(248, 249, 250),
    "accent": RGBColor(14, 121, 130),
    "accent_dark": RGBColor(15, 86, 96),
    "warn": RGBColor(191, 111, 36),
    "light": RGBColor(255, 255, 255),
}
SYUKA_RED = RGBColor(255, 0, 0)
SCREEN_BLACK = RGBColor(0, 0, 0)
CHART_DARK = RGBColor(55, 55, 55)
CONTENT_HEADLINE_BOX = (0.626, 0.39, 12.24, 0.57)
PROOF_LABELS = {
    "none": "",
    "image": "[이미지]",
    "chart": "[차트]",
    "table": "[표]",
    "article_quote": "[기사 인용]",
    "logo": "[로고]",
    "screenshot": "[기사 캡처]",
    "diagram": "[도식]",
    "map": "[지도]",
    "person_photo": "[인물 사진]",
    "generated_image_candidate": "[AI 이미지]",
}
LEFT_PROOF_TYPES = {
    "image",
    "logo",
    "person_photo",
    "screenshot",
    "diagram",
    "generated_image_candidate",
    "article_quote",
}


def _load_json(path: Path) -> dict[str, Any]:
    with path.open(encoding="utf-8") as source:
        payload = json.load(source)
    if not isinstance(payload, dict):
        raise ValueError(f"Expected JSON object at {path}")
    return payload


def _rgb_from_hex(value: str | None, fallback: RGBColor) -> RGBColor:
    if not value or not value.startswith("#") or len(value) != 7:
        return fallback
    try:
        return RGBColor(
            int(value[1:3], 16),
            int(value[3:5], 16),
            int(value[5:7], 16),
        )
    except ValueError:
        return fallback


def _cm_to_in(value: float | int | None, fallback: float) -> float:
    if value is None:
        return fallback
    return round(float(value) / 2.54, 3)


def _style_font(profile: dict[str, Any]) -> str:
    font_resolution = profile.get("font_resolution") or {}
    return (
        font_resolution.get("fallback_font")
        or (font_resolution.get("resolved_font_candidates") or [None])[0]
        or "Malgun Gothic"
    )


def load_style_profile(path: Path | None) -> PptxStyle:
    if path is None or not path.exists():
        return PptxStyle(
            profile_path=path,
            loaded=False,
            slide_width_in=SLIDE_W,
            slide_height_in=SLIDE_H,
            font_family="Malgun Gothic",
            fallback_font="Malgun Gothic",
            accent=RGBColor(14, 121, 130),
            layout_patterns={},
            raw_profile={},
        )
    profile = _load_json(path)
    slide_size = profile.get("slide_size") or {}
    layout_patterns = profile.get("layout_patterns") or {}
    common_colors = profile.get("common_colors") or []
    accent_value = common_colors[0]["value"] if common_colors else "#FF0000"
    return PptxStyle(
        profile_path=path,
        loaded=True,
        slide_width_in=float(slide_size.get("width_in") or SLIDE_W),
        slide_height_in=float(slide_size.get("height_in") or SLIDE_H),
        font_family=_style_font(profile),
        fallback_font="Malgun Gothic",
        accent=_rgb_from_hex(accent_value, RGBColor(255, 0, 0)),
        layout_patterns=layout_patterns,
        raw_profile=profile,
    )


def _layout_box(
    style: PptxStyle,
    layout_type: str,
    fallback: tuple[float, float, float, float],
) -> tuple[float, float, float, float]:
    pattern = style.layout_patterns.get(layout_type) or {}
    return (
        _cm_to_in(pattern.get("x_cm_median"), fallback[0]),
        _cm_to_in(pattern.get("y_cm_median"), fallback[1]),
        _cm_to_in(pattern.get("w_cm_median"), fallback[2]),
        _cm_to_in(pattern.get("h_cm_median"), fallback[3]),
    )


def _layout_font_size(style: PptxStyle, layout_type: str, fallback: int) -> int:
    pattern = style.layout_patterns.get(layout_type) or {}
    value = pattern.get("font_size_pt_median")
    return int(round(float(value))) if value else fallback


def _layout_color(style: PptxStyle, layout_type: str, fallback: RGBColor) -> RGBColor:
    pattern = style.layout_patterns.get(layout_type) or {}
    return _rgb_from_hex(pattern.get("font_color_top"), fallback)


def _visual_kind(slide_plan: dict[str, Any]) -> str:
    visual_plan = slide_plan.get("visual_plan") or {}
    return str(visual_plan.get("kind") or "none")


def _first_url(slide_plan: dict[str, Any], key: str) -> str | None:
    for url in _as_list(slide_plan.get(key)):
        text = str(url).strip()
        if text:
            return text
    return None


def _has_source_urls(slide_plan: dict[str, Any]) -> bool:
    return _first_url(slide_plan, "source_urls") is not None


def _is_internal_or_structural_slide(slide_plan: dict[str, Any]) -> bool:
    layout_type = str(slide_plan.get("layout_type") or "")
    slide_type = str(slide_plan.get("slide_type") or "")
    return layout_type in {
        "title",
        "section_title",
        "checklist",
        "appendix_checklist",
    } or slide_type in {"production_checklist", "source_heavy"}


def _source_name_from_url(url: str | None) -> str:
    if not url:
        return "Source"
    host = urlparse(url).netloc or url
    host = host.removeprefix("www.")
    known = {
        "bbc.com": "BBC",
        "bbc.co.uk": "BBC",
        "microsoft.com": "Microsoft",
        "unesco.org": "UNESCO",
        "oecd.org": "OECD",
        "fsc.go.kr": "금융위원회",
        "kdb.co.kr": "산업은행",
        "korea.kr": "정책브리핑",
        "bis.org": "BIS",
    }
    for domain, label in known.items():
        if domain in host:
            return label
    return host.split("/")[0]


def _short_text(value: str, limit: int = 34) -> str:
    text = re.sub(r"\s+", " ", value).strip()
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "…"


def _proof_object(slide_plan: dict[str, Any]) -> ProofObject:
    visual_plan = slide_plan.get("visual_plan") or {}
    kind = _visual_kind(slide_plan)
    layout_type = str(slide_plan.get("layout_type") or "headline_body")
    slide_type = str(slide_plan.get("slide_type") or "")
    proof_type = "none"

    if _is_chart_table_slide(slide_plan) or kind == "chart_candidate":
        proof_type = "table" if layout_type in {"table", "data_table"} else "chart"
    elif kind == "diagram":
        proof_type = "diagram"
    elif kind == "screenshot_candidate":
        proof_type = "screenshot"
    elif kind in {"photo_candidate", "image_placeholder"}:
        proof_type = "image"
    elif kind == "ai_image_prompt":
        proof_type = "generated_image_candidate"
    elif kind == "manual":
        if layout_type == "quote" or slide_type == "quote":
            proof_type = "article_quote"
        elif layout_type in {"image_placeholder", "image_heavy"}:
            proof_type = "image"
        elif layout_type in {"chart_placeholder", "table", "data_table"}:
            proof_type = "chart"
        elif _has_source_urls(slide_plan) and not _is_internal_or_structural_slide(slide_plan):
            proof_type = "article_quote"
    elif layout_type == "quote" or slide_type == "quote":
        proof_type = "article_quote"

    screen_position = "none"
    if proof_type in {"chart", "table"}:
        screen_position = "full_width_chart"
    elif proof_type in LEFT_PROOF_TYPES:
        screen_position = "left_half"
    elif proof_type != "none":
        screen_position = "center_large"

    source_url = _first_url(slide_plan, "source_urls")
    image_url = _first_url(slide_plan, "image_urls")
    required = proof_type != "none"
    return ProofObject(
        type=proof_type,
        required=required,
        screen_position=screen_position,
        display_label=PROOF_LABELS.get(proof_type, "[증거물]"),
        source_url=source_url,
        image_url=image_url,
        copyright_risk=bool(visual_plan.get("copyright_risk")),
        manual_insert_required=bool(required and (kind == "manual" or not image_url)),
    )


def _has_visual(slide_plan: dict[str, Any]) -> bool:
    return _proof_object(slide_plan).type != "none"


def _screen_placeholder_visible(slide_plan: dict[str, Any], style: PptxStyle | None) -> bool:
    proof = _proof_object(slide_plan)
    if proof.type in {"none", "chart", "table"}:
        return False
    return True


def _has_screen_visual(slide_plan: dict[str, Any], style: PptxStyle | None) -> bool:
    return _screen_placeholder_visible(slide_plan, style)


def _is_chart_table_slide(slide_plan: dict[str, Any]) -> bool:
    layout_type = slide_plan.get("layout_type")
    if layout_type in {"chart_placeholder", "table", "data_table"}:
        return True
    if layout_type == "comparison":
        text = "\n".join(_body_lines(slide_plan))
        return bool(re.search(r"\d", text)) and any(
            token in text.lower()
            for token in ["%", "달러", "조", "억", "won", "usd", "rate", "ratio"]
        )
    return False


def _prefers_left_image(slide_plan: dict[str, Any], style: PptxStyle | None) -> bool:
    if not _screen_placeholder_visible(slide_plan, style):
        return False
    return _proof_object(slide_plan).screen_position == "left_half"


def _is_english_line(line: str) -> bool:
    ascii_letters = len(re.findall(r"[A-Za-z]", line))
    hangul = len(re.findall(r"[가-힣]", line))
    return ascii_letters >= 5 and ascii_letters > hangul


def _is_korean_line(line: str) -> bool:
    return bool(re.search(r"[가-힣]", line))


def _is_bilingual_quote_slide(slide_plan: dict[str, Any]) -> bool:
    body = _body_lines(slide_plan)
    if slide_plan.get("layout_type") != "quote":
        return False
    return any(_is_english_line(line) for line in body) and any(
        _is_korean_line(line) for line in body
    )


def _body_font_size(
    style: PptxStyle,
    body: list[str],
    fallback: int,
    *,
    has_visual: bool = False,
) -> int:
    if not style.loaded:
        return fallback
    max_len = max((len(line) for line in body), default=0)
    if len(body) >= 5 or max_len > 110:
        return 20
    if has_visual and (len(body) >= 3 or max_len > 70):
        return 24
    if len(body) >= 4 or max_len > 80:
        return 24
    return 28


def _body_line_estimate(slide: dict[str, Any], font_size: int, *, has_visual: bool) -> int:
    body = _screen_body_lines(slide)
    if not body:
        return 0
    chars_per_line = 28 if has_visual else 44
    if font_size <= 20:
        chars_per_line += 8
    elif font_size >= 28:
        chars_per_line -= 6
    estimate = 0
    for line in body:
        estimate += max(1, (len(line) + chars_per_line - 1) // chars_per_line)
    return estimate


def _is_visually_dense(slide: dict[str, Any], font_size: int) -> bool:
    has_visual = _has_visual(slide)
    estimated_lines = _body_line_estimate(slide, font_size, has_visual=has_visual)
    if font_size >= 28 and estimated_lines >= 4:
        return True
    if has_visual and estimated_lines >= 3:
        return True
    return estimated_lines >= 6


def _rects_overlap(first: Rect, second: Rect) -> bool:
    return not (
        first.right <= second.left
        or second.right <= first.left
        or first.bottom <= second.top
        or second.bottom <= first.top
    )


def _as_list(value: Any) -> list[Any]:
    return value if isinstance(value, list) else []


def _body_lines(slide: dict[str, Any]) -> list[str]:
    return [str(item) for item in _as_list(slide.get("body")) if str(item).strip()]


def _fill(shape: Any, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _line(shape: Any, color: RGBColor = THEME["line"], width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def _no_line(shape: Any) -> None:
    shape.line.fill.background()


def _textbox(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    *,
    font_size: int = 24,
    bold: bool = False,
    underline: bool = False,
    color: RGBColor = THEME["ink"],
    align: PP_ALIGN | None = None,
    style: PptxStyle | None = None,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = style.font_family if style and style.loaded else "Malgun Gothic"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.underline = underline
        run.font.color.rgb = color
    return shape


def _body_box(
    slide: Any,
    body: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 22,
    color: RGBColor = THEME["ink"],
    line_colors: list[RGBColor] | None = None,
    bold: bool = False,
    underline: bool = False,
    style: PptxStyle | None = None,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    if not body:
        paragraph = frame.paragraphs[0]
        paragraph.text = ""
        return shape
    for index, line in enumerate(body):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.space_after = Pt(8)
        line_color = line_colors[index] if line_colors and index < len(line_colors) else color
        for run in paragraph.runs:
            run.font.name = style.font_family if style and style.loaded else "Malgun Gothic"
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.underline = underline
            run.font.color.rgb = line_color
    return shape


def _screen_body_lines(slide: dict[str, Any]) -> list[str]:
    body = _body_lines(slide)
    layout_type = slide.get("layout_type")
    limit = 4 if layout_type == "quote" and _is_bilingual_quote_slide(slide) else 3
    if layout_type in {"chart_placeholder", "table", "data_table"}:
        limit = 3
    proof_type = _proof_object(slide).type
    if proof_type != "none" and proof_type not in {"chart", "table", "article_quote"}:
        limit = min(limit, 2)
    if proof_type == "article_quote" and layout_type != "quote":
        limit = min(limit, 2)
    return body[:limit]


def _overflow_body_lines(slide: dict[str, Any]) -> list[str]:
    body = _body_lines(slide)
    return body[len(_screen_body_lines(slide)) :]


def _quote_line_colors(body: list[str]) -> list[RGBColor]:
    colors = []
    for line in body:
        if _is_korean_line(line) and not _is_english_line(line):
            colors.append(SYUKA_RED)
        else:
            colors.append(SCREEN_BLACK)
    return colors


def _placeholder(
    slide: Any,
    slide_plan: dict[str, Any],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    style: PptxStyle | None = None,
) -> None:
    proof = _proof_object(slide_plan)
    if proof.type == "none":
        return
    label = proof.display_label
    shape = slide.shapes.add_shape(
        1,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    placeholder_fill = (
        RGBColor(245, 245, 245) if style and style.loaded else RGBColor(235, 241, 244)
    )
    placeholder_line = (
        RGBColor(190, 190, 190) if style and style.loaded else RGBColor(160, 174, 184)
    )
    _fill(shape, placeholder_fill)
    _line(shape, placeholder_line, 1.0)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = label
    paragraph.alignment = PP_ALIGN.CENTER
    for run in paragraph.runs:
        run.font.name = style.font_family if style and style.loaded else "Malgun Gothic"
        run.font.size = Pt(13 if style and style.loaded else 16)
        run.font.color.rgb = THEME["muted"]
        run.font.bold = proof.type == "article_quote"
    if proof.type == "article_quote":
        source = _source_name_from_url(proof.source_url)
        title = _short_text(str(slide_plan.get("headline") or ""), 38)
        for text in [source, title]:
            extra = frame.add_paragraph()
            extra.text = text
            extra.alignment = PP_ALIGN.CENTER
            extra.space_before = Pt(8 if text == source else 4)
            for run in extra.runs:
                run.font.name = style.font_family if style and style.loaded else "Malgun Gothic"
                run.font.size = Pt(18 if text == source else 14)
                run.font.bold = text == source
                run.font.color.rgb = SCREEN_BLACK if text == source else THEME["muted"]


def _add_footer(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    if style.loaded:
        return
    flags = []
    if slide_plan.get("needs_source"):
        flags.append("needs_source")
    if slide_plan.get("needs_fact_check"):
        flags.append("needs_fact_check")
    if slide_plan.get("required_before_broadcast"):
        flags.append("before_broadcast")
    flag_text = " | ".join(flags) if flags else "draft skeleton"
    _textbox(
        slide,
        0.45,
        7.12,
        10.8,
        0.22,
        flag_text,
        font_size=8,
        color=THEME["muted"],
        style=style,
    )
    _textbox(
        slide,
        11.6,
        7.12,
        1.25,
        0.22,
        f"{slide_plan.get('slide_no', ''):02}",
        font_size=8,
        color=THEME["muted"],
        align=PP_ALIGN.RIGHT,
        style=style,
    )


def _slide_background(slide: Any, *, dark: bool = False) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = THEME["ink"] if dark else THEME["paper"]


def _render_title(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    _slide_background(slide, dark=False if style.loaded else True)
    left, top, width, height = _layout_box(style, "title", (0.8, 1.35, 11.8, 1.7))
    headline_color = SCREEN_BLACK if style.loaded else THEME["light"]
    _textbox(
        slide,
        left,
        top,
        width,
        height,
        str(slide_plan.get("headline") or ""),
        font_size=54 if style.loaded else _layout_font_size(style, "title", 38),
        bold=True,
        color=headline_color,
        style=style,
    )
    body = _screen_body_lines(slide_plan)
    if not style.loaded:
        _body_box(
            slide,
            body,
            left + 0.1,
            min(top + 2.1, 6.2),
            min(width, 10.8),
            1.0,
            font_size=_body_font_size(style, body, 22, has_visual=False),
            color=THEME["ink"],
            style=style,
        )
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not style.loaded and run.font.color.rgb == THEME["ink"]:
                        run.font.color.rgb = RGBColor(225, 233, 238)


def _render_section_title(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    _slide_background(slide, dark=False)
    if not style.loaded:
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        _fill(bar, THEME["accent_dark"])
        _no_line(bar)
    left, top, width, height = _layout_box(style, "section_title", (0.85, 2.35, 11.6, 1.2))
    _textbox(
        slide,
        left,
        top,
        width,
        height,
        str(slide_plan.get("headline") or ""),
        font_size=_layout_font_size(style, "section_title", 34),
        bold=True,
        color=THEME["ink"] if style.loaded else THEME["light"],
        style=style,
    )
    body = _screen_body_lines(slide_plan)
    if not style.loaded:
        _body_box(
            slide,
            body,
            left,
            min(top + 1.35, 6.2),
            min(width, 10.8),
            1.0,
            font_size=_body_font_size(style, body, 20, has_visual=False),
            color=THEME["ink"],
            style=style,
        )
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not style.loaded and run.font.color.rgb == THEME["ink"]:
                        run.font.color.rgb = RGBColor(226, 238, 241)


def _render_big_headline(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    _slide_background(slide)
    left, top, width, height = (
        CONTENT_HEADLINE_BOX
        if style.loaded
        else _layout_box(style, "big_headline", (0.7, 1.0, 11.8, 1.4))
    )
    _textbox(
        slide,
        left,
        top,
        width,
        height,
        str(slide_plan.get("headline") or ""),
        font_size=28 if style.loaded else _layout_font_size(style, "big_headline", 34),
        bold=True,
        color=SYUKA_RED if style.loaded else THEME["ink"],
        style=style,
    )
    body = _screen_body_lines(slide_plan)
    has_visual = _has_screen_visual(slide_plan, style)
    left_image = _prefers_left_image(slide_plan, style)
    body_left = 6.55 if style.loaded and left_image else 0.78
    body_top = 2.05 if style.loaded else max(top + height + 0.35, 2.8)
    if style.loaded and left_image:
        body_width = 5.8
    elif has_visual:
        body_width = 7.4
    else:
        body_width = 10.8
    _body_box(
        slide,
        body,
        body_left,
        body_top,
        body_width,
        2.9,
        font_size=_body_font_size(style, body, 24, has_visual=has_visual),
        color=SCREEN_BLACK if style.loaded else THEME["ink"],
        style=style,
    )
    if style.loaded and left_image:
        _placeholder(slide, slide_plan, 0.78, 1.55, 5.35, 4.5, style=style)
    else:
        _placeholder(slide, slide_plan, 9.25, 4.8, 2.8, 1.2, style=style)


def _render_headline_body(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    _slide_background(slide)
    has_visual = _has_screen_visual(slide_plan, style)
    left_image = _prefers_left_image(slide_plan, style)
    left, top, width, height = _layout_box(style, "headline_body", (0.65, 1.75, 7.7, 4.7))
    if style.loaded and has_visual:
        width = 5.9 if left_image else min(width, 7.35)
        left = 6.55 if left_image else left
        height = min(height, 3.85)
    headline_left, headline_top, headline_width, headline_height = (
        CONTENT_HEADLINE_BOX if style.loaded else (left, 0.55, min(width, 12.0), 0.9)
    )
    headline_color = SYUKA_RED if style.loaded else THEME["ink"]
    _textbox(
        slide,
        headline_left,
        headline_top,
        headline_width,
        headline_height,
        str(slide_plan.get("headline") or ""),
        font_size=_layout_font_size(style, "headline_body", 28),
        bold=True,
        color=headline_color,
        style=style,
    )
    body = _screen_body_lines(slide_plan)
    _body_box(
        slide,
        body,
        left,
        2.05 if style.loaded and not left_image else top,
        width,
        height,
        font_size=_body_font_size(style, body, 20, has_visual=has_visual),
        color=SCREEN_BLACK if style.loaded else THEME["ink"],
        style=style,
    )
    if style.loaded and left_image:
        _placeholder(slide, slide_plan, 0.78, 1.55, 5.35, 4.5, style=style)
    elif style.loaded and has_visual:
        _placeholder(slide, slide_plan, 8.75, 1.72, 3.65, 3.95, style=style)
    elif style.loaded:
        _placeholder(slide, slide_plan, 8.95, 1.55, 3.6, 4.65, style=style)
    else:
        _placeholder(slide, slide_plan, 8.55, 1.55, 4.05, 4.65, style=style)


def _render_quote(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    _slide_background(slide)
    left, top, width, height = _layout_box(style, "quote", (0.8, 1.0, 11.7, 4.9))
    left_image = _prefers_left_image(slide_plan, style)
    if style.loaded and left_image:
        left, top, width, height = (6.35, 1.35, 6.1, 4.85)
        _placeholder(slide, slide_plan, 0.75, 1.65, 5.25, 4.15, style=style)
    box = slide.shapes.add_shape(
        1,
        Inches(left),
        Inches(max(0.8, top - 0.75)),
        Inches(min(11.7, width + 0.6)),
        Inches(max(3.4, height + 2.2)),
    )
    _fill(box, THEME["light"])
    _line(box, THEME["line"], 1.2)
    headline_left, headline_top, headline_width, headline_height = (
        CONTENT_HEADLINE_BOX
        if style.loaded
        else (left, max(1.0, top - 0.55), min(11.0, width), 0.95)
    )
    _textbox(
        slide,
        headline_left,
        headline_top,
        headline_width,
        headline_height,
        str(slide_plan.get("headline") or ""),
        font_size=28 if style.loaded else _layout_font_size(style, "quote", 28),
        bold=True,
        color=SYUKA_RED if style.loaded else THEME["ink"],
        style=style,
    )
    proof = _proof_object(slide_plan)
    if style.loaded and proof.type == "article_quote" and proof.source_url:
        host = urlparse(proof.source_url).netloc or proof.source_url
        _textbox(
            slide,
            9.1,
            0.68,
            3.25,
            0.32,
            host,
            font_size=16,
            bold=True,
            color=SCREEN_BLACK,
            align=PP_ALIGN.RIGHT,
            style=style,
        )
    body = _screen_body_lines(slide_plan)
    line_colors = _quote_line_colors(body) if _is_bilingual_quote_slide(slide_plan) else None
    _body_box(
        slide,
        body,
        left,
        2.05 if style.loaded and not left_image else top + 0.8,
        min(10.6, width),
        2.7,
        font_size=28 if style.loaded else _body_font_size(style, body, 20, has_visual=False),
        color=SCREEN_BLACK if style.loaded else THEME["ink"],
        line_colors=line_colors,
        style=style,
    )


def _render_question(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    _slide_background(slide)
    left, top, width, height = CONTENT_HEADLINE_BOX if style.loaded else (1.0, 1.35, 11.3, 1.4)
    _textbox(
        slide,
        left,
        top,
        width,
        height,
        str(slide_plan.get("headline") or ""),
        font_size=28 if style.loaded else 32,
        bold=True,
        color=SYUKA_RED if style.loaded else THEME["ink"],
        style=style,
    )
    body = _screen_body_lines(slide_plan)
    left_image = _prefers_left_image(slide_plan, style)
    if style.loaded and left_image:
        _placeholder(slide, slide_plan, 0.78, 1.55, 5.35, 4.5, style=style)
    _body_box(
        slide,
        body,
        6.55 if style.loaded and left_image else (1.55 if not style.loaded else 0.78),
        3.2 if not style.loaded else 2.3,
        5.8 if style.loaded and left_image else (10.2 if not style.loaded else 11.65),
        2.2 if not style.loaded else 3.2,
        font_size=_body_font_size(style, body, 23, has_visual=left_image),
        color=SCREEN_BLACK if style.loaded else THEME["ink"],
        style=style,
    )


def _render_comparison(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    _slide_background(slide)
    left, top, width, height = CONTENT_HEADLINE_BOX if style.loaded else (0.65, 0.55, 12.0, 0.85)
    _textbox(
        slide,
        left,
        top,
        width,
        height,
        str(slide_plan.get("headline") or ""),
        font_size=28 if style.loaded else 27,
        bold=True,
        color=SYUKA_RED if style.loaded else THEME["ink"],
        style=style,
    )
    body = _screen_body_lines(slide_plan)
    if style.loaded and _has_screen_visual(slide_plan, style):
        _placeholder(slide, slide_plan, 0.78, 1.55, 5.35, 4.5, style=style)
        _body_box(
            slide,
            body,
            6.55,
            1.65,
            5.8,
            3.65,
            font_size=_body_font_size(style, body, 19, has_visual=True),
            color=SCREEN_BLACK,
            style=style,
        )
        return
    midpoint = max(1, (len(body) + 1) // 2)
    left_body = body[:midpoint]
    right_body = body[midpoint:]
    box_top = 1.75 if not style.loaded else 1.55
    left = slide.shapes.add_shape(1, Inches(0.8), Inches(box_top), Inches(5.75), Inches(4.45))
    right = slide.shapes.add_shape(1, Inches(6.75), Inches(box_top), Inches(5.75), Inches(4.45))
    _fill(left, THEME["light"])
    _fill(right, RGBColor(239, 246, 247))
    _line(left)
    _line(right, RGBColor(170, 202, 206))
    _body_box(
        slide,
        left_body,
        1.05,
        box_top + 0.3,
        5.25,
        3.7,
        font_size=_body_font_size(style, left_body, 19, has_visual=False),
        color=SCREEN_BLACK if style.loaded else THEME["ink"],
        style=style,
    )
    _body_box(
        slide,
        right_body,
        7.0,
        box_top + 0.3,
        5.25,
        3.7,
        font_size=_body_font_size(style, right_body, 19, has_visual=False),
        color=SCREEN_BLACK if style.loaded else THEME["ink"],
        style=style,
    )


def _source_label(slide_plan: dict[str, Any]) -> str:
    urls = [str(url) for url in _as_list(slide_plan.get("source_urls")) if str(url).strip()]
    if not urls:
        return "(출처: 확인 필요)"
    host = urlparse(urls[0]).netloc or urls[0]
    return f"(출처: {host}; speaker notes 참조)"


def _chart_title(slide_plan: dict[str, Any]) -> str:
    body = _body_lines(slide_plan)
    if body:
        return body[0]
    return str(slide_plan.get("headline") or "차트/표 후보")


def _chart_label_candidates(slide_plan: dict[str, Any]) -> list[str]:
    body = _body_lines(slide_plan)
    candidates = body[1:] if len(body) > 1 else body
    labels = [line for line in candidates if line.strip()]
    return labels[:6] or ["데이터 라벨", "비교값", "확인 필요"]


def _render_reference_chart_skeleton(
    slide: Any,
    slide_plan: dict[str, Any],
    style: PptxStyle,
) -> None:
    chart_left = 1.0
    chart_top = 2.0
    chart_width = 11.35
    chart_height = 4.2
    chart_area = slide.shapes.add_shape(
        1,
        Inches(chart_left),
        Inches(chart_top),
        Inches(chart_width),
        Inches(chart_height),
    )
    _fill(chart_area, THEME["light"])
    _line(chart_area, THEME["line"], 1.2)
    # Reference-like axes and bar/data-label placeholders. Real chart generation
    # stays out of scope; these editable boxes reserve the proof-object grammar.
    axis_y = chart_top + chart_height - 0.58
    axis_x = chart_left + 0.72
    y_axis = slide.shapes.add_shape(
        1,
        Inches(axis_x),
        Inches(chart_top + 0.52),
        Inches(0.018),
        Inches(chart_height - 1.02),
    )
    _fill(y_axis, RGBColor(210, 210, 210))
    _no_line(y_axis)
    x_axis = slide.shapes.add_shape(
        1,
        Inches(axis_x),
        Inches(axis_y),
        Inches(chart_width - 1.25),
        Inches(0.018),
    )
    _fill(x_axis, RGBColor(210, 210, 210))
    _no_line(x_axis)
    labels = _chart_label_candidates(slide_plan)
    slot = (chart_width - 1.8) / max(len(labels), 1)
    for index, label in enumerate(labels):
        bar_height = 0.85 + (len(labels) - index) * 0.32
        bar_height = min(bar_height, 2.95)
        bar_left = axis_x + 0.36 + index * slot
        bar_width = min(0.72, max(0.48, slot * 0.42))
        bar_top = axis_y - bar_height
        bar = slide.shapes.add_shape(
            1,
            Inches(bar_left),
            Inches(bar_top),
            Inches(bar_width),
            Inches(bar_height),
        )
        _fill(bar, RGBColor(72, 114, 196))
        _no_line(bar)
        _textbox(
            slide,
            bar_left - 0.28,
            max(chart_top + 0.45, bar_top - 0.34),
            bar_width + 0.56,
            0.26,
            _short_text(label, 18),
            font_size=18,
            bold=True,
            color=CHART_DARK,
            align=PP_ALIGN.CENTER,
            style=style,
        )
        _textbox(
            slide,
            bar_left - 0.32,
            axis_y + 0.08,
            bar_width + 0.64,
            0.33,
            _short_text(label.split()[0], 10),
            font_size=14,
            bold=True,
            color=CHART_DARK,
            align=PP_ALIGN.CENTER,
            style=style,
        )


def _render_chart_table_layout(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    _slide_background(slide)
    _textbox(
        slide,
        *CONTENT_HEADLINE_BOX,
        str(slide_plan.get("headline") or ""),
        font_size=28,
        bold=True,
        color=SYUKA_RED,
        style=style,
    )
    _textbox(
        slide,
        0.95,
        1.25,
        11.7,
        0.65,
        _chart_title(slide_plan),
        font_size=28,
        bold=True,
        underline=True,
        color=SCREEN_BLACK,
        align=PP_ALIGN.CENTER,
        style=style,
    )
    _render_reference_chart_skeleton(slide, slide_plan, style)
    _textbox(
        slide,
        6.7,
        6.35,
        5.75,
        0.38,
        _source_label(slide_plan),
        font_size=20,
        underline=True,
        color=SCREEN_BLACK,
        align=PP_ALIGN.RIGHT,
        style=style,
    )


def _render_placeholder_layout(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    if style.loaded and _is_chart_table_slide(slide_plan):
        _render_chart_table_layout(slide, slide_plan, style)
        return
    _slide_background(slide)
    layout_type = "image_heavy"
    if slide_plan.get("layout_type") == "chart_placeholder":
        layout_type = "chart"
    left, top, width, height = _layout_box(style, layout_type, (1.15, 1.65, 11.0, 3.0))
    if style.loaded:
        left, top, width, height = (0.78, 1.55, 5.35, 4.5)
    _textbox(
        slide,
        CONTENT_HEADLINE_BOX[0] if style.loaded else left,
        CONTENT_HEADLINE_BOX[1] if style.loaded else max(0.45, top - 0.95),
        CONTENT_HEADLINE_BOX[2] if style.loaded else min(12.0, width),
        CONTENT_HEADLINE_BOX[3] if style.loaded else 0.85,
        str(slide_plan.get("headline") or ""),
        font_size=28 if style.loaded else _layout_font_size(style, layout_type, 27),
        bold=True,
        color=SYUKA_RED if style.loaded else THEME["ink"],
        style=style,
    )
    _placeholder(slide, slide_plan, left, top, width, max(1.8, height), style=style)
    body = _screen_body_lines(slide_plan)
    _body_box(
        slide,
        body,
        6.55 if style.loaded else left,
        1.65 if style.loaded else min(top + max(2.0, height) + 0.35, 5.4),
        5.8 if style.loaded else min(11.0, width),
        3.65 if style.loaded else 1.35,
        font_size=_body_font_size(style, body, 17, has_visual=True),
        color=SCREEN_BLACK if style.loaded else THEME["ink"],
        style=style,
    )


def _render_checklist(
    slide: Any,
    slide_plan: dict[str, Any],
    *,
    appendix: bool = False,
    style: PptxStyle,
) -> None:
    _slide_background(slide)
    prefix = "[내부 체크리스트] " if appendix else ""
    left, top, width, height = _layout_box(style, "checklist", (0.65, 0.55, 12.0, 0.85))
    if style.loaded and not appendix:
        left, top, width, height = CONTENT_HEADLINE_BOX
    _textbox(
        slide,
        left,
        top,
        width,
        height,
        prefix + str(slide_plan.get("headline") or ""),
        font_size=(
            28 if style.loaded and not appendix else _layout_font_size(style, "checklist", 26)
        ),
        bold=True,
        color=THEME["warn"] if appendix else (SYUKA_RED if style.loaded else THEME["ink"]),
        style=style,
    )
    body = _screen_body_lines(slide_plan)
    _body_box(
        slide,
        body,
        0.9,
        max(top + height + 0.25, 1.65),
        11.4,
        4.9,
        font_size=_body_font_size(style, body, 18, has_visual=False),
        color=SCREEN_BLACK if style.loaded else THEME["ink"],
        style=style,
    )


def _render_slide(prs: Presentation, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    layout_type = slide_plan.get("layout_type") or "headline_body"
    if style.loaded and _is_chart_table_slide(slide_plan):
        _render_chart_table_layout(slide, slide_plan, style)
    elif layout_type == "title":
        _render_title(slide, slide_plan, style)
    elif layout_type == "section_title":
        _render_section_title(slide, slide_plan, style)
    elif layout_type == "big_headline":
        _render_big_headline(slide, slide_plan, style)
    elif layout_type == "quote":
        _render_quote(slide, slide_plan, style)
    elif layout_type in {"question", "closing_question"}:
        _render_question(slide, slide_plan, style)
    elif layout_type == "comparison":
        _render_comparison(slide, slide_plan, style)
    elif layout_type in {"chart_placeholder", "image_placeholder", "timeline"}:
        _render_placeholder_layout(slide, slide_plan, style)
    elif layout_type in {"checklist", "appendix_checklist"}:
        _render_checklist(
            slide,
            slide_plan,
            appendix=layout_type == "appendix_checklist",
            style=style,
        )
    else:
        _render_headline_body(slide, slide_plan, style)
    _add_footer(slide, slide_plan, style)
    _set_notes(slide, slide_plan, style)


def _note_json(value: Any) -> str:
    return json.dumps(value, ensure_ascii=False, indent=2)


def _set_notes(slide: Any, slide_plan: dict[str, Any], style: PptxStyle) -> None:
    visual_plan = slide_plan.get("visual_plan") or {}
    proof = _proof_object(slide_plan)
    copyright_risk = bool(visual_plan.get("copyright_risk"))
    lines = [
        f"slide_no: {slide_plan.get('slide_no')}",
        f"layout_type: {slide_plan.get('layout_type')}",
        f"style_profile_loaded: {style.loaded}",
        f"style_profile_path: {_display_path(style.profile_path) if style.profile_path else None}",
        f"applied_font_family: {style.font_family}",
        f"applied_fallback_font: {style.fallback_font}",
        f"source_slide_refs: {_note_json(_as_list(slide_plan.get('source_slide_refs')))}",
        f"proof_object_type: {proof.type}",
        f"proof_object_required: {proof.required}",
        f"proof_object_screen_position: {proof.screen_position}",
        f"manual_placeholder_hidden: {style.loaded and _visual_kind(slide_plan) == 'manual'}",
        "",
        "screen_body_lines:",
        _note_json(_screen_body_lines(slide_plan)),
        "overflow_body_lines:",
        _note_json(_overflow_body_lines(slide_plan)),
        "",
        "source_urls:",
    ]
    lines.extend(f"- {url}" for url in _as_list(slide_plan.get("source_urls")))
    lines.append("image_urls:")
    lines.extend(f"- {url}" for url in _as_list(slide_plan.get("image_urls")))
    lines.extend(
        [
            "",
            "flags:",
            f"- needs_source: {bool(slide_plan.get('needs_source'))}",
            f"- needs_fact_check: {bool(slide_plan.get('needs_fact_check'))}",
            f"- required_before_broadcast: {bool(slide_plan.get('required_before_broadcast'))}",
            f"- copyright_risk: {copyright_risk}",
            "",
            "visual_plan:",
            _note_json(visual_plan),
            "",
            "proof_object:",
            _note_json(proof.as_dict()),
            "",
            "edit_notes:",
        ]
    )
    lines.extend(f"- {note}" for note in _as_list(slide_plan.get("edit_notes")))
    if slide_plan.get("layout_type") == "appendix_checklist":
        lines.extend(["", "internal_note: appendix/internal checklist, not a broadcast slide"])
    if copyright_risk:
        lines.extend(["", "image_note: 이미지 저작권 확인 필요"])
    speaker_notes = str(slide_plan.get("speaker_notes") or "").strip()
    if speaker_notes:
        lines.extend(["", "speaker_notes:", speaker_notes])
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "\n".join(lines)


def _ordered_slides(deck_plan: dict[str, Any]) -> list[dict[str, Any]]:
    slides = list(deck_plan.get("slides", []))
    appendix = [
        slide
        for slide in slides
        if slide.get("slide_type") == "production_checklist"
        or slide.get("layout_type") == "appendix_checklist"
    ]
    main = [slide for slide in slides if slide not in appendix]
    return [*main, *appendix]


def _long_body_slides(deck_plan: dict[str, Any]) -> list[int]:
    slide_nos = []
    for slide in deck_plan.get("slides", []):
        body = _body_lines(slide)
        total_chars = sum(len(line) for line in body)
        has_visual = _has_visual(slide)
        if (
            len(body) >= 4
            or total_chars > 220
            or any(len(line) > 100 for line in body)
            or (has_visual and len(body) >= 3)
        ):
            slide_nos.append(int(slide.get("slide_no") or 0))
    return slide_nos


def _fallback_body_font_for_slide(slide: dict[str, Any]) -> int:
    layout_type = slide.get("layout_type") or "headline_body"
    if layout_type == "title":
        return 22
    if layout_type == "section_title":
        return 20
    if layout_type == "big_headline":
        return 24
    if layout_type == "quote":
        return 20
    if layout_type in {"question", "closing_question"}:
        return 23
    if layout_type == "comparison":
        return 19
    if layout_type in {"chart_placeholder", "image_placeholder", "timeline"}:
        return 17
    if layout_type in {"checklist", "appendix_checklist"}:
        return 18
    return 20


def _planned_body_font_size(slide: dict[str, Any], style: PptxStyle) -> int:
    return _body_font_size(
        style,
        _screen_body_lines(slide),
        _fallback_body_font_for_slide(slide),
        has_visual=_has_screen_visual(slide, style),
    )


def _body_line_estimates(deck_plan: dict[str, Any], style: PptxStyle) -> list[dict[str, int]]:
    estimates = []
    for slide in deck_plan.get("slides", []):
        font_size = _planned_body_font_size(slide, style)
        estimates.append(
            {
                "slide_no": int(slide.get("slide_no") or 0),
                "body_line_estimate": _body_line_estimate(
                    slide,
                    font_size,
                    has_visual=_has_screen_visual(slide, style),
                ),
                "font_size": font_size,
            }
        )
    return estimates


def _font_size_downgraded_slides(
    deck_plan: dict[str, Any],
    style: PptxStyle,
) -> list[dict[str, int]]:
    if not style.loaded:
        return []
    downgraded = []
    for slide in deck_plan.get("slides", []):
        body = _body_lines(slide)
        if not body:
            continue
        font_size = _planned_body_font_size(slide, style)
        if font_size < 28:
            downgraded.append(
                {
                    "slide_no": int(slide.get("slide_no") or 0),
                    "font_size": font_size,
                    "body_lines": len(body),
                }
            )
    return downgraded


def _visually_dense_slides(deck_plan: dict[str, Any], style: PptxStyle) -> list[int]:
    dense = []
    for slide in deck_plan.get("slides", []):
        font_size = _planned_body_font_size(slide, style)
        if _is_visually_dense(slide, font_size):
            dense.append(int(slide.get("slide_no") or 0))
    return dense


def _visual_and_long_body_slides(deck_plan: dict[str, Any], style: PptxStyle) -> list[int]:
    slides = []
    for slide in deck_plan.get("slides", []):
        if not _has_screen_visual(slide, style):
            continue
        font_size = _planned_body_font_size(slide, style)
        if _body_line_estimate(slide, font_size, has_visual=True) >= 3:
            slides.append(int(slide.get("slide_no") or 0))
    return slides


def _planned_text_placeholder_rects(
    slide: dict[str, Any],
    style: PptxStyle,
) -> tuple[Rect, Rect] | None:
    if not _has_screen_visual(slide, style):
        return None
    layout_type = slide.get("layout_type") or "headline_body"
    if _is_chart_table_slide(slide):
        return None
    if layout_type == "quote":
        left, top, width, _height = _layout_box(style, "quote", (0.8, 1.0, 11.7, 4.9))
        if style.loaded and _prefers_left_image(slide, style):
            return (
                Rect(6.35, top + 0.8, 6.1, 2.7),
                Rect(0.75, 1.65, 5.25, 4.15),
            )
        return Rect(left, 2.05 if style.loaded else top, min(10.6, width), 2.7), Rect(
            0.0,
            0.0,
            0.0,
            0.0,
        )
    if layout_type in {"question", "closing_question", "comparison"} and _prefers_left_image(
        slide,
        style,
    ):
        return Rect(6.55, 2.3 if layout_type != "comparison" else 1.65, 5.8, 3.2), Rect(
            0.78,
            1.55,
            5.35,
            4.5,
        )
    if layout_type == "big_headline":
        left, top, width, height = _layout_box(style, "big_headline", (0.7, 1.0, 11.8, 1.4))
        if style.loaded and _prefers_left_image(slide, style):
            return (
                Rect(6.55, 2.05, 5.8, 2.9),
                Rect(0.78, 1.55, 5.35, 4.5),
            )
        body_width = 7.4 if _has_screen_visual(slide, style) else 10.8
        return (
            Rect(0.78, 2.05 if style.loaded else max(top + height + 0.35, 2.8), body_width, 2.9),
            Rect(9.25, 4.8, 2.8, 1.2),
        )
    if layout_type in {"image_placeholder", "timeline"}:
        profile_layout = "image_heavy"
        left, top, width, height = _layout_box(style, profile_layout, (1.15, 1.65, 11.0, 3.0))
        if style.loaded:
            return Rect(6.55, 1.65, 5.8, 3.65), Rect(0.78, 1.55, 5.35, 4.5)
        height = max(2.0, height)
        return (
            Rect(left, min(top + height + 0.35, 5.4), min(11.0, width), 1.35),
            Rect(left, top, width, max(1.8, height)),
        )
    left, top, width, height = _layout_box(style, "headline_body", (0.65, 1.75, 7.7, 4.7))
    if style.loaded:
        if _prefers_left_image(slide, style):
            left = 6.55
            width = 5.9
            placeholder = Rect(0.78, 1.55, 5.35, 4.5)
        else:
            width = min(width, 7.35)
            top = 2.05
            placeholder = Rect(8.75, 1.72, 3.65, 3.95)
        height = min(height, 3.85)
    else:
        placeholder = Rect(8.55, 1.55, 4.05, 4.65)
    return Rect(left, top, width, height), placeholder


def _text_placeholder_overlap_warnings(
    deck_plan: dict[str, Any],
    style: PptxStyle,
) -> list[dict[str, Any]]:
    warnings = []
    for slide in deck_plan.get("slides", []):
        rects = _planned_text_placeholder_rects(slide, style)
        if not rects:
            continue
        text_rect, placeholder_rect = rects
        if _rects_overlap(text_rect, placeholder_rect):
            warnings.append(
                {
                    "slide_no": int(slide.get("slide_no") or 0),
                    "layout_type": slide.get("layout_type") or "headline_body",
                }
            )
    return warnings


def _screen_body_overflow_slides(deck_plan: dict[str, Any]) -> list[int]:
    return [
        int(slide.get("slide_no") or 0)
        for slide in deck_plan.get("slides", [])
        if _overflow_body_lines(slide)
    ]


def _split_recommended_slides(deck_plan: dict[str, Any]) -> list[int]:
    slides = []
    for slide in deck_plan.get("slides", []):
        body = _body_lines(slide)
        if len(body) >= 4 or any(len(line) > 100 for line in body):
            slides.append(int(slide.get("slide_no") or 0))
    return slides


def _slides_using_20pt(deck_plan: dict[str, Any], style: PptxStyle) -> list[int]:
    return [
        item["slide_no"]
        for item in _body_line_estimates(deck_plan, style)
        if item.get("font_size") == 20
    ]


def _headline_red_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    if not style.loaded:
        return 0
    return sum(
        1
        for slide in deck_plan.get("slides", [])
        if slide.get("layout_type") not in {"section_title", "appendix_checklist"}
    )


def _body_black_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    if not style.loaded:
        return 0
    return sum(
        1
        for slide in deck_plan.get("slides", [])
        if _screen_body_lines(slide) and not _is_bilingual_quote_slide(slide)
    )


def _manual_placeholder_hidden_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    if not style.loaded:
        return 0
    return sum(1 for slide in deck_plan.get("slides", []) if _visual_kind(slide) == "manual")


def _proof_objects(deck_plan: dict[str, Any]) -> list[ProofObject]:
    return [_proof_object(slide) for slide in deck_plan.get("slides", [])]


def _proof_object_type_counts(deck_plan: dict[str, Any]) -> dict[str, int]:
    return dict(
        Counter(proof.type for proof in _proof_objects(deck_plan) if proof.type != "none")
    )


def _layout_template(slide: dict[str, Any]) -> str:
    proof = _proof_object(slide)
    if proof.type in {"chart", "table"}:
        return "chart_table_reference"
    if proof.type == "article_quote":
        return "source_card_or_article_quote"
    if proof.type in LEFT_PROOF_TYPES:
        return "image_left_quote_right"
    return "text_only_calculation"


def _layout_template_counts(deck_plan: dict[str, Any]) -> dict[str, int]:
    return dict(Counter(_layout_template(slide) for slide in deck_plan.get("slides", [])))


def _legacy_proof_object_type(slide_plan: dict[str, Any]) -> str:
    kind = _visual_kind(slide_plan)
    layout_type = str(slide_plan.get("layout_type") or "headline_body")
    slide_type = str(slide_plan.get("slide_type") or "")
    if _is_chart_table_slide(slide_plan) or kind == "chart_candidate":
        return "table" if layout_type in {"table", "data_table"} else "chart"
    if kind == "diagram":
        return "diagram"
    if kind == "screenshot_candidate":
        return "screenshot"
    if kind in {"photo_candidate", "image_placeholder"}:
        return "image"
    if kind == "ai_image_prompt":
        return "generated_image_candidate"
    if kind == "manual" and (layout_type == "quote" or slide_type == "quote"):
        return "article_quote"
    if layout_type == "quote" or slide_type == "quote":
        return "article_quote"
    return "none"


def _legacy_proof_object_slide_count(deck_plan: dict[str, Any]) -> int:
    return sum(
        1 for slide in deck_plan.get("slides", []) if _legacy_proof_object_type(slide) != "none"
    )


def _source_backed_text_only_slides(deck_plan: dict[str, Any]) -> list[int]:
    return [
        int(slide.get("slide_no") or 0)
        for slide in deck_plan.get("slides", [])
        if _proof_object(slide).type == "none" and _has_source_urls(slide)
    ]


def _source_backed_text_only_should_have_card_slides(deck_plan: dict[str, Any]) -> list[int]:
    return [
        int(slide.get("slide_no") or 0)
        for slide in deck_plan.get("slides", [])
        if _proof_object(slide).type == "none"
        and _has_source_urls(slide)
        and not _is_internal_or_structural_slide(slide)
    ]


def _proof_object_slide_count(deck_plan: dict[str, Any]) -> int:
    return sum(1 for proof in _proof_objects(deck_plan) if proof.type != "none")


def _proof_object_area_reserved(slide: dict[str, Any], style: PptxStyle) -> bool:
    proof = _proof_object(slide)
    if proof.type == "none":
        return False
    if proof.type in {"chart", "table"}:
        return True
    return _screen_placeholder_visible(slide, style)


def _proof_object_area_reserved_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    return sum(
        1 for slide in deck_plan.get("slides", []) if _proof_object_area_reserved(slide, style)
    )


def _proof_object_required_but_missing_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    return sum(
        1
        for slide in deck_plan.get("slides", [])
        if _proof_object(slide).required and not _proof_object_area_reserved(slide, style)
    )


def _text_only_slide_count(deck_plan: dict[str, Any]) -> int:
    return sum(1 for proof in _proof_objects(deck_plan) if proof.type == "none")


def _text_only_dense_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    return sum(
        1
        for slide in deck_plan.get("slides", [])
        if _proof_object(slide).type == "none"
        and _is_visually_dense(slide, _planned_body_font_size(slide, style))
    )


def _chart_table_skeleton_count(deck_plan: dict[str, Any]) -> int:
    return sum(1 for proof in _proof_objects(deck_plan) if proof.type in {"chart", "table"})


def _article_quote_skeleton_count(deck_plan: dict[str, Any]) -> int:
    return sum(1 for proof in _proof_objects(deck_plan) if proof.type == "article_quote")


def _chart_table_style_applied_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    if not style.loaded:
        return 0
    return sum(1 for slide in deck_plan.get("slides", []) if _is_chart_table_slide(slide))


def _image_left_layout_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    if not style.loaded:
        return 0
    return sum(
        1
        for slide in deck_plan.get("slides", [])
        if _proof_object(slide).screen_position == "left_half"
        and _proof_object_area_reserved(slide, style)
    )


def _screen_footer_hidden_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    return len(deck_plan.get("slides", [])) if style.loaded else 0


def _applied_layout_count(deck_plan: dict[str, Any], style: PptxStyle) -> int:
    if not style.loaded:
        return 0
    aliases = {
        "chart_placeholder": "chart",
        "image_placeholder": "image_heavy",
        "timeline": "image_heavy",
    }
    count = 0
    for slide in deck_plan.get("slides", []):
        layout_type = slide.get("layout_type") or "headline_body"
        key = aliases.get(layout_type, layout_type)
        if key in style.layout_patterns:
            count += 1
    return count


def _parse_back_counts(output_path: Path) -> dict[str, Any]:
    parsed = parse_presentation(output_path)
    notes_text = "\n".join(str(slide.get("notes") or "") for slide in parsed["slides"])
    notes_lower = notes_text.lower()
    return {
        "parse_back_slide_count": parsed.get("slide_count"),
        "parse_back_notes_slide_count": sum(
            1 for slide in parsed["slides"] if str(slide.get("notes") or "").strip()
        ),
        "parse_back_source_url_count": parsed.get("unique_url_count"),
        "parse_back_needs_source_count": notes_lower.count("needs_source: true"),
        "parse_back_needs_fact_check_count": notes_lower.count("needs_fact_check: true"),
    }


def render_deck_plan_to_pptx(
    deck_plan: dict[str, Any],
    output_path: Path,
    *,
    style_profile: PptxStyle | None = None,
) -> dict[str, Any]:
    validation = validate_deck_plan(deck_plan)
    style = style_profile or load_style_profile(None)
    prs = Presentation()
    prs.slide_width = Inches(style.slide_width_in)
    prs.slide_height = Inches(style.slide_height_in)
    for slide_plan in _ordered_slides(deck_plan):
        _render_slide(prs, slide_plan, style)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    result = render_result(deck_plan, output_path, validation, style)
    return result


def _source_image_overlap_count(deck_plan: dict[str, Any]) -> int:
    count = 0
    for slide in deck_plan.get("slides", []):
        sources = {
            canonicalize_url(str(url))
            for url in _as_list(slide.get("source_urls"))
            if str(url).strip()
        }
        images = {
            canonicalize_url(str(url))
            for url in _as_list(slide.get("image_urls"))
            if str(url).strip()
        }
        if sources & images:
            count += 1
    return count


def render_result(
    deck_plan: dict[str, Any],
    output_path: Path,
    validation: dict[str, Any],
    style: PptxStyle,
) -> dict[str, Any]:
    slides = deck_plan.get("slides", [])
    appendix_count = sum(
        1
        for slide in slides
        if slide.get("slide_type") == "production_checklist"
        or slide.get("layout_type") == "appendix_checklist"
    )
    visual_placeholder_count = sum(
        1
        for slide in slides
        if (slide.get("visual_plan") or {}).get("kind") not in {None, "none"}
    )
    missing_notes = [
        slide.get("slide_no")
        for slide in slides
        if not str(slide.get("speaker_notes") or slide.get("notes") or "").strip()
    ]
    font_size_downgraded = _font_size_downgraded_slides(deck_plan, style)
    overlap_warnings = _text_placeholder_overlap_warnings(deck_plan, style)
    visually_dense = _visually_dense_slides(deck_plan, style)
    visual_and_long_body = _visual_and_long_body_slides(deck_plan, style)
    screen_overflow = _screen_body_overflow_slides(deck_plan)
    split_recommended = _split_recommended_slides(deck_plan)
    slides_using_20pt = _slides_using_20pt(deck_plan, style)
    proof_type_counts = _proof_object_type_counts(deck_plan)
    layout_template_counts = _layout_template_counts(deck_plan)
    legacy_proof_count = _legacy_proof_object_slide_count(deck_plan)
    proof_count = _proof_object_slide_count(deck_plan)
    source_backed_text_only = _source_backed_text_only_slides(deck_plan)
    source_backed_should_have_card = _source_backed_text_only_should_have_card_slides(deck_plan)
    warnings = list(validation.get("warnings", []))
    for warning in overlap_warnings:
        warnings.append(
            "slide {slide_no}: text/visual placeholder overlap risk ({layout_type})".format(
                **warning
            )
        )
    for slide_no in visually_dense:
        warnings.append(f"slide {slide_no}: visually dense body estimate")
    for slide_no in screen_overflow:
        warnings.append(f"slide {slide_no}: body overflow moved to speaker notes")
    for slide_no in slides_using_20pt:
        warnings.append(f"slide {slide_no}: body rendered at 20pt; consider split")
    result = {
        "deck_id": deck_plan.get("deck_id"),
        "input_deck_plan_path": deck_plan.get("_input_path"),
        "output_pptx_path": _display_path(output_path),
        "style_profile_path": _display_path(style.profile_path) if style.profile_path else None,
        "style_profile_loaded": style.loaded,
        "adaptive_font_applied": style.loaded,
        "visual_placeholder_shortened": style.loaded,
        "section_title_color_policy": (
            "theme/default black; do not force #FF0000"
            if style.loaded
            else "legacy scaffold"
        ),
        "headline_red_count": _headline_red_count(deck_plan, style),
        "body_black_count": _body_black_count(deck_plan, style),
        "bilingual_quote_slide_count": sum(
            1 for slide in slides if _is_bilingual_quote_slide(slide)
        ),
        "chart_table_style_applied_count": _chart_table_style_applied_count(
            deck_plan,
            style,
        ),
        "image_left_layout_count": _image_left_layout_count(deck_plan, style),
        "manual_placeholder_hidden_count": _manual_placeholder_hidden_count(deck_plan, style),
        "proof_object_slide_count": _proof_object_slide_count(deck_plan),
        "proof_object_type_counts": proof_type_counts,
        "layout_template_counts": layout_template_counts,
        "chart_table_reference_count": layout_template_counts.get(
            "chart_table_reference",
            0,
        ),
        "image_left_quote_right_count": layout_template_counts.get(
            "image_left_quote_right",
            0,
        ),
        "text_only_calculation_count": layout_template_counts.get(
            "text_only_calculation",
            0,
        ),
        "source_card_or_article_quote_count": layout_template_counts.get(
            "source_card_or_article_quote",
            0,
        ),
        "source_backed_text_only_count": len(source_backed_text_only),
        "source_backed_text_only_slides": source_backed_text_only,
        "source_backed_text_only_should_have_card_count": len(source_backed_should_have_card),
        "source_backed_text_only_should_have_card_slides": source_backed_should_have_card,
        "text_only_slide_count_before_after": {
            "before": len(slides) - legacy_proof_count,
            "after": len(slides) - proof_count,
        },
        "proof_object_slide_count_before_after": {
            "before": legacy_proof_count,
            "after": proof_count,
        },
        "proof_object_required_but_missing_count": _proof_object_required_but_missing_count(
            deck_plan,
            style,
        ),
        "proof_object_area_reserved_count": _proof_object_area_reserved_count(deck_plan, style),
        "text_only_slide_count": _text_only_slide_count(deck_plan),
        "text_only_dense_count": _text_only_dense_count(deck_plan, style),
        "chart_table_skeleton_count": _chart_table_skeleton_count(deck_plan),
        "article_quote_skeleton_count": _article_quote_skeleton_count(deck_plan),
        "proof_text_overlap_count": len(overlap_warnings),
        "screen_footer_hidden_count": _screen_footer_hidden_count(deck_plan, style),
        "screen_body_overflow_count": len(screen_overflow),
        "screen_body_overflow_slides": screen_overflow,
        "split_recommended_slide_count": len(split_recommended),
        "split_recommended_slides": split_recommended,
        "slides_using_20pt": slides_using_20pt,
        "applied_font_family": style.font_family,
        "applied_fallback_font": style.fallback_font,
        "applied_layout_count": _applied_layout_count(deck_plan, style),
        "font_size_downgraded_slides": font_size_downgraded,
        "slides_with_long_body": _long_body_slides(deck_plan),
        "visually_dense_slide_count": len(visually_dense),
        "visually_dense_slides": visually_dense,
        "slides_with_text_placeholder_overlap": overlap_warnings,
        "slides_with_visual_and_long_body": visual_and_long_body,
        "body_line_estimates": _body_line_estimates(deck_plan, style),
        "slides_with_font_fallback": [],
        "slide_count": len(slides),
        "section_count": len(deck_plan.get("sections", [])),
        "appendix_internal_slide_count": appendix_count,
        "needs_source_count": sum(1 for slide in slides if slide.get("needs_source")),
        "needs_fact_check_count": sum(1 for slide in slides if slide.get("needs_fact_check")),
        "visual_placeholder_count": visual_placeholder_count,
        "slides_with_missing_notes": missing_notes,
        "source_image_overlap_count": _source_image_overlap_count(deck_plan),
        "schema_valid": validation.get("schema_valid"),
        "warnings": warnings,
        "passed": (
            output_path.exists()
            and not missing_notes
            and _source_image_overlap_count(deck_plan) == 0
            and not (style.loaded and overlap_warnings)
        ),
    }
    result.update(_parse_back_counts(output_path))
    if style.loaded and style.font_family != style.fallback_font:
        result["warnings"].append(
            "Font substitution cannot be detected in Python; visually verify "
            f"{style.font_family} vs {style.fallback_font} fallback."
        )
    return result


def write_render_report(path: Path, results: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    passed = sum(1 for result in results if result.get("passed"))
    lines = [
        "# Piti PPTX Render Report",
        "",
        f"- Generated at: {datetime.now(UTC).isoformat()}",
        f"- Rendered deck count: {len(results)}",
        f"- Passed: {passed}",
        f"- Failed: {len(results) - passed}",
        "- Renderer type: scaffold PPTX skeleton",
        "- Styled draft support: optional Syukaworld style profile",
        "- Screen formatting rules: headline red 28pt; normal body black; "
        "bilingual quote Korean translation red",
        "- Chart/table rules: title 28pt bold underline; body/data labels 18pt "
        "bold; source 20pt underline",
        "- Manual visual placeholders: hidden from styled screens and preserved "
        "in speaker notes",
        "- Proof object scaffold: reserves screen areas for chart/table, article "
        "quote, image, screenshot, and diagram evidence objects without inserting assets",
        "- Reference layout templates v0: chart_table_reference, "
        "image_left_quote_right, text_only_calculation, source_card_or_article_quote",
        "- Styled screen footers: hidden from slides and preserved in speaker notes/report",
        "- Adaptive body font: enabled only as overflow protection for styled drafts",
        "- Visual placeholder text: shortened for styled drafts",
        "- Text/visual overlap policy: styled drafts fail on overlap; legacy "
        "scaffold outputs report overlap as warnings only",
        "- Production Piti agent: not implemented",
        "- Image auto collection/insertion: none",
        "- LLM/API calls: none",
        "",
        "## Decks",
        "",
        (
            "| Deck | Styled | Slides | Sections | Appendix | Needs Source | Needs Fact Check | "
            "Visuals | Proof Objects | Text Only | Dense | Overflow | Split | 20pt | "
            "Manual Hidden | Footer Hidden | Proof/Text Overlap | Missing Notes | Overlap | "
            "PPTX | Passed |"
        ),
        (
            "|---|---|---:|---:|---:|---:|---:|---:|---:|---:|---:|---:|---:|"
            "---:|---:|---:|---:|---:|---:|---|---|"
        ),
    ]
    for result in results:
        lines.append(
            "| {deck} | {styled} | {slides} | {sections} | {appendix} | {needs_source} | "
            "{needs_fact_check} | {visuals} | {proof_objects} | {text_only} | {dense} | "
            "{overflow} | {split} | {twenty_pt} | {manual_hidden} | {footer_hidden} | "
            "{proof_text_overlap} | {missing_notes} | {overlap} | {pptx} | {passed} |".format(
                deck=result.get("deck_id"),
                styled="yes" if result.get("style_profile_loaded") else "no",
                slides=result.get("slide_count"),
                sections=result.get("section_count"),
                appendix=result.get("appendix_internal_slide_count"),
                needs_source=result.get("needs_source_count"),
                needs_fact_check=result.get("needs_fact_check_count"),
                visuals=result.get("visual_placeholder_count"),
                proof_objects=result.get("proof_object_slide_count"),
                text_only=result.get("text_only_slide_count"),
                dense=result.get("visually_dense_slide_count"),
                overflow=result.get("screen_body_overflow_count"),
                split=result.get("split_recommended_slide_count"),
                twenty_pt=len(result.get("slides_using_20pt", [])),
                manual_hidden=result.get("manual_placeholder_hidden_count"),
                footer_hidden=result.get("screen_footer_hidden_count"),
                proof_text_overlap=result.get("proof_text_overlap_count"),
                missing_notes=len(result.get("slides_with_missing_notes", [])),
                overlap=result.get("source_image_overlap_count"),
                pptx=result.get("output_pptx_path"),
                passed="yes" if result.get("passed") else "no",
            )
        )
    lines.extend(["", "## Style Profile Application", ""])
    for result in results:
        lines.extend(
            [
                f"### {result.get('deck_id')} -> {result.get('output_pptx_path')}",
                "",
                f"- style_profile_path: {result.get('style_profile_path')}",
                f"- style_profile_loaded: {result.get('style_profile_loaded')}",
                f"- adaptive_font_applied: {result.get('adaptive_font_applied')}",
                f"- applied_font_family: {result.get('applied_font_family')}",
                f"- applied_fallback_font: {result.get('applied_fallback_font')}",
                f"- applied_layout_count: {result.get('applied_layout_count')}",
                "- font_size_downgraded_slides: "
                f"{result.get('font_size_downgraded_slides')}",
                f"- visual_placeholder_shortened: {result.get('visual_placeholder_shortened')}",
                f"- section_title_color_policy: {result.get('section_title_color_policy')}",
                f"- headline_red_count: {result.get('headline_red_count')}",
                f"- body_black_count: {result.get('body_black_count')}",
                "- bilingual_quote_slide_count: "
                f"{result.get('bilingual_quote_slide_count')}",
                "- chart_table_style_applied_count: "
                f"{result.get('chart_table_style_applied_count')}",
                f"- image_left_layout_count: {result.get('image_left_layout_count')}",
                "- manual_placeholder_hidden_count: "
                f"{result.get('manual_placeholder_hidden_count')}",
                f"- proof_object_slide_count: {result.get('proof_object_slide_count')}",
                f"- proof_object_type_counts: {result.get('proof_object_type_counts')}",
                f"- layout_template_counts: {result.get('layout_template_counts')}",
                "- chart_table_reference_count: "
                f"{result.get('chart_table_reference_count')}",
                "- image_left_quote_right_count: "
                f"{result.get('image_left_quote_right_count')}",
                "- text_only_calculation_count: "
                f"{result.get('text_only_calculation_count')}",
                "- source_card_or_article_quote_count: "
                f"{result.get('source_card_or_article_quote_count')}",
                "- source_backed_text_only_count: "
                f"{result.get('source_backed_text_only_count')}",
                "- source_backed_text_only_slides: "
                f"{result.get('source_backed_text_only_slides')}",
                "- source_backed_text_only_should_have_card_count: "
                f"{result.get('source_backed_text_only_should_have_card_count')}",
                "- source_backed_text_only_should_have_card_slides: "
                f"{result.get('source_backed_text_only_should_have_card_slides')}",
                "- text_only_slide_count_before_after: "
                f"{result.get('text_only_slide_count_before_after')}",
                "- proof_object_slide_count_before_after: "
                f"{result.get('proof_object_slide_count_before_after')}",
                "- proof_object_required_but_missing_count: "
                f"{result.get('proof_object_required_but_missing_count')}",
                "- proof_object_area_reserved_count: "
                f"{result.get('proof_object_area_reserved_count')}",
                f"- text_only_slide_count: {result.get('text_only_slide_count')}",
                f"- text_only_dense_count: {result.get('text_only_dense_count')}",
                f"- chart_table_skeleton_count: {result.get('chart_table_skeleton_count')}",
                "- article_quote_skeleton_count: "
                f"{result.get('article_quote_skeleton_count')}",
                f"- proof_text_overlap_count: {result.get('proof_text_overlap_count')}",
                f"- screen_footer_hidden_count: {result.get('screen_footer_hidden_count')}",
                f"- screen_body_overflow_count: {result.get('screen_body_overflow_count')}",
                f"- screen_body_overflow_slides: {result.get('screen_body_overflow_slides')}",
                "- split_recommended_slide_count: "
                f"{result.get('split_recommended_slide_count')}",
                f"- split_recommended_slides: {result.get('split_recommended_slides')}",
                f"- slides_using_20pt: {result.get('slides_using_20pt')}",
                f"- slides_with_long_body: {result.get('slides_with_long_body')}",
                f"- visually_dense_slides: {result.get('visually_dense_slides')}",
                "- slides_with_text_placeholder_overlap: "
                f"{result.get('slides_with_text_placeholder_overlap')}",
                "- slides_with_visual_and_long_body: "
                f"{result.get('slides_with_visual_and_long_body')}",
                f"- body_line_estimates: {result.get('body_line_estimates')}",
                f"- slides_with_font_fallback: {result.get('slides_with_font_fallback')}",
                f"- parse_back_slide_count: {result.get('parse_back_slide_count')}",
                f"- parse_back_notes_slide_count: {result.get('parse_back_notes_slide_count')}",
                f"- parse_back_source_url_count: {result.get('parse_back_source_url_count')}",
                f"- parse_back_needs_source_count: {result.get('parse_back_needs_source_count')}",
                "- parse_back_needs_fact_check_count: "
                f"{result.get('parse_back_needs_fact_check_count')}",
                "",
            ]
        )
    lines.extend(["", "## Warnings", ""])
    for result in results:
        warnings = result.get("warnings", [])
        if not warnings:
            lines.append(f"- {result.get('deck_id')}: none")
        else:
            for warning in warnings[:12]:
                lines.append(f"- {result.get('deck_id')}: {warning}")
            if len(warnings) > 12:
                lines.append(f"- {result.get('deck_id')}: ... +{len(warnings) - 12} more")
    lines.extend(
        [
            "",
            "## Readiness",
            "",
            "- ready_for_ppt_generation: true (scaffold only)",
            "- ready_for_production_piti_agent: false",
            "- ready_for_broadcast: false",
        ]
    )
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def render_default_decks(
    report_path: Path = DEFAULT_REPORT_PATH,
    *,
    style_profile_path: Path | None = DEFAULT_STYLE_PROFILE_PATH,
) -> list[Path]:
    output_paths: list[Path] = []
    results: list[dict[str, Any]] = []
    cases = [
        *DEFAULT_CASES,
        *[
            PptxRenderCase(
                deck_plan_path=case.deck_plan_path,
                output_path=case.output_path,
                style_profile_path=style_profile_path,
            )
            for case in STYLED_CASES
        ],
    ]
    for case in cases:
        deck_plan = _load_json(case.deck_plan_path)
        deck_plan["_input_path"] = str(case.deck_plan_path.relative_to(paths.REPO_ROOT))
        style = load_style_profile(case.style_profile_path)
        result = render_deck_plan_to_pptx(deck_plan, case.output_path, style_profile=style)
        result["input_deck_plan_path"] = str(case.deck_plan_path.relative_to(paths.REPO_ROOT))
        results.append(result)
        output_paths.append(case.output_path)
    write_render_report(report_path, results)
    output_paths.append(report_path)
    return output_paths


@app.callback(invoke_without_command=True)
def main(
    input_path: Annotated[
        Path | None,
        typer.Option("--input", help="Single Piti deck plan JSON to render."),
    ] = None,
    output_path: Annotated[
        Path | None,
        typer.Option("--output", help="PPTX output path for single render."),
    ] = None,
    report_path: Annotated[
        Path,
        typer.Option("--report", help="Markdown render report path for default renders."),
    ] = DEFAULT_REPORT_PATH,
    style_profile_path: Annotated[
        Path | None,
        typer.Option(
            "--style-profile",
            help="Optional Syukaworld style profile JSON for styled output.",
        ),
    ] = DEFAULT_STYLE_PROFILE_PATH,
) -> None:
    """Render deck plan JSON into editable PPTX skeletons."""
    if input_path:
        if output_path is None:
            raise typer.BadParameter("--output is required with --input")
        deck_plan = _load_json(input_path)
        deck_plan["_input_path"] = str(input_path)
        style = load_style_profile(style_profile_path)
        result = render_deck_plan_to_pptx(deck_plan, output_path, style_profile=style)
        console.print(
            f"[green]Rendered {output_path} "
            f"(slides={result['slide_count']}, passed={result['passed']}).[/green]"
        )
        raise typer.Exit(0 if result["passed"] else 1)
    rendered = render_default_decks(report_path=report_path, style_profile_path=style_profile_path)
    console.print(f"[green]Rendered {len(rendered)} Piti PPTX artifacts.[/green]")


if __name__ == "__main__":
    app()
